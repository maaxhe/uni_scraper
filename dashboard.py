"""
Stud.IP Dashboard — lokales Lern-Panel.
Starten: python dashboard.py  →  http://localhost:5001
"""

import json
import os
import re
import subprocess
import sys
from datetime import datetime, timedelta
from pathlib import Path

import markdown2
from dotenv import load_dotenv
from flask import Flask, Response, jsonify, render_template_string, request, send_file

load_dotenv()

COURSES_DIR      = Path(os.environ.get("COURSES_DIR", "/Users/maxmacbookpro/Documents/Uni/Cognitive Science [Course]/Courses"))
PYTHON           = sys.executable
SUMMARIZE_SCRIPT = str(Path(__file__).parent / "summarize.py")
SCRAPER_SCRIPT   = str(Path(__file__).parent / "scraper.py")
PIPELINE_LOG     = str(Path(__file__).parent / "pipeline.log")
OUTPUT_FILENAME  = "_zusammenfassung.md"
SUMMARY_RE       = re.compile(r'^_zusammenfassung(?!.*_qa).*\.md$')
NOTES_FILENAME      = "_notizen.md"
FILE_NOTES_FILENAME  = "_file_notes.json"
CUSTOM_INFO_FILENAME = "_custom_info.json"
PROGRESS_FILE    = Path(__file__).parent / "progress.json"
USER_STATE_FILE  = Path(__file__).parent / "user_state.json"
SUPPORTED_EXT    = {".pdf", ".doc", ".docx", ".txt", ".md", ".pptx"}

app = Flask(__name__)

# ---------------------------------------------------------------------------
# Progress store
# ---------------------------------------------------------------------------

def load_progress() -> dict:
    if PROGRESS_FILE.exists():
        return json.loads(PROGRESS_FILE.read_text())
    return {}

def save_progress(data: dict):
    PROGRESS_FILE.write_text(json.dumps(data, indent=2))

# ---------------------------------------------------------------------------
# User state store (favourites, read files, archived, hidden, pins, todos)
# ---------------------------------------------------------------------------

def load_user_state() -> dict:
    if USER_STATE_FILE.exists():
        return json.loads(USER_STATE_FILE.read_text())
    return {}

def save_user_state(data: dict):
    USER_STATE_FILE.write_text(json.dumps(data, indent=2))

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def list_summaries(d: Path) -> list[dict]:
    """Return all summary files in d, sorted newest first."""
    files = sorted(
        [f for f in d.iterdir() if f.is_file() and SUMMARY_RE.match(f.name)],
        key=lambda f: f.stat().st_mtime,
        reverse=True,
    )
    return [{"name": f.name, "mtime": f.stat().st_mtime} for f in files]

def get_latest_summary(d: Path) -> Path | None:
    """Return path to the most recently modified summary file, or None."""
    summaries = list_summaries(d)
    return (d / summaries[0]["name"]) if summaries else None

def _course_info(rel_path: str, d: Path, progress: dict) -> dict:
    """Build info dict for a single course directory."""
    files = list_files(d)
    summary_path = get_latest_summary(d) or (d / OUTPUT_FILENAME)
    notes_path   = d / NOTES_FILENAME
    sync_path    = d / "_last_sync"
    p = progress.get(rel_path, progress.get(d.name, {}))
    summary_mtime = int(summary_path.stat().st_mtime) if summary_path.exists() else None
    sync_mtime    = int(sync_path.read_text().strip()) if sync_path.exists() else None
    # Count files newer than the last sync
    new_files = 0
    new_file_names = []
    if sync_mtime:
        for fname in files:
            fpath = d / fname
            if fpath.exists() and int(fpath.stat().st_mtime) > sync_mtime:
                new_files += 1
                new_file_names.append(fname)
    return {
        "name":           d.name,
        "path":           rel_path,
        "has_summary":    summary_path.exists(),
        "has_notes":      notes_path.exists(),
        "file_count":     len(files),
        "new_files":      new_files,
        "new_file_names": new_file_names,
        "summary_age": summary_mtime,
        "sync_age":    sync_mtime,
        "progress":    {
            "total":        p.get("total", 0),
            "known":        p.get("known", 0),
            "last_studied": p.get("last_studied"),
        },
        "is_group":    False,
    }

SEMESTER_RE = re.compile(r'(?:SoSe|WiSe|SS|WS)\s*\d{2}', re.IGNORECASE)

def _semester_sort_key(name: str) -> int:
    """
    Returns a numeric sort key so semesters order newest-first.
    SoSe YYYY → YYYY * 2 + 1   (e.g. SoSe 2026 = 4053)
    WiSe YYYY_ZZ → YYYY * 2 + 2 (e.g. WiSe 2025_26 = 4052)
    Higher = newer.
    """
    m = re.search(r'(\d{4})', name)
    if not m:
        return 0
    year = int(m.group(1))
    if re.search(r'(?:SoSe|SS)', name, re.IGNORECASE):
        return year * 2 + 1
    else:  # WiSe / WS
        return year * 2 + 2

def get_courses():
    progress = load_progress()
    dirs = [d for d in COURSES_DIR.iterdir() if d.is_dir()]

    semester_items = []
    other_items = []

    for d in sorted(dirs, key=lambda d: d.name):
        is_semester = bool(SEMESTER_RE.search(d.name))
        sub_dirs = sorted([sd for sd in d.iterdir() if sd.is_dir() and not sd.name.startswith('.')])
        if is_semester or len(sub_dirs) >= 2:
            sub_courses = [
                _course_info(f"{d.name}/{sd.name}", sd, progress)
                for sd in sub_dirs
            ]
            entry = {
                "name":        d.name,
                "path":        d.name,
                "is_group":    True,
                "is_semester": is_semester,
                "courses":     sub_courses,
            }
            if is_semester:
                semester_items.append(entry)
            else:
                other_items.append(entry)
        else:
            other_items.append(_course_info(d.name, d, progress))

    # Semester groups newest-first using proper calendar order
    semester_items.sort(key=lambda x: _semester_sort_key(x["name"]), reverse=True)
    return semester_items + other_items

def list_files(course_dir: Path) -> list[str]:
    return sorted([
        str(f.relative_to(course_dir)) for f in course_dir.rglob("*")
        if f.is_file()
        and f.suffix.lower() != ".json"
        and not SUMMARY_RE.match(f.name)
        and f.name != NOTES_FILENAME
        and f.name != "_last_sync"
        and ".summary" not in f.name
    ])

def read_file_text(course_name: str, filename: str) -> str:
    path = COURSES_DIR / course_name / filename
    if not path.exists():
        return ""
    try:
        return _extract_file_text(path) or ""
    except Exception as e:
        return f"Fehler beim Lesen: {e}"

def get_qa_file(summary_path: Path) -> Path:
    """Derive the _qa.md path from a summary path."""
    return summary_path.parent / (summary_path.stem + "_qa.md")

def parse_flashcards(qa_md: str) -> list[dict]:
    """Extract Q&A pairs from a _qa.md file."""
    cards = []
    sections = re.split(r'\n## ', qa_md)
    for section in sections:
        lines = section.strip().split('\n')
        section_title = lines[0].strip('# ').strip() if lines else "Unknown"

        q_block = re.search(r'\*\*(?:Training Questions|Trainingsfragen)\*\*\s*\n(.*?)(?=\n\*\*|\Z)', section, re.DOTALL)
        a_block = re.search(r'\*\*(?:Answers|Antworten)\*\*\s*\n(.*?)(?=\n\*\*|\Z)', section, re.DOTALL)

        if not q_block:
            continue

        questions = re.findall(r'\d+\.\s+(.+)', q_block.group(1))
        answers   = re.findall(r'\d+\.\s+(.+)', a_block.group(1)) if a_block else []

        for i, q in enumerate(questions):
            cards.append({
                "id":      f"{section_title}_{i}",
                "section": section_title,
                "q":       q.strip(),
                "a":       answers[i].strip() if i < len(answers) else "–",
            })
    return cards

_TAIL_BYTES = 8192  # read only the last 8 KB of the log — enough for recent timestamps

def get_pipeline_status() -> dict:
    log_path = Path(PIPELINE_LOG)
    last_run = None
    last_ok  = None
    if log_path.exists():
        size = log_path.stat().st_size
        with log_path.open("rb") as fh:
            fh.seek(max(0, size - _TAIL_BYTES))
            tail = fh.read().decode("utf-8", errors="replace")
        runs = re.findall(r'(\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}) Pipeline gestartet', tail)
        ends = re.findall(r'(\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}) Pipeline beendet', tail)
        last_run = runs[-1] if runs else None
        last_ok  = ends[-1] if ends else None
    return {"last_run": last_run, "last_ok": last_ok}

# ---------------------------------------------------------------------------
# API
# ---------------------------------------------------------------------------

@app.route("/api/courses")
def api_courses():
    data = get_courses()
    # Inject global last-sync timestamp (courses.json mtime as fallback)
    last_sync = int(COURSES_JSON.stat().st_mtime) if COURSES_JSON.exists() else None
    return jsonify({"tree": data, "last_sync": last_sync})

def _old_material_paths(course_name: str) -> set:
    """
    For a semester-scoped course (e.g. 'SoSe 2026/Neurodynamics') return the
    set of relative file paths that also exist under 'Alle Kurse/<basename>'.
    These are files the professor carried over from previous offerings.
    Returns an empty set for courses not nested in a semester folder.
    """
    parts = course_name.rstrip("/").split("/")
    if len(parts) < 2:
        return set()
    basename = parts[-1]
    alle_dir = COURSES_DIR / "Alle Kurse" / basename
    if not alle_dir.is_dir():
        return set()
    return set(list_files(alle_dir))

@app.route("/api/files/<path:course_name>")
def api_files(course_name):
    d = COURSES_DIR / course_name
    if not d.is_dir():
        return jsonify({"files": [], "old_files": []})
    all_files = list_files(d)
    old_set   = _old_material_paths(course_name)
    old_files     = [f for f in all_files if f in old_set]
    current_files = [f for f in all_files if f not in old_set]
    return jsonify({"files": current_files, "old_files": old_files})

@app.route("/api/file-text/<path:course_name>/<path:filename>")
def api_file_text(course_name, filename):
    text = read_file_text(course_name, filename)
    return jsonify({"text": text})

@app.route("/api/file-raw/<path:course_name>/<path:filename>")
def api_file_raw(course_name, filename):
    path = COURSES_DIR / course_name / filename
    if not path.exists():
        return "Not found", 404
    return send_file(str(path))

@app.route("/api/summaries/<path:course_name>")
def api_summaries(course_name):
    d = COURSES_DIR / course_name
    return jsonify(list_summaries(d))

def _protect_latex(md: str):
    """Replace LaTeX blocks with placeholders so markdown2 can't mangle them."""
    blocks: dict[str, str] = {}
    counter = [0]
    def store(m):
        key = f"\x00LATEX{counter[0]}\x00"
        blocks[key] = m.group(0)
        counter[0] += 1
        return key
    # display math first ($$...$$), then inline ($...$)
    md = re.sub(r'\$\$[\s\S]+?\$\$', store, md)
    md = re.sub(r'\$[^$\n]+?\$', store, md)
    return md, blocks

def _restore_latex(html: str, blocks: dict) -> str:
    for key, val in blocks.items():
        html = html.replace(key, val)
    return html

@app.route("/api/summary/<path:course_name>")
def api_summary(course_name):
    filename = request.args.get("file")
    d = COURSES_DIR / course_name
    path = (d / filename) if filename else get_latest_summary(d)
    if not path or not path.exists():
        return jsonify({"html": None, "md": None})
    md   = path.read_text(encoding="utf-8")
    protected_md, latex_blocks = _protect_latex(md)
    html = markdown2.markdown(protected_md, extras=["fenced-code-blocks", "tables"])
    html = _restore_latex(html, latex_blocks)
    return jsonify({"html": html, "md": md, "file": path.name})

@app.route("/api/summary-delete/<path:course_name>", methods=["POST"])
def api_summary_delete(course_name):
    filename = (request.json or {}).get("file")
    if not filename:
        return jsonify({"error": "No file specified"}), 400
    p = COURSES_DIR / course_name / filename
    if not p.exists() or not SUMMARY_RE.match(p.name):
        return jsonify({"error": "File not found"}), 404
    p.unlink()
    return jsonify({"ok": True})

@app.route("/api/summary-rename/<path:course_name>", methods=["POST"])
def api_summary_rename(course_name):
    body    = request.json or {}
    old_name = body.get("file")
    new_label = body.get("label", "").strip()
    if not old_name or not new_label:
        return jsonify({"error": "Missing params"}), 400
    d = COURSES_DIR / course_name
    old_p = d / old_name
    if not old_p.exists() or not SUMMARY_RE.match(old_p.name):
        return jsonify({"error": "File not found"}), 404
    # Sanitise label → filename
    safe = re.sub(r'[^\w\- ]', '', new_label).strip().replace(' ', '_')
    if not safe:
        return jsonify({"error": "Invalid name"}), 400
    new_name = f"_zusammenfassung_{safe}.md"
    new_p = d / new_name
    if new_p.exists():
        return jsonify({"error": "Name already in use"}), 409
    old_p.rename(new_p)
    return jsonify({"ok": True, "file": new_name})

@app.route("/api/summary-save/<path:course_name>", methods=["POST"])
def api_summary_save(course_name):
    body     = request.json or {}
    filename = body.get("file")
    content  = body.get("content", "")
    if not filename:
        return jsonify({"error": "No file"}), 400
    p = COURSES_DIR / course_name / filename
    if not p.exists() or not SUMMARY_RE.match(p.name):
        return jsonify({"error": "File not found"}), 404
    p.write_text(content, encoding="utf-8")
    return jsonify({"ok": True})

@app.route("/api/summary-raw/<path:course_name>")
def api_summary_raw(course_name):
    filename = request.args.get("file")
    d = COURSES_DIR / course_name
    path = (d / filename) if filename else get_latest_summary(d)
    if not path or not path.exists():
        return "Not found", 404
    return send_file(str(path), as_attachment=True, download_name=f"{course_name}_{path.name}")

@app.route("/api/download-zip/<path:course_name>")
def api_download_zip(course_name):
    import zipfile, io
    course_dir = COURSES_DIR / course_name
    buf = io.BytesIO()
    include_notes = request.args.get('notes', '0') == '1'
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in course_dir.rglob('*'):
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXT and not SUMMARY_RE.match(f.name) and '_qa' not in f.name:
                zf.write(f, f.relative_to(course_dir))
        if include_notes:
            notes_path = course_dir / NOTES_FILENAME
            if notes_path.exists():
                zf.write(notes_path, NOTES_FILENAME)
    buf.seek(0)
    safe = re.sub(r'[^\w\-]', '_', course_name.split('/')[-1])
    return send_file(buf, mimetype='application/zip', as_attachment=True, download_name=f'{safe}_files.zip')

@app.route("/api/flashcards/<path:course_name>")
def api_flashcards(course_name):
    summary = get_latest_summary(COURSES_DIR / course_name)
    if not summary:
        return jsonify([])
    qa_path = get_qa_file(summary)
    if not qa_path.exists():
        return jsonify([])
    cards = parse_flashcards(qa_path.read_text(encoding="utf-8"))
    return jsonify(cards)

def _collect_flashcards(rel_path: str, d: Path) -> list:
    summary = get_latest_summary(d)
    if not summary or not summary.exists():
        return []
    qa_path = get_qa_file(summary)
    if not qa_path.exists():
        return []
    cards = parse_flashcards(qa_path.read_text(encoding="utf-8"))
    for c in cards:
        c["course"] = d.name
        c["id"] = f"{rel_path}::{c['id']}"
    return cards

CUSTOM_CARDS_FILE  = "_custom_cards.json"
CHAT_HISTORY_FILE  = "_chat_history.json"

@app.route("/api/custom-cards/<path:course_name>", methods=["GET"])
def api_get_custom_cards(course_name):
    p = COURSES_DIR / course_name / CUSTOM_CARDS_FILE
    if not p.exists():
        return jsonify([])
    return jsonify(json.loads(p.read_text(encoding="utf-8")))

@app.route("/api/custom-cards/<path:course_name>", methods=["POST"])
def api_save_custom_cards(course_name):
    cards = request.json
    p = COURSES_DIR / course_name / CUSTOM_CARDS_FILE
    p.write_text(json.dumps(cards, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"ok": True})

@app.route("/api/srs/<path:course_name>", methods=["GET"])
def api_get_srs(course_name):
    p = COURSES_DIR / course_name / "_srs.json"
    if not p.exists():
        return jsonify({})
    return jsonify(json.loads(p.read_text(encoding="utf-8")))

@app.route("/api/srs/<path:course_name>", methods=["POST"])
def api_save_srs(course_name):
    p = COURSES_DIR / course_name / "_srs.json"
    p.write_text(json.dumps(request.json, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"ok": True})

@app.route("/api/chat-history/<path:course_name>", methods=["GET"])
def api_get_chat_history(course_name):
    p = COURSES_DIR / course_name / CHAT_HISTORY_FILE
    if not p.exists():
        return jsonify([])
    return jsonify(json.loads(p.read_text(encoding="utf-8")))

@app.route("/api/chat-history/<path:course_name>", methods=["POST"])
def api_save_chat_history(course_name):
    data = request.json  # {title, messages: [{role, content}]}
    p = COURSES_DIR / course_name / CHAT_HISTORY_FILE
    history = json.loads(p.read_text(encoding="utf-8")) if p.exists() else []
    entry = {
        "id":       datetime.now().strftime("%Y%m%d_%H%M%S"),
        "title":    data.get("title", "Conversation"),
        "date":     datetime.now().strftime("%Y-%m-%d %H:%M"),
        "messages": data.get("messages", []),
    }
    history.insert(0, entry)   # newest first
    p.write_text(json.dumps(history, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"ok": True, "id": entry["id"]})

@app.route("/api/chat-history/<path:course_name>/rename", methods=["PATCH"])
def api_rename_chat_history(course_name):
    data = request.json  # {id, title}
    p = COURSES_DIR / course_name / CHAT_HISTORY_FILE
    if not p.exists():
        return jsonify({"error": "not found"}), 404
    history = json.loads(p.read_text(encoding="utf-8"))
    for entry in history:
        if entry.get("id") == data.get("id"):
            entry["title"] = data.get("title", entry["title"]).strip() or entry["title"]
            break
    p.write_text(json.dumps(history, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"ok": True})

@app.route("/api/chat-history/<path:course_name>/delete", methods=["DELETE"])
def api_delete_chat_history(course_name):
    conv_id = request.json.get("id")
    p = COURSES_DIR / course_name / CHAT_HISTORY_FILE
    if not p.exists():
        return jsonify({"error": "not found"}), 404
    history = json.loads(p.read_text(encoding="utf-8"))
    history = [e for e in history if e.get("id") != conv_id]
    p.write_text(json.dumps(history, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"ok": True})

@app.route("/api/generate-cards/<path:course_name>", methods=["POST"])
def api_generate_cards(course_name):
    course_dir   = COURSES_DIR / course_name
    summary_path = get_latest_summary(course_dir)
    if not summary_path or not summary_path.exists():
        return jsonify({"error": "No summary found"}), 404
    body  = request.json or {}
    count = max(3, min(int(body.get("count", 10)), 30))
    # Allow specifying a particular summary file
    specific = body.get("file")
    if specific:
        sp = COURSES_DIR / course_name / specific
        if sp.exists() and SUMMARY_RE.match(sp.name):
            summary_path = sp
    summary = summary_path.read_text(encoding="utf-8")[:12000]
    prompt = f"""Create exactly {count} high-quality flashcards from the following course summary.

Use a MIX of these card types (spread them across the content, don't cluster):
- "recall"      — direct fact or definition: "What is X?" / "Define X."
- "mechanism"   — cause, process, or reasoning: "Why does X happen?" / "How does X work?"
- "contrast"    — comparison: "What is the difference between X and Y?"
- "application" — applying a concept: "Given [scenario], what would [result/approach] be?"
- "cloze"       — fill-in-the-blank: write a key sentence from the material with ONE important term replaced by ___ ; the answer is ONLY that missing term/phrase (keep it short).

Target distribution: ~25% recall, ~25% mechanism, ~20% contrast, ~15% application, ~15% cloze.
Avoid trivial or overly broad questions. Prefer specific, testable knowledge.

LaTeX rules (apply to ALL types, both q and a):
- Inline math: $...$ (e.g. $E = mc^2$, $\\alpha$, $\\nabla f$)
- Block math: $$...$$ for standalone equations
- Never write math as plain text.

Reply ONLY as a JSON array, no explanation:
[{{"type": "recall", "q": "...", "a": "..."}}, {{"type": "cloze", "q": "The ___ controls ...", "a": "missing term"}}, ...]

Summary:
{summary}"""

    anthropic_key = os.environ.get("ANTHROPIC_API_KEY")
    openai_key    = os.environ.get("OPENAI_API_KEY")

    if anthropic_key:
        import anthropic as _anthropic
        model = os.environ.get("ANTHROPIC_MODEL", "claude-haiku-4-5-20251001")
        client = _anthropic.Anthropic(api_key=anthropic_key)
        try:
            msg = client.messages.create(
                model=model, max_tokens=2000,
                messages=[{"role": "user", "content": prompt}],
            )
        except _anthropic.AuthenticationError:
            return jsonify({"error": "❌ API key invalid — create a new key at https://console.anthropic.com"}), 401
        except _anthropic.PermissionDeniedError:
            return jsonify({"error": "❌ API key lacks permission — check your Anthropic account and plan."}), 403
        except _anthropic.RateLimitError:
            return jsonify({"error": "❌ Rate limit hit — too many requests. Wait a moment and try again."}), 429
        except _anthropic.APIStatusError as e:
            if e.status_code in (402, 403):
                return jsonify({"error": f"❌ API quota exhausted or billing issue (HTTP {e.status_code}). Check https://console.anthropic.com"}), 402
            return jsonify({"error": f"❌ Anthropic API error (HTTP {e.status_code}): {e.message}"}), 500
        except _anthropic.APIConnectionError:
            return jsonify({"error": "❌ Cannot reach Anthropic API — check your internet connection."}), 503
        except _anthropic.APITimeoutError:
            return jsonify({"error": "❌ Request timed out — the API took too long to respond."}), 504
        raw = msg.content[0].text.strip()
    elif openai_key:
        from openai import OpenAI
        model    = os.environ.get("OPENAI_MODEL", "gpt-4o-mini")
        base_url = os.environ.get("OPENAI_BASE_URL") or None
        client   = OpenAI(api_key=openai_key, base_url=base_url)
        try:
            resp = client.chat.completions.create(
                model=model, max_tokens=2000,
                messages=[{"role": "user", "content": prompt}],
            )
        except Exception as e:
            err = str(e).lower()
            if "401" in err or "authentication" in err:
                return jsonify({"error": "❌ API key invalid — check OPENAI_API_KEY in your .env file."}), 401
            elif "429" in err or "rate limit" in err:
                return jsonify({"error": "❌ Rate limit hit — wait a moment and try again."}), 429
            elif "402" in err or "quota" in err or "billing" in err or "insufficient_quota" in err:
                return jsonify({"error": "❌ API quota exhausted or billing issue — check your account."}), 402
            elif "connection" in err or "timeout" in err:
                return jsonify({"error": "❌ Cannot reach API — check your internet connection."}), 503
            return jsonify({"error": f"❌ API error: {e}"}), 500
        raw = resp.choices[0].message.content.strip()
    else:
        return jsonify({"error": "❌ No API key configured — add ANTHROPIC_API_KEY or OPENAI_API_KEY to your .env file."}), 500

    raw = re.sub(r'^```[a-z]*\n?', '', raw).rstrip('`').strip()
    cards = json.loads(raw)
    return jsonify(cards)

@app.route("/api/all-flashcards")
def api_all_flashcards():
    """Return all flashcards from all courses combined."""
    all_cards = []
    for d in sorted(COURSES_DIR.iterdir()):
        if not d.is_dir():
            continue
        sub_dirs = [sd for sd in d.iterdir() if sd.is_dir() and not sd.name.startswith('.')]
        if len(sub_dirs) >= 2:
            for sd in sorted(sub_dirs):
                all_cards.extend(_collect_flashcards(f"{d.name}/{sd.name}", sd))
        else:
            all_cards.extend(_collect_flashcards(d.name, d))
    return jsonify(all_cards)

@app.route("/api/file-note/<path:course_name>/<path:filename>", methods=["GET", "POST"])
def api_file_note(course_name, filename):
    notes_file = COURSES_DIR / course_name / FILE_NOTES_FILENAME
    notes = json.loads(notes_file.read_text(encoding="utf-8")) if notes_file.exists() else {}
    if request.method == "POST":
        text = request.json.get("text", "")
        if text.strip():
            notes[filename] = text
        elif filename in notes:
            del notes[filename]
        notes_file.write_text(json.dumps(notes, indent=2, ensure_ascii=False), encoding="utf-8")
        return jsonify({"ok": True})
    return jsonify({"text": notes.get(filename, "")})

@app.route("/api/file-note-download/<path:course_name>/<path:filename>")
def api_file_note_download(course_name, filename):
    notes_file = COURSES_DIR / course_name / FILE_NOTES_FILENAME
    notes = json.loads(notes_file.read_text(encoding="utf-8")) if notes_file.exists() else {}
    text = notes.get(filename, "")
    stem = Path(filename).stem
    md_name = f"{stem}_notizen.md"
    return Response(text, mimetype="text/markdown",
                    headers={"Content-Disposition": f"attachment; filename=\"{md_name}\""})

@app.route("/api/notes/<path:course_name>", methods=["GET", "POST"])
def api_notes(course_name):
    path = COURSES_DIR / course_name / NOTES_FILENAME
    if request.method == "POST":
        text = request.json.get("text", "")
        path.write_text(text, encoding="utf-8")
        return jsonify({"ok": True})
    text = path.read_text(encoding="utf-8") if path.exists() else ""
    return jsonify({"text": text})

def _update_streak(prog: dict) -> dict:
    """Update the global learning streak stored under _streak key."""
    today = datetime.now().date().isoformat()
    yesterday = (datetime.now().date() - timedelta(days=1)).isoformat()
    s = prog.get("_streak", {"count": 0, "last_date": None})
    if s["last_date"] == today:
        pass  # already counted today
    elif s["last_date"] == yesterday:
        s["count"] += 1
    else:
        s["count"] = 1
    s["last_date"] = today
    prog["_streak"] = s
    return prog

@app.route("/api/progress/<path:course_name>", methods=["GET", "POST"])
def api_progress(course_name):
    prog = load_progress()
    if request.method == "POST":
        data = request.json
        data["last_studied"] = datetime.now().isoformat(timespec="seconds")
        prog[course_name] = data
        prog = _update_streak(prog)
        save_progress(prog)
        return jsonify({"ok": True})
    return jsonify(prog.get(course_name, {"total": 0, "known": 0, "cards": {}}))

@app.route("/api/streak")
def api_streak():
    prog = load_progress()
    return jsonify(prog.get("_streak", {"count": 0, "last_date": None}))

@app.route("/api/user-state", methods=["GET"])
def api_user_state_get():
    return jsonify(load_user_state())

@app.route("/api/user-state", methods=["POST"])
def api_user_state_post():
    data = request.json or {}
    state = load_user_state()
    state.update(data)
    save_user_state(state)
    return jsonify({"ok": True})

@app.route("/api/progress-reset/<path:course_name>", methods=["POST"])
def api_progress_reset(course_name):
    prog = load_progress()
    if course_name in prog:
        del prog[course_name]
        save_progress(prog)
    return jsonify({"ok": True})

@app.route("/api/summarize", methods=["POST"])
def api_summarize():
    data     = request.json
    course   = data.get("course", "")
    limit    = data.get("limit", 3)
    force    = data.get("force", False)
    files    = data.get("files", [])
    lang     = data.get("lang", "en")
    new_file = data.get("new_file", False)
    length   = data.get("length", "short")

    # Resolve to absolute path so nested courses (e.g. "Archiv/Machine Learning") work
    course_dir = COURSES_DIR / course
    cmd = [PYTHON, SUMMARIZE_SCRIPT, "--dir", str(course_dir), "--limit", str(limit), "--lang", lang, "--length", length]
    if force:
        cmd.append("--force")
    if new_file:
        # Find next available number: _zusammenfassung.md = 1, _zusammenfassung_2.md = 2, …
        existing = list_summaries(course_dir)
        nums = []
        for s in existing:
            if s["name"] == "_zusammenfassung.md":
                nums.append(1)
            else:
                m = re.match(r'^_zusammenfassung_(\d+)\.md$', s["name"])
                if m:
                    nums.append(int(m.group(1)))
        next_num = max(nums, default=0) + 1
        out_name = "_zusammenfassung.md" if next_num == 1 else f"_zusammenfassung_{next_num}.md"
        cmd += ["--out", out_name, "--force"]
    if files:
        cmd += ["--files"] + files

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        return jsonify({"success": result.returncode == 0, "log": result.stdout + result.stderr})
    except subprocess.TimeoutExpired:
        return jsonify({"success": False, "log": "Timeout after 5 minutes."})
    except Exception as e:
        return jsonify({"success": False, "log": str(e)})

@app.route("/api/summarize-all", methods=["POST"])
def api_summarize_all():
    """Summarise all courses that don't have a summary yet."""
    lang  = request.json.get("lang", "en")
    limit = request.json.get("limit", 3)
    log_lines = []
    errors = 0
    courses = get_courses()
    pending = []
    for item in courses:
        if item.get("is_group"):
            for sub in item.get("courses", []):
                if not sub["has_summary"] and sub["file_count"] > 0:
                    pending.append(sub["path"])
        else:
            if not item["has_summary"] and item["file_count"] > 0:
                pending.append(item["path"])

    for path in pending:
        course_dir = COURSES_DIR / path
        cmd = [PYTHON, SUMMARIZE_SCRIPT, "--dir", str(course_dir), "--limit", str(limit), "--lang", lang]
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            log_lines.append(f"✓ {path}" if r.returncode == 0 else f"✗ {path}: {r.stderr[:120]}")
            if r.returncode != 0:
                errors += 1
        except Exception as e:
            log_lines.append(f"✗ {path}: {e}")
            errors += 1

    return jsonify({"success": errors == 0, "done": len(pending) - errors, "count": len(pending), "log": "\n".join(log_lines) or "Nichts zu tun."})

# ---------------------------------------------------------------------------
# Non-blocking scrape / sync — subprocess writes directly to PIPELINE_LOG,
# dashboard holds zero output in RAM. JS polls /api/sync-status.
# ---------------------------------------------------------------------------
_active_proc: "subprocess.Popen | None" = None
_active_log_fh = None  # file handle for PIPELINE_LOG, kept open while proc runs

def _start_proc(cmd: list) -> dict:
    global _active_proc, _active_log_fh
    if _active_proc and _active_proc.poll() is None:
        return {"started": False, "error": "Already running"}
    if _active_log_fh:
        try: _active_log_fh.close()
        except Exception: pass
    _active_log_fh = open(PIPELINE_LOG, "w", encoding="utf-8", buffering=1)
    _active_proc = subprocess.Popen(cmd, stdout=_active_log_fh, stderr=subprocess.STDOUT)
    return {"started": True}

@app.route("/api/sync-status")
def api_sync_status():
    global _active_proc, _active_log_fh
    running = _active_proc is not None and _active_proc.poll() is None
    success = None
    if _active_proc and not running:
        success = _active_proc.returncode == 0
        if _active_log_fh:
            try: _active_log_fh.close()
            except Exception: pass
            _active_log_fh = None
    log = ""
    try:
        log_path = Path(PIPELINE_LOG)
        if log_path.exists():
            log = log_path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        pass
    return jsonify({"running": running, "success": success, "log": log})

@app.route("/api/scrape", methods=["POST"])
def api_scrape():
    return jsonify(_start_proc([PYTHON, SCRAPER_SCRIPT]))

COURSES_JSON = Path(__file__).parent / "courses.json"

@app.route("/api/sync-course", methods=["POST"])
def api_sync_course():
    """Re-sync a single course by its local relative path using courses.json."""
    course_path = (request.json or {}).get("course", "")
    if not course_path:
        return jsonify({"started": False, "error": "No course path provided."})
    if not COURSES_JSON.exists():
        return jsonify({"started": False, "error": "courses.json not found. Run a full scrape first."})
    registry = json.loads(COURSES_JSON.read_text(encoding="utf-8"))
    if course_path not in registry:
        return jsonify({"started": False, "error": f"Course '{course_path}' not in registry. Run a full scrape first."})
    return jsonify(_start_proc([PYTHON, SCRAPER_SCRIPT, "--course", course_path]))

@app.route("/api/course-registry")
def api_course_registry():
    if not COURSES_JSON.exists():
        return jsonify({})
    return jsonify(json.loads(COURSES_JSON.read_text(encoding="utf-8")))

@app.route("/api/course-info/<path:course_name>")
def api_course_info(course_name):
    """Return stored metadata for a course from courses.json."""
    if not COURSES_JSON.exists():
        return jsonify({})
    registry = json.loads(COURSES_JSON.read_text(encoding="utf-8"))
    # Try exact match first, then suffix match (registry keys may differ from dashboard paths)
    entry = registry.get(course_name)
    if not entry:
        for key, val in registry.items():
            if key.endswith("/" + course_name.split("/")[-1]):
                entry = val
                break
    return jsonify(entry.get("meta", {}) if entry else {})

@app.route("/api/custom-info/<path:course_name>")
def api_get_custom_info(course_name):
    path = COURSES_DIR / course_name / CUSTOM_INFO_FILENAME
    if not path.exists():
        return jsonify([])
    return jsonify(json.loads(path.read_text(encoding="utf-8")))

@app.route("/api/custom-info/<path:course_name>", methods=["POST"])
def api_save_custom_info(course_name):
    d = COURSES_DIR / course_name
    if not d.is_dir():
        return jsonify({"success": False})
    fields = request.json  # list of {label, value}
    (d / CUSTOM_INFO_FILENAME).write_text(json.dumps(fields, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"success": True})

# ---------------------------------------------------------------------------
# Background service (macOS LaunchAgent)
# ---------------------------------------------------------------------------
_PLIST_LABEL = "com.uniscraper.dashboard"
_PLIST_PATH  = Path.home() / "Library" / "LaunchAgents" / f"{_PLIST_LABEL}.plist"
_DASHBOARD_SCRIPT = str(Path(__file__).resolve())
_DASHBOARD_DIR    = str(Path(__file__).parent.resolve())
_DASHBOARD_LOG    = str(Path(__file__).parent / "dashboard.log")

def _plist_content() -> str:
    return f"""<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE plist PUBLIC "-//Apple//DTD PLIST 1.0//EN" "http://www.apple.com/DTDs/PropertyList-1.0.dtd">
<plist version="1.0">
<dict>
    <key>Label</key><string>{_PLIST_LABEL}</string>
    <key>ProgramArguments</key>
    <array>
        <string>{PYTHON}</string>
        <string>{_DASHBOARD_SCRIPT}</string>
    </array>
    <key>RunAtLoad</key><true/>
    <key>KeepAlive</key><false/>
    <key>WorkingDirectory</key><string>{_DASHBOARD_DIR}</string>
    <key>StandardOutPath</key><string>{_DASHBOARD_LOG}</string>
    <key>StandardErrorPath</key><string>{_DASHBOARD_LOG}</string>
</dict>
</plist>"""

@app.route("/api/background-status")
def api_background_status():
    enabled = _PLIST_PATH.exists()
    return jsonify({
        "enabled": enabled,
        "plist_path": str(_PLIST_PATH),
        "start_cmd": f"cd '{_DASHBOARD_DIR}' && {PYTHON} dashboard.py",
    })

@app.route("/api/background-enable", methods=["POST"])
def api_background_enable():
    try:
        _PLIST_PATH.parent.mkdir(parents=True, exist_ok=True)
        _PLIST_PATH.write_text(_plist_content(), encoding="utf-8")
        subprocess.run(["launchctl", "load", str(_PLIST_PATH)], check=False)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

@app.route("/api/background-disable", methods=["POST"])
def api_background_disable():
    try:
        if _PLIST_PATH.exists():
            subprocess.run(["launchctl", "unload", str(_PLIST_PATH)], check=False)
            _PLIST_PATH.unlink()
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

@app.route("/api/pipeline-status")
def api_pipeline():
    return jsonify(get_pipeline_status())

def _extract_file_text(path: Path) -> str:
    """Extract plain text from a file."""
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="replace")
        elif suffix == ".pdf":
            import fitz
            with fitz.open(str(path)) as doc:
                return "\n\n".join(page.get_text() for page in doc)
        elif suffix in {".doc", ".docx"}:
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts)
    except Exception:
        pass
    return ""

def _snippets(text: str, q: str, max_matches: int = 3) -> list[str]:
    matches, start, lower = [], 0, text.lower()
    while True:
        idx = lower.find(q, start)
        if idx == -1 or len(matches) >= max_matches:
            break
        snippet = text[max(0, idx-80):idx+120].replace('\n', ' ')
        matches.append(re.sub(r'\s+', ' ', snippet).strip())
        start = idx + 1
    return matches

def _search_dir(q: str, rel_path: str, d: Path, results: list):
    """Search summary and notes in a single course directory."""
    hits = []
    source = None
    # Search summary first
    summary = get_latest_summary(d)
    if summary and summary.exists():
        md = summary.read_text(encoding="utf-8", errors="replace")
        snips = _snippets(md, q)
        if snips:
            hits = snips
            source = "summary"
    # Also search notes
    notes = d / NOTES_FILENAME
    if notes.exists():
        nd = notes.read_text(encoding="utf-8", errors="replace")
        nsnips = _snippets(nd, q)
        if nsnips and not hits:
            hits = nsnips
            source = "notes"
        elif nsnips:
            source = "both"
    if hits:
        results.append({
            "course":  rel_path,
            "name":    d.name,
            "snippet": hits[0],
            "count":   len(hits),
            "source":  source,
        })

@app.route("/api/search")
def api_search():
    q = request.args.get("q", "").lower().strip()
    if not q:
        return jsonify([])
    results = []
    for d in sorted(COURSES_DIR.iterdir()):
        if not d.is_dir():
            continue
        sub_dirs = [sd for sd in d.iterdir() if sd.is_dir() and not sd.name.startswith('.')]
        if len(sub_dirs) >= 2:
            for sd in sorted(sub_dirs):
                _search_dir(q, f"{d.name}/{sd.name}", sd, results)
        else:
            _search_dir(q, d.name, d, results)
    return jsonify(results)

@app.route("/api/search-files")
def api_search_files():
    """Full-text search inside actual course files (PDFs, DOCX, TXT, MD)."""
    q = request.args.get("q", "").lower().strip()
    if len(q) < 2:
        return jsonify([])
    results = []
    for d in sorted(COURSES_DIR.iterdir()):
        if not d.is_dir():
            continue
        sub_dirs = [sd for sd in d.iterdir() if sd.is_dir() and not sd.name.startswith('.')]
        dirs = [(f"{d.name}/{sd.name}", sd) for sd in sorted(sub_dirs)] if len(sub_dirs) >= 2 else [(d.name, d)]
        for rel, course_dir in dirs:
            files = [
                f for f in course_dir.rglob("*")
                if f.is_file() and f.suffix.lower() in {".pdf", ".docx", ".txt", ".md"}
                and not SUMMARY_RE.match(f.name) and f.name != NOTES_FILENAME
            ]
            for f in files[:20]:  # cap per course to avoid timeouts
                text = _extract_file_text(f)
                if not text:
                    continue
                snips = _snippets(text, q, max_matches=2)
                if snips:
                    results.append({
                        "course": rel, "name": course_dir.name,
                        "file": f.name, "snippet": snips[0], "count": len(snips),
                        "source": "file",
                    })
    return jsonify(results)

# ---------------------------------------------------------------------------
# Frontend
# ---------------------------------------------------------------------------

HTML = r"""<!DOCTYPE html>
<html lang="de">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Stud.IP Dashboard</title>
<style>
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

:root {
  --bg:        #0e0e11;
  --bg2:       #16161a;
  --bg3:       #1e1e24;
  --bg4:       #27272f;
  --bg5:       #32323c;
  --border:    #252530;
  --border2:   #35353f;
  --text:      #eaeaf2;
  --text2:     #8888a8;
  --text3:     #48486a;
  --blue:      #5b8ef0;
  --blue2:     #3b6ed8;
  --blue3:     #2452a8;
  --green:     #34d399;
  --yellow:    #fbbf24;
  --red:       #f87171;
  --purple:    #a78bfa;
  --orange:    #fb923c;
  --radius:    8px;
  --radius-lg: 12px;
  --radius-xl: 16px;
  --shadow-sm: 0 1px 3px rgba(0,0,0,.5);
  --shadow:    0 4px 14px rgba(0,0,0,.5);
  --shadow-lg: 0 8px 28px rgba(0,0,0,.6);
  --transition: 140ms ease;
  --glow-blue: 0 0 0 3px rgba(91,142,240,.2);
}

:root.light {
  --bg:        #f5f5f7;
  --bg2:       #ffffff;
  --bg3:       #ebebef;
  --bg4:       #dfdfe8;
  --bg5:       #d2d2dc;
  --border:    #e2e2ea;
  --border2:   #ccccd8;
  --text:      #111116;
  --text2:     #44445a;
  --text3:     #7878a0;
  --blue:      #3b6ed8;
  --blue2:     #2456b8;
  --blue3:     #1a3e90;
  --green:     #059669;
  --yellow:    #d97706;
  --red:       #dc2626;
  --purple:    #6d28d9;
  --orange:    #ea580c;
  --shadow-sm: 0 1px 2px rgba(0,0,0,.07);
  --shadow:    0 4px 12px rgba(0,0,0,.09);
  --shadow-lg: 0 8px 24px rgba(0,0,0,.12);
  --glow-blue: 0 0 0 3px rgba(59,110,216,.18);
}
:root.light .preview-body.pdf-wrap { background: #888; }
:root.light .dot-ok { box-shadow: none; }

body {
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Helvetica, sans-serif;
  background: var(--bg);
  color: var(--text);
  height: 100vh;
  display: flex;
  flex-direction: column;
  overflow: hidden;
  font-size: 14px;
  -webkit-font-smoothing: antialiased;
}

/* ── Topbar ── */
#topbar {
  height: 46px;
  background: var(--bg2);
  border-bottom: 1px solid var(--border);
  display: flex;
  align-items: center;
  padding: 0 18px;
  gap: 10px;
  flex-shrink: 0;
  z-index: 10;
}
#topbar-logo {
  font-size: 14px; font-weight: 700; color: var(--text2);
  display: flex; align-items: center; gap: 7px; letter-spacing: -.01em;
}
#topbar-logo .logo-icon { font-size: 16px; }
#topbar-logo span { color: var(--blue); }

#search-global {
  flex: 1; max-width: 340px;
  padding: 8px 14px 8px 36px;
  background: var(--bg3);
  border: 1px solid var(--border);
  border-radius: 22px;
  color: var(--text);
  font-size: 13px;
  outline: none;
  transition: border-color var(--transition), box-shadow var(--transition);
  background-image: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='14' height='14' viewBox='0 0 24 24' fill='none' stroke='%2349597a' stroke-width='2.5'%3E%3Ccircle cx='11' cy='11' r='8'/%3E%3Cpath d='m21 21-4.35-4.35'/%3E%3C/svg%3E");
  background-repeat: no-repeat;
  background-position: 12px center;
}
#search-global:focus { border-color: var(--blue); box-shadow: var(--glow-blue); }
#search-global::placeholder { color: var(--text3); }

/* ── Topbar responsiveness ── */
@media (max-width: 980px) {
  #topbar { gap: 8px; padding: 0 12px; }
  #search-global { max-width: 240px; }
}
@media (max-width: 780px) {
  .logo-wordmark { display: none; }
  #search-global { max-width: none; flex: 1; }
}
@media (max-width: 600px) {
  .tbtn-label { display: none; }
  .topbar-shortcuts { display: none; }
  .tbtn { padding: 7px 10px; }
  #search-global { font-size: 12px; padding: 7px 10px 7px 30px; }
}

.tbtn {
  padding: 6px 14px;
  border-radius: var(--radius);
  border: none;
  cursor: pointer;
  font-size: 12.5px;
  font-weight: 600;
  transition: opacity var(--transition), background var(--transition), box-shadow var(--transition);
  white-space: nowrap;
  letter-spacing: -.01em;
}
.tbtn:hover:not(:disabled) { opacity: .85; }
.tbtn:active:not(:disabled) { opacity: .7; }
.tbtn:disabled { opacity: .3; cursor: not-allowed; }
.btn-blue   { background: var(--blue); color: #fff; }
.btn-gray   { background: var(--bg4); color: var(--text2); border: 1px solid var(--border2); }
.btn-green  { background: #059669; color: #d1fae5; }
.btn-red    { background: #dc2626; color: #fee2e2; }
.btn-purple { background: #6d28d9; color: #ede9fe; }

/* ── Layout ── */
#layout { display: flex; flex: 1; overflow: hidden; }

/* ── Sidebar toggle button ── */
#sidebar-toggle {
  width: 30px; height: 30px; border-radius: var(--radius);
  background: none; border: 1px solid var(--border);
  color: var(--text3); cursor: pointer; font-size: 14px; line-height: 1;
  display: flex; align-items: center; justify-content: center;
  transition: background var(--transition), color var(--transition), border-color var(--transition);
  flex-shrink: 0;
}
#sidebar-toggle:hover { background: var(--bg3); color: var(--text2); border-color: var(--border2); }

/* ── Sidebar ── */
#sidebar {
  width: 248px; min-width: 150px; max-width: 480px;
  background: var(--bg);
  border-right: 1px solid var(--border);
  transition: width var(--transition), min-width var(--transition), padding var(--transition);
  display: flex; flex-direction: column;
  overflow: hidden;
  flex-shrink: 0;
}

#sidebar-top {
  padding: 10px 10px 8px;
  border-bottom: 1px solid var(--border);
  display: flex; flex-direction: column; gap: 7px;
  background: var(--bg);
}
#sidebar-search {
  width: 100%; padding: 7px 10px;
  background: var(--bg3); border: 1px solid transparent;
  border-radius: var(--radius); color: var(--text); font-size: 12px; outline: none;
  transition: border-color var(--transition);
}
#sidebar-search:focus { border-color: var(--border2); }
#sidebar-search::placeholder { color: var(--text3); }

/* Filter pills */
.filter-pills { display: flex; gap: 3px; flex-wrap: wrap; }
.pill {
  padding: 3px 9px; border-radius: 6px; font-size: 11px; font-weight: 500;
  cursor: pointer; border: none; background: none;
  color: var(--text3); transition: all var(--transition);
}
.pill:hover { color: var(--text2); background: var(--bg3); }
.pill.active { background: var(--bg4); color: var(--text); }

/* Sort row */
.sort-row { display: flex; align-items: center; gap: 6px; }
.sort-row label { font-size: 10px; color: var(--text3); flex-shrink: 0; }
.sort-select {
  flex: 1; font-size: 11px; background: transparent; border: none;
  color: var(--text3); padding: 3px 0; outline: none; cursor: pointer;
}

.sidebar-section-label {
  padding: 10px 12px 4px;
  font-size: 10px; font-weight: 600;
  color: var(--text3); letter-spacing: .04em;
}

#course-list { overflow-y: auto; flex: 1; padding: 4px 6px 4px; }

.citem {
  padding: 8px 10px; border-radius: var(--radius);
  cursor: pointer; display: flex; align-items: center; gap: 8px;
  margin-bottom: 1px;
  transition: background var(--transition);
  position: relative;
}
.citem:hover  { background: var(--bg3); }
.citem.active { background: var(--bg3); }
.citem.active .citem-name { color: var(--text); }

.citem-dot { width: 5px; height: 5px; border-radius: 50%; flex-shrink: 0; }
.dot-ok     { background: var(--green); }
.dot-missing{ background: var(--bg4); }

.citem-body { flex: 1; min-width: 0; }
.citem-name { font-size: 12px; font-weight: 500; color: var(--text2); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; transition: color var(--transition); }
.citem.active .citem-name { color: var(--text); }
.citem-meta { font-size: 10px; color: var(--text3); display: flex; gap: 6px; margin-top: 2px; flex-wrap: wrap; }

.progress-mini { height: 2px; background: var(--border); border-radius: 2px; margin-top: 5px; }
.progress-mini-fill {
  height: 100%; border-radius: 2px; transition: width .4s ease;
  background: linear-gradient(90deg, var(--blue), var(--purple));
}

/* Favorite star button */
.fav-btn {
  background: none; border: none; cursor: pointer;
  font-size: 12px; padding: 3px 5px; flex-shrink: 0;
  opacity: 0.25; transition: opacity var(--transition), transform var(--transition);
  line-height: 1;
}
.fav-btn:hover { opacity: .8; transform: scale(1.25); }
.fav-btn.is-fav { opacity: 1; }

/* Context menu */
#ctx-menu {
  position: fixed; z-index: 9999;
  background: var(--bg2); border: 1px solid var(--border2);
  border-radius: var(--radius); padding: 4px 0;
  box-shadow: 0 4px 16px rgba(0,0,0,.35);
  min-width: 160px; display: none;
}
#ctx-menu.visible { display: block; }
.ctx-item {
  padding: 7px 14px; font-size: 12px; color: var(--text2);
  cursor: pointer; white-space: nowrap;
  transition: background var(--transition);
}
.ctx-item:hover { background: var(--bg3); color: var(--text); }
.ctx-item.danger { color: #f87171; }
.ctx-item.danger:hover { background: rgba(248,113,113,.12); }

/* Archive section header */
.archive-header {
  display: flex; align-items: center; gap: 7px;
  padding: 7px 10px 7px; cursor: pointer;
  border-radius: var(--radius);
  transition: background var(--transition);
  margin: 8px 0 1px; user-select: none;
  border-left: 2px solid transparent;
  font-size: 10px; color: var(--text3);
  font-weight: 600; text-transform: uppercase; letter-spacing: .05em;
}
.archive-header:hover { background: var(--bg3); color: var(--text2); }
.citem.archived { opacity: 0.55; }

/* Group headers in sidebar */
.group-header {
  display: flex; align-items: center; gap: 7px; padding: 6px 10px;
  cursor: pointer; border-radius: var(--radius);
  transition: background var(--transition);
  margin: 6px 0 1px; user-select: none;
}
.group-header:hover { background: var(--bg3); }
.group-chevron { font-size: 7px; color: var(--text3); width: 9px; flex-shrink: 0; }
.group-name { font-size: 10px; font-weight: 600; color: var(--text3); letter-spacing: .04em; flex: 1; }
.group-meta { font-size: 10px; color: var(--text3); }

/* Semester group header */
.semester-header { margin-top: 10px; }
.semester-header .group-name { color: var(--text2); }

/* Global learn button at sidebar bottom */
#global-learn-btn {
  margin: 6px 8px 10px; flex-shrink: 0;
  padding: 8px 12px; font-size: 12px; font-weight: 600;
  background: var(--bg3);
  border: 1px solid var(--border); color: var(--text3);
  border-radius: var(--radius); cursor: pointer; text-align: center;
  transition: all var(--transition);
}
#global-learn-btn:hover {
  color: var(--text2); border-color: var(--border2); background: var(--bg4);
}

/* Credits footer at sidebar bottom */
#sidebar-credits {
  flex-shrink: 0; padding: 6px 12px 10px;
  font-size: 10px; color: var(--text3); text-align: center;
  opacity: 0.45; letter-spacing: 0.01em; line-height: 1.5;
}
#sidebar-credits a {
  color: var(--text3); text-decoration: none; font-weight: 500;
  transition: color var(--transition);
}
#sidebar-credits a:hover { color: var(--blue); }
#sidebar.collapsed #sidebar-credits { display: none; }

/* ── Resize divider ── */
.resize-divider {
  width: 5px; background: transparent; cursor: col-resize;
  flex-shrink: 0; position: relative; z-index: 5; transition: background var(--transition);
}
.resize-divider:hover, .resize-divider.dragging { background: var(--blue); opacity: .6; }
/* Subtle grip dots to aid discoverability */
.resize-divider::after {
  content: ''; position: absolute; top: 50%; left: 50%;
  transform: translate(-50%, -50%);
  width: 3px; height: 28px; border-radius: 3px;
  background: var(--border2); opacity: 0;
  transition: opacity var(--transition);
}
.resize-divider:hover::after, .resize-divider.dragging::after { opacity: 1; }

/* Collapsed sidebar */
#sidebar.collapsed {
  width: 0 !important; min-width: 0 !important; overflow: hidden;
  border-right: none;
}
#sidebar.collapsed + .resize-divider { display: none; }

/* ── Mobile sidebar overlay ── */
#sidebar-backdrop {
  display: none; position: fixed; inset: 0; top: 46px;
  background: rgba(0,0,0,.55); z-index: 199;
  opacity: 0; pointer-events: none;
  transition: opacity .25s;
}
#sidebar-backdrop.visible { opacity: 1; pointer-events: auto; }

@media (max-width: 768px) {
  #sidebar {
    position: fixed; top: 46px; left: 0; bottom: 0; z-index: 200;
    width: 280px !important; min-width: 280px !important;
    transform: translateX(-100%);
    transition: transform .25s ease, box-shadow .25s;
    border-right: 1px solid var(--border);
  }
  #sidebar.mobile-open {
    transform: translateX(0);
    box-shadow: 4px 0 24px rgba(0,0,0,.6);
  }
  #sidebar.collapsed { transform: translateX(-100%); width: 280px !important; }
  .resize-divider#divider-sidebar { display: none; }
  #sidebar-backdrop { display: block; }
}

/* ── Main ── */
#main { flex: 1; display: flex; flex-direction: column; overflow: hidden; min-width: 0; }

/* ── Tabs ── */
#tabs {
  display: flex; align-items: stretch;
  padding: 0 14px; gap: 0;
  background: var(--bg2); border-bottom: 1px solid var(--border);
  height: 40px; flex-shrink: 0;
}
.tab {
  padding: 0 12px; border-radius: 0; cursor: pointer;
  font-size: 12.5px; font-weight: 500; color: var(--text3);
  transition: color var(--transition); border: none; background: none;
  display: flex; align-items: center; gap: 5px;
  position: relative;
}
.tab:hover { color: var(--text2); }
.tab.active { color: var(--text); }
.tab.active::after {
  content: ''; position: absolute; bottom: 0; left: 8px; right: 8px;
  height: 2px; background: var(--blue); border-radius: 2px 2px 0 0;
}
.tab-spacer { flex: 1; }
#tabs-course-label {
  font-size: 11px; color: var(--text3); font-weight: 400;
  max-width: 200px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap;
  padding: 0 4px;
}

/* ── Course title bar ── */
#course-title-bar {
  display: flex; align-items: center; gap: 10px;
  padding: 10px 20px 6px;
  background: var(--bg2); border-bottom: 1px solid var(--border);
  flex-shrink: 0;
}
#course-title-text {
  font-size: 15px; font-weight: 700; color: var(--text);
  overflow: hidden; text-overflow: ellipsis; white-space: nowrap;
  flex: 1;
}

/* ── Content ── */
#content { flex: 1; overflow: hidden; position: relative; }

/* IMPORTANT: panels use display:none by default — never add display:flex/block to panel IDs */
.panel { position: absolute; inset: 0; overflow-y: auto; display: none; padding: 28px 36px; }
.panel.active { display: block; animation: panelIn .18s ease; }
@keyframes panelIn { from { opacity: 0; transform: translateY(4px); } to { opacity: 1; transform: none; } }
/* Special panels that need flex layout: use inner wrapper with class .panel-inner-flex */
.panel-inner-flex { display: flex; flex-direction: column; height: 100%; }

/* Home panel */
.stats-grid { display: grid; grid-template-columns: repeat(2, 1fr); gap: 10px; margin-bottom: 20px; }
.stat-card {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 16px 18px;
  transition: border-color var(--transition);
}
.stat-card:hover { border-color: var(--border2); }
.stat-label { font-size: 10px; color: var(--text3); letter-spacing: .04em; margin-bottom: 8px; font-weight: 600; }
.stat-value { font-size: 28px; font-weight: 700; color: var(--text); letter-spacing: -.02em; line-height: 1; }
.stat-sub   { font-size: 11px; color: var(--text3); margin-top: 5px; }
.stat-bar { height: 2px; background: var(--border); border-radius: 2px; margin-top: 10px; overflow: hidden; }
.stat-bar-fill { height: 100%; border-radius: 2px; transition: width .6s cubic-bezier(.4,0,.2,1); }

.section-title {
  font-size: 11px; font-weight: 600; color: var(--text3); margin-bottom: 12px;
  letter-spacing: .02em;
}

.pipeline-card {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 12px 18px; display: flex; align-items: center; gap: 12px; margin-bottom: 20px;
}
.pipeline-dot { width: 9px; height: 9px; border-radius: 50%; flex-shrink: 0; }

/* Background service card */
.bg-service-card {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 16px 18px; margin-top: 24px;
}
.bg-service-header { display: flex; align-items: center; gap: 10px; margin-bottom: 12px; }
.bg-service-title { font-size: 12px; font-weight: 600; color: var(--text2); flex: 1; }
.bg-toggle { position: relative; width: 36px; height: 20px; flex-shrink: 0; }
.bg-toggle input { opacity: 0; width: 0; height: 0; }
.bg-toggle-slider {
  position: absolute; inset: 0; background: var(--bg4); border-radius: 20px;
  cursor: pointer; transition: background .2s;
}
.bg-toggle input:checked + .bg-toggle-slider { background: var(--blue); }
.bg-toggle-slider::before {
  content: ''; position: absolute; width: 14px; height: 14px; border-radius: 50%;
  background: #fff; left: 3px; top: 3px; transition: transform .2s;
}
.bg-toggle input:checked + .bg-toggle-slider::before { transform: translateX(16px); }
.bg-instructions {
  font-size: 11px; color: var(--text3); line-height: 1.7;
  border-top: 1px solid var(--border); padding-top: 10px; margin-top: 4px;
}
.bg-instructions code {
  background: var(--bg3); border: 1px solid var(--border); border-radius: 4px;
  padding: 2px 7px; font-size: 11px; color: var(--text2); font-family: "SF Mono", monospace;
  cursor: pointer; user-select: all;
}
.bg-status-dot { width: 7px; height: 7px; border-radius: 50%; flex-shrink: 0; display: inline-block; margin-right: 4px; }

/* Sticky notes pinboard */
.pinboard {
  display: flex; flex-wrap: wrap; gap: 12px; margin-bottom: 24px;
}
.sticky {
  width: 160px; min-height: 120px; border-radius: 3px; padding: 10px;
  display: flex; flex-direction: column; gap: 6px; position: relative;
  box-shadow: 2px 3px 8px rgba(0,0,0,.35);
  transition: box-shadow var(--transition), transform var(--transition);
  cursor: grab;
}
.sticky:active { cursor: grabbing; }
.sticky textarea { cursor: text; }
.sticky:hover { box-shadow: 3px 5px 14px rgba(0,0,0,.5); transform: translateY(-1px); }
.sticky textarea {
  flex: 1; background: transparent; border: none; outline: none;
  font-size: 12.5px; line-height: 1.6; resize: none; color: #1a1a1a;
  font-family: inherit; min-height: 80px;
}
.sticky textarea::placeholder { color: rgba(0,0,0,.35); }
.sticky-del {
  position: absolute; top: 5px; right: 6px;
  background: none; border: none; cursor: pointer;
  font-size: 12px; color: rgba(0,0,0,.3); padding: 0; line-height: 1;
  transition: color .15s;
}
.sticky-del:hover { color: rgba(0,0,0,.7); }
.pinboard-add {
  width: 160px; min-height: 120px; border-radius: 3px;
  border: 2px dashed var(--border2); display: flex; align-items: center;
  justify-content: center; cursor: pointer; color: var(--text3); font-size: 22px;
  transition: border-color var(--transition), color var(--transition);
}
.pinboard-add:hover { border-color: var(--text3); color: var(--text2); }

/* To-do widget */
.todo-widget {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 14px 18px; margin-bottom: 24px; max-width: 700px;
}
.todo-widget-header {
  display: flex; align-items: center; gap: 8px; margin-bottom: 10px;
}
.todo-widget-title { font-size: 11px; font-weight: 600; color: var(--text3); letter-spacing: .04em; flex: 1; }
.todo-clear-btn {
  background: none; border: none; color: var(--text3); cursor: pointer;
  font-size: 11px; padding: 2px 6px; border-radius: 4px;
  transition: color var(--transition), background var(--transition);
}
.todo-clear-btn:hover { color: var(--text2); background: var(--bg3); }
.todo-list { display: flex; flex-direction: column; gap: 1px; margin-bottom: 8px; }
.todo-item {
  display: flex; align-items: flex-start; gap: 8px; padding: 4px 2px;
  border-radius: 5px; position: relative;
}
.todo-item.drag-over { outline: 2px dashed var(--blue); outline-offset: -1px; border-radius: 5px; }
.todo-drag-handle {
  color: var(--text3); font-size: 11px; cursor: grab; flex-shrink: 0; padding-top: 4px;
  opacity: 0; transition: opacity .15s; user-select: none; line-height: 1;
}
.todo-item:hover .todo-drag-handle { opacity: 1; }
.todo-drag-handle:active { cursor: grabbing; }
.todo-cb {
  appearance: none; -webkit-appearance: none;
  width: 14px; height: 14px; border: 1.5px solid var(--text3);
  border-radius: 3px; cursor: pointer; flex-shrink: 0; margin-top: 3px;
  background: transparent; transition: background var(--transition), border-color var(--transition);
}
.todo-cb:checked {
  background: var(--blue); border-color: var(--blue);
  background-image: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 10 10'%3E%3Cpath d='M1.5 5l2.5 2.5 4.5-4.5' stroke='white' stroke-width='1.5' fill='none' stroke-linecap='round' stroke-linejoin='round'/%3E%3C/svg%3E");
  background-size: 10px; background-repeat: no-repeat; background-position: center;
}
.todo-text {
  flex: 1; background: none; border: none; outline: none;
  color: var(--text2); font-size: 13px; line-height: 1.5;
  font-family: inherit; resize: none; padding: 0; min-height: 20px;
  transition: color var(--transition);
}
.todo-text.done { color: var(--text3); text-decoration: line-through; }
.todo-text::placeholder { color: var(--text3); }
.todo-del-btn {
  background: none; border: none; color: transparent; cursor: pointer;
  font-size: 12px; padding: 0 3px; border-radius: 3px; line-height: 1; flex-shrink: 0;
  transition: color var(--transition);
}
.todo-item:hover .todo-del-btn { color: var(--text3); }
.todo-del-btn:hover { color: var(--red) !important; }
.todo-add-btn {
  background: none; border: none; color: var(--text3); cursor: pointer;
  font-size: 12px; padding: 3px 0; display: flex; align-items: center; gap: 5px;
  transition: color var(--transition);
}
.todo-add-btn:hover { color: var(--text2); }
.todo-link-btn {
  background: none; border: 1px solid var(--border); color: var(--text3); cursor: pointer;
  font-size: 11px; padding: 1px 5px; border-radius: 4px; flex-shrink: 0; line-height: 1.4;
  transition: color var(--transition), border-color var(--transition);
}
.todo-link-btn:hover { color: var(--blue); border-color: var(--blue); }
.todo-due {
  font-size: 10px; padding: 1px 6px; border-radius: 10px; white-space: nowrap;
  flex-shrink: 0; cursor: pointer; border: none; font-family: inherit; line-height: 1.6;
  transition: opacity .15s;
}
.todo-due.overdue  { background: rgba(248,113,113,.18); color: var(--red); }
.todo-due.today    { background: rgba(251,191,36,.18);  color: var(--yellow); }
.todo-due.upcoming { background: rgba(91,142,240,.12);  color: var(--blue); }
.todo-due:hover { opacity: .75; }
.todo-sort-btn {
  background: none; border: none; cursor: pointer; font-size: 11px; color: var(--text3);
  padding: 1px 5px; border-radius: 4px; transition: color .15s;
}
.todo-sort-btn:hover { color: var(--text2); }
.todo-sort-btn.active { color: var(--blue); }
.todo-course-pill {
  font-size: 10px; background: rgba(79,142,247,.15); color: var(--blue);
  border-radius: 10px; padding: 2px 7px; white-space: nowrap; flex-shrink: 0;
  cursor: pointer; transition: background var(--transition);
}
.todo-course-pill:hover { background: rgba(79,142,247,.3); }
.todo-course-picker {
  position: absolute; z-index: 200; background: var(--bg2); border: 1px solid var(--border2);
  border-radius: var(--radius-lg); box-shadow: var(--shadow-lg); min-width: 200px; max-width: 280px;
  right: 0; top: 100%; margin-top: 4px;
}
.tcp-search-wrap { padding: 8px 10px 4px; }
.tcp-search {
  width: 100%; background: var(--bg3); border: 1px solid var(--border); border-radius: 5px;
  color: var(--text2); font-size: 12px; padding: 4px 8px; outline: none; box-sizing: border-box;
}
.tcp-list { max-height: 200px; overflow-y: auto; padding: 4px 0 6px; }
.tcp-item {
  padding: 6px 12px; font-size: 12px; color: var(--text2); cursor: pointer;
  transition: background var(--transition);
}
.tcp-item:hover, .tcp-item.tcp-active { background: var(--bg3); }
/* Course-level todo widget */
.course-todos-wrap {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 10px 14px; margin-bottom: 16px; display: none;
}
.course-todos-header {
  font-size: 11px; font-weight: 600; color: var(--text3); letter-spacing: .04em; margin-bottom: 8px;
  display: flex; align-items: center; gap: 8px;
}
.course-todos-header a {
  font-size: 10px; color: var(--blue); cursor: pointer; margin-left: auto; text-decoration: none;
}
.course-todos-header a:hover { text-decoration: underline; }

.recent-list { display: flex; flex-direction: column; gap: 2px; }
.recent-item {
  background: transparent; border-radius: var(--radius);
  padding: 9px 12px; cursor: pointer; display: flex; align-items: center; gap: 12px;
  transition: background var(--transition);
}
.recent-item:hover { background: var(--bg3); }
.recent-item-name { flex: 1; font-size: 13px; color: var(--text); font-weight: 500; }
.recent-item-meta { font-size: 11px; color: var(--text3); }
.recent-item-progress { min-width: 80px; }
.recent-progress-bar { height: 3px; background: var(--border); border-radius: 3px; overflow: hidden; margin-top: 4px; }
.recent-progress-fill { height: 100%; background: linear-gradient(90deg, var(--blue), var(--purple)); border-radius: 3px; }

/* Files panel */
.files-layout { display: flex; height: 100%; overflow: hidden; gap: 0; }
.files-list-col { width: 260px; min-width: 140px; max-width: 500px; flex-shrink: 0; overflow-y: auto; padding: 0 2px 0 0; transition: width 180ms ease, min-width 180ms ease; }
.files-list-col.collapsed { width: 0 !important; min-width: 0 !important; overflow: hidden; }
.files-list-col.collapsed ~ #divider-files { display: none; }
#filelist-toggle-tab {
  width: 18px; flex-shrink: 0; display: flex; align-items: flex-start; justify-content: center;
  padding-top: 6px; background: none; border: none; cursor: pointer;
  color: var(--text3); font-size: 11px; transition: color var(--transition);
}
#filelist-toggle-tab:hover { color: var(--text2); }
.files-preview-col { flex: 1; min-width: 0; overflow: hidden; }

.file-item {
  display: flex; align-items: center; gap: 8px; padding: 6px 8px;
  border-radius: var(--radius); cursor: pointer;
  transition: background var(--transition), border-color var(--transition);
  border: 1px solid transparent; margin-bottom: 1px;
}
.file-item:hover { background: var(--bg3); }
.file-item.active { background: var(--bg4); border-color: rgba(79,142,247,.4); }
.file-item.new-file { border-color: rgba(251,191,36,.35); background: rgba(251,191,36,.06); }
.file-item input[type=checkbox] { accent-color: var(--blue); flex-shrink: 0; cursor: pointer; display: none; }
#file-list.selection-mode .file-item input[type=checkbox] { display: inline-block; }
.file-icon { font-size: 14px; flex-shrink: 0; }
.file-name { font-size: 12px; color: var(--text2); flex: 1; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.file-date { font-size: 10px; color: var(--text3); flex-shrink: 0; }
.file-read-check { font-size: 11px; color: var(--green, #4ade80); flex-shrink: 0; opacity: 0.85; }
.file-item.file-read .file-name { color: var(--text3); }
.read-toggle-btn {
  font-size: 11px; padding: 2px 8px; border-radius: 5px; border: 1px solid var(--border);
  background: none; color: var(--text3); cursor: pointer; flex-shrink: 0; transition: all .15s;
  white-space: nowrap;
}
.read-toggle-btn.is-read { color: var(--green, #4ade80); border-color: var(--green, #4ade80); }
.preview-header { position: relative; }

/* Folder tree */
.folder-item {
  display: flex; align-items: center; gap: 6px;
  padding: 5px 8px; border-radius: var(--radius);
  cursor: pointer; user-select: none;
  transition: background var(--transition);
  margin-bottom: 1px;
}
.folder-item:hover { background: var(--bg3); }
.folder-chevron {
  font-size: 8px; color: var(--text3); width: 10px; flex-shrink: 0;
  transition: transform var(--transition);
}
.folder-item.collapsed .folder-chevron { transform: rotate(-90deg); }
.folder-icon { font-size: 13px; flex-shrink: 0; }
.folder-name { font-size: 11px; font-weight: 600; color: var(--text3); flex: 1;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.folder-count { font-size: 10px; color: var(--text3); background: var(--bg3);
  padding: 1px 5px; border-radius: 8px; flex-shrink: 0; }
.folder-contents { overflow: hidden; }
.folder-contents.collapsed { display: none; }
.folder-item.has-new .folder-name { color: var(--yellow); }
.folder-new-badge { font-size: 9px; background: rgba(234,179,8,.18); color: var(--yellow);
  border: 1px solid rgba(234,179,8,.4); border-radius: 8px; padding: 1px 5px; flex-shrink: 0; }

.file-actions {
  margin-top: 14px; padding-top: 12px; border-top: 1px solid var(--border);
  display: flex; flex-direction: column; gap: 7px;
}
.limit-row { display: flex; align-items: center; gap: 8px; font-size: 12px; color: var(--text3); }
.limit-input {
  width: 55px; padding: 5px 8px; background: var(--bg3);
  border: 1px solid var(--border); border-radius: var(--radius);
  color: var(--text); font-size: 12px; text-align: center; outline: none;
  transition: border-color var(--transition);
}
.limit-input:focus { border-color: var(--blue); }

.preview-box {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  height: 100%; overflow: hidden; display: flex; flex-direction: column;
  box-shadow: var(--shadow-sm);
}
.preview-header {
  padding: 10px 16px; border-bottom: 1px solid var(--border);
  font-size: 12px; color: var(--text3); background: var(--bg3);
  border-radius: var(--radius-lg) var(--radius-lg) 0 0; flex-shrink: 0;
  display: flex; align-items: center; gap: 8px;
}
.preview-header-name { flex: 1; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; font-weight: 500; color: var(--text2); }
.preview-body {
  flex: 1; overflow-y: auto; padding: 20px;
  font-family: "SF Mono", "JetBrains Mono", monospace; font-size: 12px;
  color: var(--text2); line-height: 1.75; white-space: pre-wrap; word-break: break-word;
}
.preview-body.pdf-wrap {
  padding: 0; overflow-y: auto; overflow-x: hidden;
  background: #444; display: flex; flex-direction: column; align-items: center; gap: 8px;
}
.pdf-page-canvas {
  display: block; box-shadow: 0 2px 8px rgba(0,0,0,.6);
}
/* PDF text layer — transparent overlay for text selection */
.pdf-page-wrapper {
  position: relative; display: inline-block;
  box-shadow: 0 2px 8px rgba(0,0,0,.6);
}
.pdf-page-wrapper .pdf-page-canvas { box-shadow: none; }
.textLayer {
  position: absolute; left: 0; top: 0; right: 0; bottom: 0;
  overflow: hidden; line-height: 1;
  user-select: text; pointer-events: auto;
}
.textLayer span, .textLayer br {
  color: transparent; position: absolute; white-space: pre;
  cursor: text; transform-origin: 0% 0%;
}
.textLayer ::selection { background: rgba(58,130,246,.35); color: transparent; }
/* PDF find highlights */
.pdf-hl     { background: rgba(255,213,0,.55) !important; border-radius: 2px; }
.pdf-hl-cur { background: rgba(255,110,0,.75) !important; }
/* PDF in-page find bar */
#pdf-find-bar {
  display: none; align-items: center; gap: 6px; flex-shrink: 0;
  padding: 6px 12px; border-bottom: 1px solid var(--border);
  background: var(--bg3);
}
#pdf-find-input {
  flex: 1; max-width: 200px;
  padding: 5px 10px; border-radius: 6px;
  border: 1px solid var(--border); background: var(--bg2);
  color: var(--text); font-size: 12px; outline: none;
  transition: border-color var(--transition);
}
#pdf-find-input:focus { border-color: var(--blue); }
#pdf-find-count { font-size: 11px; min-width: 54px; }

/* Fake fullscreen — covers entire viewport, includes notes panel */
#preview-area {
  display: contents; /* no layout effect when not fullscreen */
}
#preview-area.fullscreen {
  display: flex;
  position: fixed; inset: 0; z-index: 500;
  background: var(--bg);
}
#preview-area.fullscreen .files-preview-col {
  flex: 1; min-width: 0; overflow: hidden; display: flex; flex-direction: column;
}
#preview-area.fullscreen .preview-box {
  flex: 1; border-radius: 0; border: none; height: 100%;
}
#preview-area.fullscreen .preview-body {
  height: calc(100vh - 42px);
}
#preview-area.fullscreen .files-notes-col.open {
  height: 100vh; border-top: none; border-left: 1px solid var(--border);
}
#preview-area.fullscreen .files-notes-col #fnotes-editor {
  height: calc(100vh - 44px);
}

/* Zoom controls in preview header */
.zoom-controls { display: flex; align-items: center; gap: 3px; flex-shrink: 0; }
.zoom-btn { background: none; border: 1px solid var(--border); color: var(--text2);
  width: 22px; height: 22px; border-radius: 5px; cursor: pointer; font-size: 14px;
  display: flex; align-items: center; justify-content: center; line-height: 1;
  transition: background var(--transition), border-color var(--transition); }
.zoom-btn:hover { background: var(--bg4); border-color: var(--border2); }
.zoom-label { font-size: 11px; color: var(--text3); min-width: 34px; text-align: center; }
#fullscreen-btn { background: none; border: none; cursor: pointer; font-size: 15px;
  color: var(--text3); padding: 2px 5px; border-radius: 5px; line-height: 1;
  transition: color var(--transition), background var(--transition); }
#fullscreen-btn:hover { color: var(--text); background: var(--bg4); }
.preview-placeholder {
  height: 100%; display: flex; align-items: center; justify-content: center;
  flex-direction: column; gap: 10px; color: var(--text3);
}
.preview-placeholder .icon { font-size: 40px; opacity: .5; }

/* Summary panel */
.summary-toolbar {
  display: flex; align-items: center; gap: 8px; margin-bottom: 24px;
  padding-bottom: 18px; border-bottom: 1px solid var(--border); flex-wrap: wrap;
}
.summary-file-pill {
  display: inline-flex; align-items: center; gap: 4px;
  background: var(--bg3); border: 1px solid var(--border);
  border-radius: 20px; padding: 3px 10px 3px 12px;
  font-size: 11px; color: var(--text3); cursor: pointer;
  transition: border-color var(--transition), color var(--transition);
}
.summary-file-pill:hover, .summary-file-pill.active {
  border-color: var(--blue); color: var(--text);
}
.summary-file-pill.active { background: rgba(91,142,240,.1); color: var(--blue); }
.summary-pill-del {
  background: none; border: none; cursor: pointer; color: var(--text3);
  padding: 0 0 0 2px; font-size: 11px; line-height: 1;
  transition: color var(--transition);
}
.summary-pill-del:hover { color: var(--red); }
.summary-pill-rename {
  background: none; border: none; cursor: pointer; color: var(--text3);
  padding: 0; font-size: 10px; line-height: 1;
  transition: color var(--transition);
}
.summary-pill-rename:hover { color: var(--blue); }
#summary-editor {
  width: 100%; box-sizing: border-box; flex: 1; min-height: 400px;
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius);
  color: var(--text); font-size: 13px; font-family: "SF Mono", "Fira Code", monospace;
  line-height: 1.7; padding: 16px; resize: vertical;
  transition: border-color var(--transition);
}
#summary-editor:focus { outline: none; border-color: var(--blue); }
.md-content { max-width: 820px; }
.md-content h1 { font-size: 24px; color: var(--text); margin-bottom: 8px; line-height: 1.25; font-weight: 800; letter-spacing: -.02em; }
.md-content h2 { font-size: 17px; color: #93c5fd; margin: 32px 0 12px; padding-bottom: 8px; border-bottom: 1px solid var(--border); font-weight: 700; }
.md-content h3 { font-size: 14px; color: #a5b4fc; margin: 20px 0 8px; font-weight: 600; }
.md-content p  { color: var(--text2); line-height: 1.9; margin-bottom: 18px; text-align: justify; hyphens: auto; }
.md-content ul, .md-content ol { color: var(--text2); line-height: 1.9; margin: 8px 0 20px 24px; text-align: justify; hyphens: auto; }
.md-content li { margin-bottom: 8px; }
.md-content strong { color: #93c5fd; font-weight: 700; }
.md-content em { color: var(--text3); }
.md-content hr { border: none; border-top: 1px solid var(--border); margin: 32px 0; }
.md-content .katex-display { margin: 22px 0; overflow-x: auto; }
.md-content .katex { font-size: 1.05em; }
.md-content code { background: var(--bg3); padding: 2px 7px; border-radius: 5px; font-size: 12px; color: #6ee7b7; font-family: "SF Mono", monospace; border: 1px solid var(--border); }
.md-content blockquote { border-left: 3px solid var(--blue); padding: 8px 16px; color: var(--text3); margin: 14px 0; background: rgba(79,142,247,.05); border-radius: 0 var(--radius) var(--radius) 0; }
.md-content table { border-collapse: collapse; width: 100%; margin: 14px 0; font-size: 13px; border-radius: var(--radius); overflow: hidden; }
.md-content th { background: var(--bg3); padding: 9px 14px; text-align: left; color: #93c5fd; border: 1px solid var(--border); font-weight: 600; }
.md-content td { padding: 8px 14px; border: 1px solid var(--border); color: var(--text2); }
.md-content tr:nth-child(even) td { background: rgba(255,255,255,.02); }

/* Flashcard panel */
.flash-layout { max-width: 700px; margin: 0 auto; }
.flash-header { display: flex; align-items: center; gap: 10px; margin-bottom: 22px; flex-wrap: wrap; }
.flash-header-title { font-size: 14px; font-weight: 600; color: var(--text2); flex: 1; }
.flash-progress-bar { background: var(--border); border-radius: 10px; height: 5px; margin-bottom: 14px; overflow: hidden; }
.flash-progress-fill {
  height: 100%; border-radius: 10px; transition: width .5s cubic-bezier(.4,0,.2,1);
  background: linear-gradient(90deg, var(--blue), var(--purple));
}
.flash-meta {
  font-size: 12px; color: var(--text3); text-align: center; margin-bottom: 24px;
  display: flex; justify-content: center; align-items: center; gap: 16px;
}
.flash-meta span { display: flex; align-items: center; gap: 4px; }
.flash-timer { font-size: 12px; color: var(--text3); font-variant-numeric: tabular-nums; letter-spacing: .03em; }

.flash-card {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: var(--radius-xl); padding: 36px 32px; text-align: center;
  min-height: 220px; display: flex; flex-direction: column;
  align-items: center; justify-content: center; gap: 16px;
  margin-bottom: 20px;
  transition: border-color .2s;
}
.flash-card.flash-card-known   { border-color: var(--green); }
.flash-card.flash-card-unknown { border-color: var(--red); }
.flash-section { font-size: 10px; color: var(--text3); text-transform: uppercase; letter-spacing: .08em; font-weight: 600; }
.flash-course-badge { font-size: 10px; color: var(--purple); background: rgba(167,139,250,.12); padding: 3px 10px; border-radius: 12px; border: 1px solid rgba(167,139,250,.25); }
.flash-question { font-size: 19px; color: var(--text); line-height: 1.5; font-weight: 600; letter-spacing: -.01em; }
.flash-answer {
  background: var(--bg3); border: 1px solid var(--border2);
  border-radius: var(--radius); padding: 16px 22px;
  font-size: 14px; color: var(--text2); line-height: 1.7;
  display: none; text-align: left; width: 100%;
  animation: fadeIn .2s ease;
}
@keyframes fadeIn { from { opacity: 0; transform: translateY(4px); } to { opacity: 1; transform: none; } }

.flash-btns { display: flex; gap: 12px; justify-content: center; flex-wrap: wrap; }
.flash-btn {
  padding: 11px 26px; border-radius: var(--radius); border: none;
  cursor: pointer; font-size: 13px; font-weight: 600;
  transition: opacity var(--transition);
  letter-spacing: -.01em;
}
.flash-btn:hover { opacity: .85; }
.fb-reveal  { background: var(--bg4); color: var(--text2); border: 1px solid var(--border2); padding: 11px 32px; }
.fb-known   { background: #059669; color: #d1fae5; }
.fb-unknown { background: #dc2626; color: #fee2e2; }

.flash-done {
  text-align: center; padding: 48px 40px;
  display: none; flex-direction: column; align-items: center; gap: 14px;
}
.flash-done .big-icon { font-size: 64px; }
.flash-done h2 { font-size: 24px; color: var(--text); font-weight: 800; }
.flash-done p  { color: var(--text3); font-size: 14px; line-height: 1.6; }

/* Manage-cards view */
.manage-cards-wrap { max-width: 700px; margin: 0 auto; }
.manage-cards-header { display: flex; align-items: center; gap: 10px; margin-bottom: 20px; flex-wrap: wrap; }
.manage-cards-header h2 { font-size: 16px; font-weight: 700; color: var(--text); flex: 1; margin: 0; }
.mc-add-form {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 16px 18px; margin-bottom: 20px; display: flex; flex-direction: column; gap: 10px;
}
.mc-add-form textarea {
  width: 100%; resize: vertical; background: var(--bg3); border: 1px solid var(--border);
  border-radius: var(--radius); color: var(--text); font-size: 13px; padding: 8px 10px;
  font-family: inherit; min-height: 52px; box-sizing: border-box;
  transition: border-color var(--transition);
}
.mc-add-form textarea:focus { outline: none; border-color: var(--border2); }
.mc-add-form textarea::placeholder { color: var(--text3); }
.mc-add-row { display: flex; gap: 8px; align-items: flex-end; }
.mc-list { display: flex; flex-direction: column; gap: 8px; }
.mc-card {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius);
  padding: 12px 14px; display: flex; gap: 12px; align-items: flex-start;
}
.mc-card-body { flex: 1; min-width: 0; }
.mc-card-q { font-size: 13px; font-weight: 600; color: var(--text); margin-bottom: 4px; }
.mc-card-a {
  font-size: 12px; color: transparent; background: var(--bg3);
  border-radius: 3px; cursor: pointer; user-select: none;
  transition: color .15s, background .15s;
}
.mc-card-a:hover { color: var(--text3); background: transparent; }
.mc-card-del { background: none; border: none; color: var(--text3); cursor: pointer; font-size: 14px; padding: 2px 4px; flex-shrink: 0; }
.mc-card-del:hover { color: var(--red); }
.card-type-badge { font-size: 9px; text-transform: uppercase; letter-spacing: .09em; padding: 2px 7px; border-radius: 10px; font-weight: 700; display: inline-block; margin-bottom: 5px; }
.ct-recall      { color: #93c5fd; background: rgba(147,197,253,.1); border: 1px solid rgba(147,197,253,.2); }
.ct-mechanism   { color: #c4b5fd; background: rgba(196,181,253,.1); border: 1px solid rgba(196,181,253,.2); }
.ct-contrast    { color: #fdba74; background: rgba(253,186,116,.1); border: 1px solid rgba(253,186,116,.2); }
.ct-application { color: #6ee7b7; background: rgba(110,231,183,.1); border: 1px solid rgba(110,231,183,.2); }
.ct-cloze       { color: #67e8f9; background: rgba(103,232,249,.1); border: 1px solid rgba(103,232,249,.2); }
.cloze-blank { display: inline-block; min-width: 90px; height: 1.1em; border-bottom: 2px solid var(--blue); background: rgba(79,142,247,.08); border-radius: 3px 3px 0 0; margin: 0 3px; vertical-align: middle; }
.cloze-fill  { color: var(--blue); font-weight: 700; border-bottom: 2px solid var(--blue); padding: 0 3px; }
.mc-section-label { font-size: 11px; font-weight: 600; color: var(--text3); text-transform: uppercase; letter-spacing: .06em; margin: 16px 0 8px; }
.mc-gen-section {
  display: flex; align-items: center; gap: 14px;
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: var(--radius-lg); padding: 14px 16px; margin-bottom: 14px;
}
.mc-gen-section-disabled { opacity: 0.5; font-size: 13px; color: var(--text3); gap: 10px; }
.mc-gen-info { display: flex; align-items: center; gap: 12px; flex: 1; min-width: 0; }
.mc-gen-icon { font-size: 22px; flex-shrink: 0; }
.mc-gen-title { font-size: 13px; font-weight: 600; color: var(--text); }
.mc-gen-desc { font-size: 11px; color: var(--text3); margin-top: 2px; }

.flash-kbd-hint { text-align: center; font-size: 11px; color: var(--text3); margin-top: 10px; }
.flash-kbd-hint kbd {
  background: var(--bg3); border: 1px solid var(--border2); border-radius: 5px;
  padding: 2px 6px; font-size: 10px; font-family: "SF Mono", monospace; color: var(--text2);
}

/* ── SRS Study System ── */
.srs-overview { max-width: 700px; margin: 0 auto; padding: 4px 0; }
.srs-overview-header { display: flex; align-items: center; gap: 10px; margin-bottom: 18px; }
.srs-overview-header h2 { font-size: 18px; font-weight: 700; color: var(--text); flex: 1; margin: 0; }
.srs-stats { display: grid; grid-template-columns: repeat(4, 1fr); gap: 8px; margin-bottom: 16px; }
.srs-stat { background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg); padding: 14px 10px; text-align: center; }
.srs-stat-num { font-size: 30px; font-weight: 800; line-height: 1; margin-bottom: 5px; }
.srs-stat-label { font-size: 10px; color: var(--text3); text-transform: uppercase; letter-spacing: .06em; }
.srs-stat-due .srs-stat-num { color: var(--yellow); }
.srs-stat-new .srs-stat-num { color: var(--blue); }
.srs-stat-learning .srs-stat-num { color: var(--orange); }
.srs-stat-mastered .srs-stat-num { color: var(--purple); }
.srs-action-btns { display: flex; gap: 8px; margin-bottom: 18px; flex-wrap: wrap; }
.srs-badge { display: inline-flex; align-items: center; font-size: 9px; font-weight: 700; letter-spacing: .06em; text-transform: uppercase; padding: 2px 7px; border-radius: 10px; }
.srs-badge-new      { background: rgba(91,142,240,.15);  color: var(--blue);   }
.srs-badge-due      { background: rgba(251,191,36,.15);  color: var(--yellow); }
.srs-badge-overdue  { background: rgba(248,113,113,.18); color: var(--red);    }
.srs-badge-learning { background: rgba(251,146,60,.15);  color: var(--orange); }
.srs-badge-review   { background: rgba(52,211,153,.12);  color: var(--green);  }
.srs-badge-mastered { background: rgba(167,139,250,.15); color: var(--purple); }
.rating-grid { display: grid; grid-template-columns: repeat(4, 1fr); gap: 8px; margin-top: 16px; }
.rb { border: 1px solid transparent; border-radius: var(--radius); padding: 12px 6px 10px; cursor: pointer; font-size: 13px; font-weight: 600; display: flex; flex-direction: column; align-items: center; gap: 5px; transition: opacity .12s, transform .1s; }
.rb:hover { opacity: .85; transform: translateY(-1px); }
.rb:active { transform: scale(.97); }
.rb-label { font-size: 13px; font-weight: 600; }
.rb-hint  { font-size: 10px; opacity: .6; background: rgba(0,0,0,.12); padding: 1px 6px; border-radius: 3px; font-family: "SF Mono", monospace; }
.rb-again { background: rgba(248,113,113,.15); color: var(--red);    border-color: rgba(248,113,113,.3); }
.rb-hard  { background: rgba(251,146,60,.15);  color: var(--orange); border-color: rgba(251,146,60,.3);  }
.rb-good  { background: rgba(91,142,240,.15);  color: var(--blue);   border-color: rgba(91,142,240,.3);  }
.rb-easy  { background: rgba(52,211,153,.15);  color: var(--green);  border-color: rgba(52,211,153,.3);  }
.srs-summary { max-width: 700px; margin: 0 auto; text-align: center; padding: 8px 0; }
.srs-summary h2 { font-size: 24px; font-weight: 800; color: var(--text); margin: 0 0 6px; }
.srs-summary-sub { font-size: 13px; color: var(--text3); margin-bottom: 22px; }
.srs-rating-breakdown { display: grid; grid-template-columns: repeat(4, 1fr); gap: 8px; margin-bottom: 20px; }
.srb-item { background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg); padding: 14px 8px; }
.srb-num { font-size: 26px; font-weight: 800; line-height: 1; margin-bottom: 4px; }
.srb-label { font-size: 11px; color: var(--text3); }
.srs-next-due { background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg); padding: 12px 16px; margin-bottom: 18px; text-align: left; font-size: 13px; color: var(--text2); }
.srs-next-due-title { font-size: 10px; font-weight: 700; color: var(--text3); text-transform: uppercase; letter-spacing: .06em; margin-bottom: 6px; }
.srs-streak {
  display: inline-block; font-size: 15px; font-weight: 700; color: #f97316;
  background: rgba(249,115,22,.12); border: 1px solid rgba(249,115,22,.3);
  border-radius: 20px; padding: 4px 16px; margin: 0 auto 14px;
  box-shadow: 0 0 10px rgba(249,115,22,.25);
}
.srs-streak-day1 { color: var(--text2); background: var(--bg2); border-color: var(--border); box-shadow: none; }
.srs-hard-cards {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 10px 14px; margin-bottom: 18px; text-align: left; font-size: 12px; color: var(--text3);
}
.srs-hard-cards summary { cursor: pointer; font-weight: 600; color: var(--orange, #f97316); font-size: 12px; }
.srs-hard-cards ul { margin: 8px 0 2px 16px; padding: 0; line-height: 1.7; }
.srs-hard-cards li { color: var(--text2); }
@media (max-width: 500px) { .srs-stats, .rating-grid, .srs-rating-breakdown { grid-template-columns: repeat(2, 1fr); } }

/* Notes panel (course-level) */
.notes-panel { display: flex; flex-direction: column; height: calc(100vh - 158px); }
.notes-toolbar { display: flex; align-items: center; gap: 8px; margin-bottom: 12px; }
#notes-editor {
  flex: 1; background: var(--bg2); border: 1px solid var(--border);
  border-radius: var(--radius-lg); padding: 20px 22px; color: var(--text2);
  font-size: 14px; line-height: 1.8; resize: none; outline: none;
  font-family: inherit; transition: border-color var(--transition), box-shadow var(--transition);
}
#notes-editor:focus { border-color: var(--blue); box-shadow: var(--glow-blue); }
#notes-preview {
  flex: 1; background: var(--bg2); border: 1px solid var(--border);
  border-radius: var(--radius-lg); padding: 20px 22px; overflow-y: auto;
  display: none;
}
#notes-saved { font-size: 12px; display: none; }

/* File-level notes sliding panel */
.files-notes-col {
  width: 0; overflow: hidden; flex-shrink: 0;
  transition: width 200ms ease;
  border-left: 1px solid transparent;
  display: flex; flex-direction: column;
}
.files-notes-col.open {
  width: 290px; min-width: 200px;
  border-left-color: var(--border);
}
.fnotes-header {
  padding: 10px 14px 8px;
  display: flex; align-items: center; gap: 8px;
  border-bottom: 1px solid var(--border); flex-shrink: 0;
}
.fnotes-title {
  font-size: 11px; font-weight: 600; color: var(--text3); flex: 1;
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}
.fnotes-dl-btn {
  background: none; border: none; color: var(--text3); cursor: pointer;
  font-size: 13px; padding: 2px 5px; border-radius: 4px; line-height: 1;
  transition: color var(--transition), background var(--transition);
}
.fnotes-dl-btn:hover { color: var(--text2); background: var(--bg3); }
#fnotes-saved { font-size: 10px; color: var(--text3); white-space: nowrap; }
#fnotes-editor {
  flex: 1; background: transparent; border: none; outline: none;
  color: var(--text2); font-size: 12.5px; line-height: 1.75;
  resize: none; padding: 14px 14px;
  font-family: inherit;
}
#fnotes-editor::placeholder { color: var(--text3); }
#fnotes-rendered {
  flex: 1; overflow-y: auto; padding: 14px 14px;
  display: none; font-size: 12.5px; line-height: 1.75; color: var(--text2);
}
#fnotes-rendered h1,#fnotes-rendered h2,#fnotes-rendered h3 {
  color: var(--text1); margin: 12px 0 6px; font-size: 13px; font-weight: 600;
}
#fnotes-rendered h1 { font-size: 15px; }
#fnotes-rendered p { margin: 4px 0 8px; }
#fnotes-rendered ul,#fnotes-rendered ol { padding-left: 16px; margin: 4px 0 8px; }
#fnotes-rendered li { margin: 3px 0; }
#fnotes-rendered code { background: var(--bg3); padding: 1px 5px; border-radius: 3px; font-size: 11.5px; }
#fnotes-rendered pre { background: var(--bg3); padding: 10px; border-radius: var(--radius); overflow-x: auto; margin: 8px 0; }
#fnotes-rendered pre code { background: none; padding: 0; }
/* Task list checkboxes — only strip bullet from task-list items, not all lists */
#fnotes-rendered ul { padding-left: 18px; margin: 4px 0 8px; }
#fnotes-rendered li { margin: 3px 0; }
#fnotes-rendered li.task-list-item { display: flex; align-items: flex-start; gap: 7px; padding: 2px 0; }
#fnotes-rendered li.task-list-item::marker { content: ''; }
#fnotes-rendered li input.task-list-item-checkbox {
  appearance: none; -webkit-appearance: none;
  width: 14px; height: 14px; border: 1.5px solid var(--text3);
  border-radius: 3px; cursor: pointer; flex-shrink: 0; margin-top: 4px;
  background: transparent; transition: background var(--transition), border-color var(--transition);
}
#fnotes-rendered li input.task-list-item-checkbox:checked {
  background: var(--blue); border-color: var(--blue);
  background-image: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 10 10'%3E%3Cpath d='M1.5 5l2.5 2.5 4.5-4.5' stroke='white' stroke-width='1.5' fill='none' stroke-linecap='round' stroke-linejoin='round'/%3E%3C/svg%3E");
  background-size: 10px; background-repeat: no-repeat; background-position: center;
}

/* Search results */
.search-results-header { font-size: 12px; color: var(--text3); margin-bottom: 16px; }
.search-result {
  border-radius: var(--radius); padding: 12px 14px; margin-bottom: 2px;
  cursor: pointer; transition: background var(--transition);
}
.search-result:hover { background: var(--bg3); }
.search-result-course { font-size: 12px; color: var(--blue); margin-bottom: 5px; font-weight: 600; display: flex; align-items: center; gap: 8px; }
.search-result-count { font-size: 10px; background: rgba(79,142,247,.15); color: var(--blue); padding: 1px 6px; border-radius: 10px; }
.search-result-snippet { font-size: 13px; color: var(--text2); line-height: 1.6; }
.search-result-snippet mark { background: rgba(79,142,247,.2); color: #93c5fd; border-radius: 3px; padding: 0 3px; }
.search-empty { text-align: center; padding: 60px; color: var(--text3); font-size: 14px; }

/* Log */
#log-box {
  position: fixed; bottom: 16px; right: 16px; width: 440px;
  background: rgba(10,12,22,.97); border: 1px solid var(--border); border-radius: var(--radius-lg);
  font-family: "SF Mono", monospace; font-size: 11px;
  color: #4ade80; display: none; z-index: 100; box-shadow: var(--shadow-lg);
  flex-direction: column; max-height: 280px; backdrop-filter: blur(12px);
}
#log-header { padding: 9px 14px; border-bottom: 1px solid var(--border); display: flex; align-items: center; gap: 8px; }
#log-title { flex: 1; font-size: 11px; color: var(--text3); font-family: sans-serif; }
#log-close { cursor: pointer; color: var(--text3); font-size: 14px; background: none; border: none; transition: color var(--transition); }
#log-close:hover { color: var(--text); }
#log-content { padding: 12px 14px; white-space: pre-wrap; overflow-y: auto; flex: 1; line-height: 1.6; }

/* Toast */
#toast-container {
  position: fixed; bottom: 20px; left: 50%; transform: translateX(-50%);
  display: flex; flex-direction: column; gap: 6px; z-index: 200;
  pointer-events: none; align-items: center;
}
.toast {
  background: rgba(24,27,46,.95); border: 1px solid var(--border2); border-radius: 20px;
  padding: 9px 18px; font-size: 13px; color: var(--text2);
  box-shadow: var(--shadow-lg); backdrop-filter: blur(12px);
  animation: toastIn .2s cubic-bezier(.4,0,.2,1);
  display: flex; align-items: center; gap: 8px; white-space: nowrap;
}
.toast.toast-ok  { border-color: rgba(52,211,153,.4); color: #6ee7b7; }
.toast.toast-err { border-color: rgba(248,113,113,.4); color: #fca5a5; }
@keyframes toastIn { from { opacity: 0; transform: translateY(8px) scale(.95); } to { opacity: 1; transform: none; } }

/* Keyboard shortcuts overlay */
#shortcuts-overlay {
  position: fixed; inset: 0; background: rgba(0,0,0,.7); z-index: 300;
  display: none; align-items: center; justify-content: center;
  backdrop-filter: blur(4px);
}
#shortcuts-overlay.open { display: flex; }
.shortcuts-box {
  background: var(--bg2); border: 1px solid var(--border2); border-radius: var(--radius-xl);
  padding: 30px 34px; max-width: 500px; width: 90%;
  box-shadow: var(--shadow-lg);
  animation: panelIn .2s ease;
}
.shortcuts-box h3 { font-size: 17px; color: var(--text); margin-bottom: 20px; font-weight: 700; }
.shortcut-row { display: flex; align-items: center; gap: 14px; margin-bottom: 10px; }
.shortcut-row kbd {
  background: var(--bg4); border: 1px solid var(--border2); border-radius: 6px;
  padding: 3px 9px; font-size: 11px; font-family: "SF Mono", monospace;
  color: var(--text2); min-width: 40px; text-align: center; flex-shrink: 0;
  box-shadow: 0 2px 0 var(--border);
}
.shortcut-desc { font-size: 13px; color: var(--text3); }
.shortcuts-section { font-size: 9px; text-transform: uppercase; letter-spacing: .1em; color: var(--text3); margin: 16px 0 8px; font-weight: 700; }

/* Confirm modal */
#confirm-overlay {
  position: fixed; inset: 0; background: rgba(0,0,0,.7); z-index: 300;
  display: none; align-items: center; justify-content: center; backdrop-filter: blur(4px);
}
#confirm-overlay.open { display: flex; }
.confirm-box {
  background: var(--bg2); border: 1px solid var(--border2); border-radius: var(--radius-xl);
  padding: 30px 34px; max-width: 380px; width: 90%; text-align: center;
  box-shadow: var(--shadow-lg); animation: panelIn .2s ease;
}
.confirm-box h3 { font-size: 17px; color: var(--text); margin-bottom: 10px; font-weight: 700; }
.confirm-box p { font-size: 13px; color: var(--text3); margin-bottom: 24px; line-height: 1.6; }
.confirm-btns { display: flex; gap: 10px; justify-content: center; }

/* Diff modal */
#diff-overlay {
  position: fixed; inset: 0; background: rgba(0,0,0,.7); z-index: 300;
  display: none; align-items: center; justify-content: center; backdrop-filter: blur(4px);
}
#diff-overlay.open { display: flex; }
#diff-box {
  background: var(--bg2); border: 1px solid var(--border2); border-radius: var(--radius-xl);
  padding: 24px 28px; width: min(820px, 94vw); max-height: 85vh;
  display: flex; flex-direction: column; box-shadow: var(--shadow-lg); animation: panelIn .2s ease;
}
.diff-view {
  flex: 1; overflow-y: auto; font-size: 12px; font-family: "SF Mono", monospace;
  border: 1px solid var(--border); border-radius: var(--radius); background: var(--bg3);
  min-height: 200px; max-height: 60vh;
}
.diff-line { padding: 2px 10px; white-space: pre-wrap; word-break: break-word; line-height: 1.6; }
.diff-line-add    { background: rgba(34,197,94,.15); color: #4ade80; }
.diff-line-remove { background: rgba(239,68,68,.15); color: #f87171; }
.diff-line-same   { color: var(--text2); }

/* Spinner */
.spin {
  display: inline-block; width: 13px; height: 13px;
  border: 2px solid var(--border2); border-top-color: var(--blue);
  border-radius: 50%; animation: spin .6s linear infinite; vertical-align: middle;
}
@keyframes spin { to { transform: rotate(360deg); } }

/* Empty state */
.empty-state {
  height: 100%; display: flex; flex-direction: column;
  align-items: center; justify-content: center;
  color: var(--text3); gap: 14px; text-align: center; padding: 40px;
}
.empty-state .icon { font-size: 56px; opacity: .6; }
.empty-state h3 { font-size: 17px; color: var(--text2); font-weight: 600; }
.empty-state p { font-size: 13px; max-width: 300px; line-height: 1.7; }

/* Scrollbar */
::-webkit-scrollbar { width: 4px; height: 4px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: var(--border2); border-radius: 10px; }
::-webkit-scrollbar-thumb:hover { background: var(--text3); }

/* files panel */
#panel-files { padding: 16px 20px; }

/* New-files badge */
.new-badge {
  display: inline-block; font-size: 9px; font-weight: 700;
  background: var(--yellow); color: #000; border-radius: 8px;
  padding: 1px 5px; vertical-align: middle; margin-left: 5px; letter-spacing: .02em;
}

/* File meta */
.file-meta-info { font-size: 10px; color: var(--text3); margin-left: auto; flex-shrink: 0; padding-left: 6px; }

/* Chat panel — NO display override here, uses .panel base class */
#panel-chat { padding: 0; overflow: hidden; }
#panel-chat.active { display: flex; flex-direction: column; }
#chat-body { flex: 1; display: flex; flex-direction: column; min-height: 0; }

/* Course info panel */
.info-card { background: var(--bg2); border: 1px solid var(--border); border-radius: 10px; padding: 20px 24px; margin-bottom: 14px; }
.info-title { font-size: 17px; font-weight: 700; color: var(--text1); margin-bottom: 4px; line-height: 1.3; }
.info-subtitle { font-size: 13px; color: var(--text3); margin-bottom: 16px; }
.info-grid { display: grid; grid-template-columns: 140px 1fr; gap: 8px 16px; font-size: 13px; }
.info-label { color: var(--text3); font-weight: 600; font-size: 11px; text-transform: uppercase; letter-spacing: .04em; padding-top: 2px; }
.info-value { color: var(--text2); line-height: 1.5; }
.info-desc { font-size: 13px; color: var(--text2); line-height: 1.65; margin-top: 16px; padding-top: 16px; border-top: 1px solid var(--border); white-space: pre-wrap; }
.custom-info-row { display: grid; grid-template-columns: 140px 1fr auto; gap: 6px; align-items: start; margin-bottom: 6px; }
.custom-info-row input, .custom-info-row textarea {
  background: var(--bg3); border: 1px solid var(--border); border-radius: 6px;
  color: var(--text); font-size: 12px; padding: 5px 8px; font-family: inherit; outline: none;
  transition: border-color .15s;
}
.custom-info-row input:focus, .custom-info-row textarea:focus { border-color: var(--blue); }
.custom-info-row input { height: 30px; }
.custom-info-row textarea { resize: vertical; min-height: 30px; line-height: 1.5; }
.custom-info-del { background: none; border: none; cursor: pointer; color: var(--text3); font-size: 13px; padding: 4px; border-radius: 4px; transition: color .15s; margin-top: 2px; }
.custom-info-del:hover { color: var(--red); }
.chat-layout { display: flex; flex-direction: column; flex: 1; min-height: 0; }
.chat-messages {
  flex: 1; overflow-y: auto; padding: 28px 40px; display: flex; flex-direction: column; gap: 18px;
}
.chat-msg { display: flex; gap: 12px; max-width: 760px; animation: fadeIn .2s ease; }
.chat-msg.user { align-self: flex-end; flex-direction: row-reverse; }
.chat-avatar {
  width: 32px; height: 32px; border-radius: 50%; flex-shrink: 0;
  display: flex; align-items: center; justify-content: center; font-size: 15px;
  background: var(--bg4); border: 1px solid var(--border2);
}
.chat-msg.user .chat-avatar { background: linear-gradient(135deg, var(--blue2), var(--blue3)); }
.chat-bubble {
  background: var(--bg2); border: 1px solid var(--border); border-radius: var(--radius-lg);
  padding: 13px 18px; font-size: 13.5px; color: var(--text2); line-height: 1.7;
  max-width: 600px; box-shadow: var(--shadow-sm);
}
.chat-msg.user .chat-bubble {
  background: rgba(91,142,240,.12);
  border-color: rgba(91,142,240,.25); color: var(--text);
  border-radius: var(--radius-lg) var(--radius-lg) 4px var(--radius-lg);
}
.chat-msg:not(.user) .chat-bubble { border-radius: var(--radius-lg) var(--radius-lg) var(--radius-lg) 4px; }
.chat-bubble.streaming::after { content: '▋'; animation: blink .8s step-end infinite; color: var(--blue); }
@keyframes blink { 50% { opacity: 0; } }
.chat-bubble p { margin-bottom: 8px; }
.chat-bubble p:last-child { margin-bottom: 0; }
.chat-input-row {
  padding: 14px 24px; border-top: 1px solid var(--border); background: var(--bg2);
  display: flex; gap: 10px; align-items: flex-end; flex-shrink: 0;
}
.chat-input-row textarea {
  flex: 1; background: var(--bg3); border: 1px solid var(--border); border-radius: var(--radius-lg);
  color: var(--text); font-size: 13px; padding: 11px 16px; resize: none; outline: none;
  font-family: inherit; line-height: 1.6; max-height: 120px; min-height: 44px;
  transition: border-color var(--transition), box-shadow var(--transition);
}
.chat-input-row textarea:focus { border-color: var(--blue); box-shadow: var(--glow-blue); }
.chat-input-row textarea::placeholder { color: var(--text3); }
.chat-input-btn {
  padding: 10px 16px; background: var(--blue); color: #fff; border: none;
  border-radius: var(--radius-lg); cursor: pointer; font-size: 13px; font-weight: 600;
  transition: opacity var(--transition);
  flex-shrink: 0;
}
.chat-input-btn:hover { opacity: .85; }
.chat-input-btn:disabled { opacity: .3; cursor: not-allowed; }
.chat-suggestions { display: flex; flex-wrap: wrap; gap: 6px; padding: 0 40px 14px; }
.chat-toolbar { display: flex; gap: 6px; padding: 8px 16px 0; justify-content: flex-end; }
.chat-toolbar button { font-size: 11px; }
.chat-context-bar {
  display: flex; align-items: center; gap: 8px; padding: 6px 24px 0;
  flex-shrink: 0;
}
.chat-context-label { font-size: 11px; color: var(--text3); white-space: nowrap; }
.chat-context-select {
  flex: 1; background: var(--bg3); border: 1px solid var(--border); border-radius: var(--radius-lg);
  color: var(--text); font-size: 12px; padding: 5px 10px; outline: none; cursor: pointer;
  transition: border-color var(--transition);
}
.chat-context-select:focus { border-color: var(--blue); }
.chat-file-pill {
  display: inline-flex; align-items: center; gap: 5px;
  background: var(--blue); color: #fff; font-size: 11px; font-weight: 600;
  border-radius: 20px; padding: 2px 10px 2px 8px; margin: 0 24px 6px;
  flex-shrink: 0; max-width: calc(100% - 48px); overflow: hidden;
}
.chat-file-pill span { overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
.chat-file-pill button { background: none; border: none; color: rgba(255,255,255,.7); cursor: pointer; font-size: 12px; padding: 0; line-height: 1; }
.chat-file-pill button:hover { color: #fff; }
.chat-history-panel { position: absolute; top: 0; right: 0; bottom: 0; width: 280px; background: var(--bg2); border-left: 1px solid var(--border); display: flex; flex-direction: column; z-index: 10; transform: translateX(100%); transition: transform .2s ease; }
.chat-history-panel.open { transform: translateX(0); }
.chat-history-header { display: flex; align-items: center; padding: 14px 16px; border-bottom: 1px solid var(--border); gap: 8px; }
.chat-history-header h3 { flex: 1; font-size: 13px; font-weight: 700; color: var(--text); margin: 0; }
.chat-history-list { flex: 1; overflow-y: auto; padding: 8px; display: flex; flex-direction: column; gap: 4px; }
.chat-history-item { padding: 8px 10px; border-radius: var(--radius); border: 1px solid var(--border); cursor: pointer; transition: background .15s; font-size: 12px; font-weight: 600; color: var(--text); display: flex; align-items: center; gap: 6px; }
.chat-history-item:hover { background: var(--bg3); }
.chat-history-item-label { flex: 1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
.chat-history-rename-btn { flex-shrink: 0; background: none; border: none; color: var(--text3); cursor: pointer; font-size: 12px; padding: 2px 4px; border-radius: 4px; opacity: 0; transition: opacity .15s, color .15s; }
.chat-history-item:hover .chat-history-rename-btn { opacity: 1; }
.chat-history-rename-btn:hover { color: var(--blue); }
.chat-history-rename-input { flex: 1; background: var(--bg); border: 1px solid var(--blue); border-radius: 4px; color: var(--text); font-size: 12px; font-weight: 600; padding: 2px 6px; outline: none; min-width: 0; }
.chat-suggestion {
  font-size: 12px; padding: 5px 12px; background: var(--bg3); border: 1px solid var(--border);
  border-radius: 16px; cursor: pointer; color: var(--text3);
  transition: all var(--transition);
}
.chat-suggestion:hover { border-color: var(--blue); color: var(--text2); background: var(--bg4); transform: translateY(-1px); }

/* Recommendations */

/* Sync tile */
#sync-tile {
  display: inline-flex; flex-direction: column; gap: 0;
  font-size: 11px; color: var(--text3);
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 10px; padding: 6px 12px; margin-bottom: 18px;
}
#sync-tile-header { display: flex; align-items: center; gap: 6px; }
#sync-tile .sync-dot {
  width: 6px; height: 6px; border-radius: 50%; background: var(--green); flex-shrink: 0;
}
#sync-tile-files {
  margin-top: 6px; display: flex; flex-direction: column; gap: 2px;
  border-top: 1px solid var(--border); padding-top: 6px;
}
#sync-tile-files .sync-file {
  display: flex; align-items: center; gap: 6px; font-size: 10px; color: var(--text2);
}
#sync-tile-files .sync-file-course {
  color: var(--text3); flex-shrink: 0;
}
#sync-tile-files .sync-file-name {
  overflow: hidden; text-overflow: ellipsis; white-space: nowrap;
}

/* Courses overview */
.co-semester { margin-bottom: 32px; }
.co-semester-title {
  font-size: 11px; font-weight: 700; letter-spacing: .1em; text-transform: uppercase;
  color: var(--text3); margin-bottom: 12px; padding-bottom: 6px;
  border-bottom: 1px solid var(--border);
}
.co-grid {
  display: grid;
  grid-template-columns: repeat(auto-fill, minmax(200px, 1fr));
  gap: 12px;
}
.co-card {
  background: var(--bg2); border: 1px solid var(--border); border-radius: 10px;
  padding: 14px 14px 12px; cursor: pointer; position: relative; overflow: hidden;
  transition: border-color .15s, transform .15s, box-shadow .15s;
  display: flex; flex-direction: column; gap: 6px;
}
.co-card:hover { border-color: var(--blue); transform: translateY(-2px); box-shadow: var(--shadow); }
.co-card-accent { position: absolute; top: 0; left: 0; right: 0; height: 3px; border-radius: 10px 10px 0 0; }
.co-card-name { font-size: 13px; font-weight: 600; color: var(--text); line-height: 1.35; margin-top: 4px; }
.co-card-badges { display: flex; flex-wrap: wrap; gap: 4px; margin-top: 2px; }
.co-badge {
  font-size: 10px; padding: 1px 6px; border-radius: 20px; font-weight: 600;
  background: var(--bg3); color: var(--text3);
}
.co-badge.green  { background: rgba(52,211,153,.15); color: var(--green); }
.co-badge.blue   { background: rgba(91,142,240,.15); color: var(--blue); }
.co-badge.yellow { background: rgba(251,191,36,.15);  color: var(--yellow); }
.co-badge.purple { background: rgba(167,139,250,.15); color: var(--purple); }
.co-card-bar { height: 3px; border-radius: 2px; background: var(--bg3); margin-top: 4px; overflow: hidden; }
.co-card-bar-fill { height: 100%; border-radius: 2px; background: linear-gradient(90deg,var(--blue),var(--purple)); }

/* Last-studied label in sidebar */
.citem-last { font-size: 9px; color: var(--text3); margin-top: 1px; }

/* Search source badge */
.search-source { font-size: 10px; padding: 2px 7px; border-radius: 8px; font-weight: 500; }
.search-source-summary { background: rgba(79,142,247,.12); color: var(--blue); }
.search-source-notes   { background: rgba(167,139,250,.12); color: var(--purple); }
.search-source-both    { background: rgba(52,211,153,.12);  color: var(--green); }

/* ── Command Palette ── */
#cmd-overlay {
  position: fixed; inset: 0; background: rgba(0,0,0,.65); z-index: 400;
  display: none; align-items: flex-start; justify-content: center;
  padding-top: 14vh; backdrop-filter: blur(5px);
}
#cmd-overlay.open { display: flex; }
#cmd-box {
  background: var(--bg2); border: 1px solid var(--border2); border-radius: var(--radius-xl);
  width: 100%; max-width: 560px; box-shadow: var(--shadow-lg);
  overflow: hidden; animation: panelIn .15s ease;
}
#cmd-input {
  width: 100%; padding: 16px 20px; background: transparent;
  border: none; border-bottom: 1px solid var(--border);
  color: var(--text); font-size: 15px; outline: none; font-family: inherit;
}
#cmd-input::placeholder { color: var(--text3); }
#cmd-results { max-height: 380px; overflow-y: auto; padding: 6px; }
.cmd-item {
  display: flex; align-items: center; gap: 12px; padding: 10px 14px;
  border-radius: var(--radius); cursor: pointer;
  transition: background var(--transition); user-select: none;
}
.cmd-item.cmd-selected { background: rgba(79,142,247,.14); }
.cmd-item-icon { font-size: 15px; width: 22px; text-align: center; flex-shrink: 0; }
.cmd-item-label { flex: 1; font-size: 13px; font-weight: 500; color: var(--text2); }
.cmd-item.cmd-selected .cmd-item-label { color: var(--text); }
.cmd-item-meta { font-size: 11px; color: var(--text3); flex-shrink: 0; }
.cmd-section { font-size: 9px; text-transform: uppercase; letter-spacing: .1em; color: var(--text3); padding: 8px 14px 3px; font-weight: 700; }
#cmd-empty { text-align: center; padding: 30px; color: var(--text3); font-size: 13px; }
#cmd-footer {
  padding: 8px 16px; border-top: 1px solid var(--border);
  display: flex; gap: 16px; align-items: center;
}
.cmd-hint { font-size: 10px; color: var(--text3); display: flex; align-items: center; gap: 4px; }
.cmd-hint kbd {
  background: var(--bg4); border: 1px solid var(--border2); border-radius: 4px;
  padding: 1px 5px; font-size: 9px; font-family: "SF Mono", monospace; color: var(--text3);
}

/* ── Summary ToC ── */
.summary-layout { display: flex; gap: 28px; align-items: flex-start; }
.summary-toc {
  width: 200px; min-width: 160px; flex-shrink: 0;
  position: sticky; top: 16px; align-self: flex-start;
  max-height: calc(100vh - 120px); overflow-y: auto; padding-bottom: 20px;
}
.summary-toc-title {
  font-size: 9px; text-transform: uppercase; letter-spacing: .1em;
  color: var(--text3); font-weight: 700; margin-bottom: 8px; padding-left: 8px;
}
.toc-link {
  display: block; font-size: 12px; color: var(--text3); text-decoration: none;
  padding: 4px 8px; border-radius: 5px; border-left: 2px solid transparent;
  transition: color var(--transition), border-color var(--transition), background var(--transition);
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis; margin-bottom: 1px; line-height: 1.4;
}
.toc-link:hover { color: var(--text2); border-left-color: var(--border2); background: var(--bg3); }
.toc-link.toc-h3 { padding-left: 18px; font-size: 11px; }
.toc-link.toc-active { color: var(--blue); border-left-color: var(--blue); background: rgba(79,142,247,.07); }
.summary-content-wrap { flex: 1; min-width: 0; }
@media (max-width: 960px) { .summary-toc { display: none; } }

/* ── Welcome modal ── */
#welcome-overlay {
  position: fixed; inset: 0; z-index: 9999;
  background: rgba(0,0,0,.55); backdrop-filter: blur(6px);
  display: flex; align-items: center; justify-content: center;
  animation: fadeIn .25s ease;
}
#welcome-overlay.hidden { display: none; }
@keyframes fadeIn { from { opacity: 0; } to { opacity: 1; } }
#welcome-modal {
  background: var(--bg2); border: 1px solid var(--border2);
  border-radius: 18px; padding: 40px 44px 36px;
  max-width: 480px; width: calc(100% - 40px);
  box-shadow: 0 24px 64px rgba(0,0,0,.35);
  animation: slideUp .28s cubic-bezier(.22,1,.36,1);
}
@keyframes slideUp { from { transform: translateY(18px); opacity: 0; } to { transform: translateY(0); opacity: 1; } }
#welcome-modal .welcome-from {
  font-size: 11px; color: var(--text3); margin: 0 0 10px;
  text-transform: uppercase; letter-spacing: .07em;
}
#welcome-modal h1 {
  font-size: 21px; font-weight: 700; color: var(--text1);
  margin: 0 0 18px; letter-spacing: -.01em; line-height: 1.25;
}
#welcome-modal .welcome-body {
  font-size: 14px; color: var(--text2); line-height: 1.7;
  margin: 0 0 28px;
}
#welcome-modal .welcome-btn {
  width: 100%; padding: 11px; font-size: 14px; font-weight: 600;
  background: var(--blue); color: #fff; border: none;
  border-radius: var(--radius-lg); cursor: pointer;
  transition: background var(--transition);
}
#welcome-modal .welcome-btn:hover { background: var(--blue2); }
#welcome-modal .welcome-footer {
  margin-top: 14px; text-align: center;
  font-size: 11px; color: var(--text3);
}
#welcome-modal .welcome-footer a {
  color: var(--text3); text-decoration: none;
}
#welcome-modal .welcome-footer a:hover { color: var(--blue); }

/* ── Selection toolbar ── */
#selection-toolbar {
  background: rgba(79,142,247,.07); border: 1px solid rgba(79,142,247,.25);
  border-radius: var(--radius); padding: 8px 10px; margin-top: 8px;
  display: none; flex-direction: column; gap: 5px;
}
#selection-toolbar.visible { display: flex; }
#sel-count { font-size: 11px; color: var(--blue); font-weight: 600; }
.sel-toolbar-btns { display: flex; gap: 5px; }
</style>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/katex.min.css">
<script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/katex.min.js"></script>
<script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.11/dist/contrib/auto-render.min.js"
  onload="window._katexReady=true"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js"></script>
<script>pdfjsLib.GlobalWorkerOptions.workerSrc = 'https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js';</script>
</head>
<body>

<!-- Topbar -->
<div id="topbar">
  <button id="sidebar-toggle" onclick="toggleSidebar()" title="Toggle course list (B)">☰</button>
  <div id="topbar-logo" onclick="goHome()" style="cursor:pointer"><span class="logo-icon">📚</span><span class="logo-wordmark"> <span>Stud.IP</span> Dashboard</span></div>
  <input id="search-global" type="text" placeholder="🔍  Search all summaries… (Ctrl+K)" oninput="handleGlobalSearch(event)">
  <button class="tbtn btn-gray" id="pdf-find-topbar-btn" onclick="togglePdfFind()"
    style="display:none;font-size:12px;padding:6px 10px" title="Search in PDF (Ctrl+F)">🔍 in PDF</button>
  <button class="tbtn btn-gray" id="courses-overview-btn" onclick="goCoursesOverview()" title="Course overview">All courses</button>
  <button class="tbtn btn-gray topbar-shortcuts" onclick="showShortcuts()" title="Keyboard shortcuts (?)">⌨️</button>
  <button class="tbtn btn-gray" id="theme-toggle-btn" onclick="toggleTheme()" title="Toggle light/dark">🌙</button>
  <button class="tbtn btn-gray" id="settings-btn" onclick="toggleSettingsPopover()" title="Settings">⚙</button>
  <button class="tbtn btn-blue" id="scrape-btn" onclick="runScraper()">↓<span class="tbtn-label"> Sync</span></button>
</div>

<!-- Settings popover -->
<div id="settings-popover" style="display:none;position:fixed;top:48px;right:12px;z-index:2000;
  background:var(--bg2);border:1px solid var(--border2);border-radius:10px;
  padding:16px 18px;width:300px;box-shadow:0 8px 32px rgba(0,0,0,.35)">
  <div style="font-size:12px;font-weight:700;color:var(--text2);margin-bottom:14px">⚙ Settings</div>

  <!-- Background toggle -->
  <div style="display:flex;align-items:center;gap:10px;margin-bottom:10px">
    <span style="font-size:12px;color:var(--text2);flex:1">Auto-start at login</span>
    <label class="bg-toggle">
      <input type="checkbox" id="bg-toggle-input" onchange="bgToggle(this.checked)">
      <span class="bg-toggle-slider"></span>
    </label>
  </div>
  <div id="bg-status-line" style="font-size:11px;color:var(--text3);margin-bottom:14px"></div>

  <!-- Restart instructions -->
  <div style="border-top:1px solid var(--border);padding-top:12px">
    <div style="font-size:11px;font-weight:600;color:var(--text3);margin-bottom:6px">HOW TO RESTART</div>
    <div id="bg-instructions" style="font-size:11px;color:var(--text3);line-height:1.7">Loading…</div>
  </div>
</div>

<!-- Welcome modal (shown once) -->
<div id="welcome-overlay" class="hidden">
  <div id="welcome-modal">
    <p class="welcome-from">by Max (Coxi)</p>
    <h1>Hey, glad you're using it&nbsp;:)</h1>
    <p class="welcome-body">
      I built this because Stud.IP was driving me crazy —
      slow as hell, everything hidden, nothing in one place.<br><br>
      Here you get all your course files, notes, todos,
      AI summaries and flashcards. Runs completely locally,
      nothing leaves your machine.<br><br>
      Good luck this semester ✌️
    </p>
    <button class="welcome-btn" onclick="closeWelcome()">Thanks, let's go</button>
  </div>
</div>

<!-- Mobile sidebar backdrop -->
<div id="sidebar-backdrop" onclick="toggleSidebar()"></div>

<!-- Course context menu -->
<div id="ctx-menu">
  <div class="ctx-item" id="ctx-archive-btn" onclick="ctxArchive()">📦 Archive</div>
  <div class="ctx-item" id="ctx-restore-btn" onclick="ctxRestore()" style="display:none">↩ Restore</div>
  <div class="ctx-item danger" id="ctx-remove-btn" onclick="ctxRemove()">✕ Remove</div>
</div>

<!-- Layout -->
<div id="layout">

  <!-- Sidebar -->
  <div id="sidebar">
    <div id="sidebar-top">
      <input id="sidebar-search" type="text" placeholder="Search courses…" oninput="filterAndRenderSidebar()">
      <div class="filter-pills">
        <button class="pill active" data-filter="all"       onclick="setFilter('all')">All</button>
        <button class="pill"        data-filter="fav"       onclick="setFilter('fav')">⭐ Favourites</button>
      </div>
      <div class="sort-row">
        <label>Sort:</label>
        <select class="sort-select" onchange="setSortAndRender(this.value)">
          <option value="name">Name</option>
          <option value="files">File count</option>
          <option value="progress">Learn progress</option>
          <option value="recent">Last synced</option>
        </select>
      </div>
    </div>
    <div id="course-list"></div>
    <div id="sidebar-credits">Built by <a href="https://maximilianherrmann.com" target="_blank" rel="noopener">Max (Coxi) Herrmann</a> &mdash; free to use</div>
  </div>

  <!-- Sidebar / Main resize divider -->
  <div class="resize-divider" id="divider-sidebar" title="Ziehen zum Anpassen"></div>

  <!-- Main -->
  <div id="main">

    <!-- Course title bar (hidden on home) -->
    <div id="course-title-bar" style="display:none">
      <span id="course-title-text"></span>
      <button id="course-sync-btn" onclick="syncCourse()" title="Sync new folders & files from Stud.IP"
        style="background:none;border:1px solid var(--border);cursor:pointer;color:var(--text3);font-size:11px;padding:2px 8px;border-radius:5px;flex-shrink:0;margin-left:auto">↓ Sync</button>
      <button onclick="goHome()" title="Back to home (Esc)" style="background:none;border:none;cursor:pointer;color:var(--text3);font-size:12px;padding:2px 6px;border-radius:5px;flex-shrink:0">✕</button>
    </div>

    <!-- Tabs (hidden on home) -->
    <div id="tabs" style="display:none">
      <button class="tab active" data-tab="files"   onclick="switchTab('files')">📁 Files</button>
      <button class="tab"        data-tab="info"    onclick="switchTab('info')">ℹ️ Info</button>
      <button class="tab"        data-tab="summary" onclick="switchTab('summary')">📄 Summary</button>
      <button class="tab"        data-tab="learn"   onclick="switchTab('learn')">🧠 Study</button>
      <button class="tab"        data-tab="notes"   onclick="switchTab('notes')">✏️ Notes</button>
      <button class="tab"        data-tab="chat"    onclick="switchTab('chat')">💬 Chat</button>
      <div class="tab-spacer"></div>
      <span id="tab-spinner" style="display:none" class="spin"></span>
    </div>

    <!-- Content panels -->
    <div id="content">

      <!-- Home -->
      <div class="panel active" id="panel-home">
        <div id="sync-tile"></div>
        <!-- Sticky notes pinboard -->
        <div class="section-title" style="margin-bottom:12px">Pinboard</div>
        <div class="pinboard" id="pinboard"></div>
        <!-- To-do widget -->
        <div class="todo-widget">
          <div class="todo-widget-header">
            <span class="todo-widget-title">TO DO</span>
            <button class="todo-sort-btn" id="todo-sort-btn" onclick="todoToggleSort()" title="Sort by due date">📅 Sort</button>
            <button class="todo-clear-btn" id="todo-hide-done-btn" onclick="todoToggleHideDone()" title="Hide/show completed">Hide done</button>
            <button class="todo-clear-btn" onclick="todoClearDone()" title="Remove completed">Remove completed</button>
          </div>
          <div class="todo-list" id="todo-list"></div>
          <button class="todo-add-btn" onclick="todoAddItem()">+ Add task</button>
        </div>

      </div>

      <!-- Courses overview -->
      <div class="panel" id="panel-courses">
        <div id="courses-overview-body"></div>
      </div>

      <!-- Info -->
      <div class="panel" id="panel-info">
        <div id="info-body" style="max-width:680px">
          <div class="empty-state">
            <div class="icon">ℹ️</div>
            <h3>Loading course info…</h3>
          </div>
        </div>
      </div>

      <!-- Files -->
      <div class="panel" id="panel-files">
        <!-- Course-linked todos -->
        <div class="course-todos-wrap" id="course-todos-wrap">
          <div class="course-todos-header">
            TO DO
            <a onclick="goHome()">+ Add on home</a>
          </div>
          <div id="course-todos-list"></div>
        </div>
        <div class="files-layout" id="files-layout">
          <div class="files-list-col" id="files-list-col">
            <div style="display:flex;align-items:center;gap:6px;margin-bottom:6px;">
              <span style="font-size:12px;font-weight:600;color:var(--text3);flex:1">FILES</span>
              <select id="file-sort-select" class="sort-select" style="font-size:10px;padding:2px 4px;max-width:90px" onchange="setFileSort(this.value)">
                <option value="name">Name</option>
                <option value="date">Datum</option>
              </select>
              <button id="selection-toggle-btn" style="font-size:11px;color:var(--text3);background:none;border:1px solid var(--border);border-radius:5px;padding:2px 8px;cursor:pointer;transition:all .15s" onclick="toggleSelectionMode()">Select</button>
            </div>
            <div style="margin-bottom:8px;position:relative">
              <input id="file-search" type="text" placeholder="🔍 Filter files…"
                oninput="filterFileList(this.value)"
                style="width:100%;box-sizing:border-box;background:var(--bg3);border:1px solid var(--border);border-radius:6px;color:var(--text);font-size:12px;padding:5px 28px 5px 8px;outline:none;transition:border-color .15s"
                onfocus="this.style.borderColor='var(--blue)'" onblur="this.style.borderColor='var(--border)'">
              <button id="file-search-clear" onclick="clearFileSearch()" title="Clear"
                style="display:none;position:absolute;right:6px;top:50%;transform:translateY(-50%);background:none;border:none;cursor:pointer;color:var(--text3);font-size:13px;padding:0;line-height:1">✕</button>
            </div>
            <div id="file-list"></div>
            <div class="file-actions">
              <div id="selection-toolbar">
                <span id="sel-count">0 selected</span>
                <div class="sel-toolbar-btns">
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="selectAllFiles()">All</button>
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="selectNoFiles()">None</button>
                </div>
                <div class="sel-toolbar-btns">
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="markSelectedFilesRead(true)">✓ Read</button>
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="markSelectedFilesRead(false)">↺ Unread</button>
                </div>
                <div style="display:flex;gap:4px;align-items:center;margin-bottom:4px">
                  <span style="font-size:10px;color:var(--text3);flex:1">Length</span>
                  <button id="sum-len-short" class="tbtn btn-blue" style="font-size:10px;padding:2px 10px" onclick="setSummaryLength('short')">Short</button>
                  <button id="sum-len-long"  class="tbtn btn-gray" style="font-size:10px;padding:2px 10px" onclick="setSummaryLength('long')">Long</button>
                </div>
                <button class="tbtn btn-blue" id="sel-summarize-btn" style="width:100%;font-size:11px" onclick="generateSummary()">Summarize</button>
              </div>
              <div id="bulk-read-row" style="display:flex;flex-direction:column;gap:4px;margin-bottom:4px">
                <div style="display:flex;gap:6px">
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="markAllFilesRead(true)">✓ All read</button>
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="markAllFilesRead(false)">↺ All unread</button>
                </div>
                <div style="display:flex;gap:6px">
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="downloadCourseZip()">⬇ All files</button>
                  <button class="tbtn btn-gray" style="flex:1;font-size:10px" onclick="downloadCourseZip(true)">⬇ Files + Notes</button>
                </div>
              </div>
            </div>
          </div>
          <!-- Always-visible file list toggle tab -->
          <button id="filelist-toggle-tab" onclick="toggleFileList()" title="Toggle file list">◀</button>
          <!-- Files / Preview resize divider -->
          <div class="resize-divider" id="divider-files" title="Ziehen zum Anpassen"></div>
          <!-- preview-area wraps both preview and notes so fake-fullscreen includes both -->
          <div id="preview-area">
            <div class="files-preview-col" id="files-preview-col">
              <div class="preview-box">
                <div class="preview-header" id="preview-header">
                  <span class="preview-header-name">Select a file to preview</span>
                </div>
                <div id="pdf-find-bar">
                  <input id="pdf-find-input" type="text" placeholder="Search in PDF…"
                    oninput="pdfFindRun(this.value)"
                    onkeydown="if(event.key==='Enter'){event.shiftKey?pdfFindPrev():pdfFindNext()}else if(event.key==='Escape'){pdfFindClose()}">
                  <span id="pdf-find-count"></span>
                  <button class="zoom-btn" onclick="pdfFindPrev()" title="Previous (Shift+Enter)">↑</button>
                  <button class="zoom-btn" onclick="pdfFindNext()" title="Next (Enter)">↓</button>
                  <button onclick="pdfFindClose()" title="Close"
                    style="background:none;border:none;cursor:pointer;color:var(--text3);font-size:14px;padding:2px 6px;border-radius:4px;margin-left:auto;transition:background var(--transition)"
                    onmouseover="this.style.background='var(--bg4)'" onmouseout="this.style.background='none'">✕</button>
                </div>
                <div class="preview-body" id="preview-body">
                  <div class="preview-placeholder">
                    <div class="icon">👆</div>
                    <div>Click a file</div>
                  </div>
                </div>
              </div>
            </div>
            <!-- Resize handle between preview and notes panel -->
            <div class="resize-divider" id="divider-notes" style="display:none" title="Ziehen zum Anpassen"></div>
            <!-- File notes panel (slides in from right) -->
            <div class="files-notes-col" id="files-notes-col">
              <div class="fnotes-header">
                <span class="fnotes-title" id="fnotes-title">Notes</span>
                <span id="fnotes-saved"></span>
                <button class="fnotes-dl-btn" id="fnotes-mode-btn" title="Preview" onclick="toggleFnotesMode()">👁</button>
                <button class="fnotes-dl-btn" title="Download as Markdown" onclick="downloadFileNote()">⬇</button>
                <button class="fnotes-dl-btn" title="Close" onclick="toggleFileNotes()">✕</button>
              </div>
              <textarea id="fnotes-editor" placeholder="Notes for this file… (Markdown)"></textarea>
              <div id="fnotes-rendered"></div>
            </div>
          </div>
        </div>
      </div>

      <!-- Summary -->
      <div class="panel" id="panel-summary">
        <div id="summary-body">
          <div class="empty-state">
            <div class="icon">📄</div>
            <h3>No summary yet</h3>
            <p>Go to "Files" and click "Summarize"</p>
          </div>
        </div>
      </div>

      <!-- Learn / Flashcards -->
      <div class="panel" id="panel-learn">
        <div id="learn-body">
          <div class="empty-state">
            <div class="icon">🧠</div>
            <h3>No flashcards</h3>
            <p>Create a summary first</p>
          </div>
        </div>
      </div>

      <!-- Notes -->
      <div class="panel" id="panel-notes">
        <div class="notes-panel">
          <div class="notes-toolbar">
            <span style="font-size:13px;font-weight:600;color:var(--text2);">Notes</span>
            <div style="flex:1"></div>
            <span id="notes-saved"></span>
            <button class="tbtn btn-gray" id="notes-preview-btn" onclick="toggleNotesPreview()">Preview</button>
          </div>
          <textarea id="notes-editor" placeholder="Your notes, questions, connections… (Markdown supported)"></textarea>
          <div id="notes-preview" class="md-content"></div>
        </div>
      </div>

      <!-- Chat panel -->
      <div class="panel" id="panel-chat">
        <div id="chat-body">
          <div class="empty-state">
            <div class="icon">💬</div>
            <h3>Ask questions</h3>
            <p>Load a course and ask questions about the summary</p>
          </div>
        </div>
      </div>

      <!-- Search results overlay -->
      <div class="panel" id="panel-search">
        <div id="search-results-body"></div>
      </div>

    </div>
  </div>
</div>

<!-- Log box -->
<div id="log-box" style="display:none;flex-direction:column">
  <div id="log-header">
    <span id="log-title">Log</span>
    <button id="log-close" onclick="document.getElementById('log-box').style.display='none'">✕</button>
  </div>
  <div id="log-content"></div>
</div>

<!-- Toast container -->
<div id="toast-container"></div>

<!-- Keyboard shortcuts overlay -->
<div id="shortcuts-overlay" onclick="hideShortcuts()">
  <div class="shortcuts-box" onclick="event.stopPropagation()">
    <h3>⌨️ Keyboard shortcuts</h3>
    <div class="shortcuts-section">Global</div>
    <div class="shortcut-row"><kbd>?</kbd>          <span class="shortcut-desc">Show shortcuts</span></div>
    <div class="shortcut-row"><kbd>⌘K</kbd>         <span class="shortcut-desc">Open command palette</span></div>
    <div class="shortcut-row"><kbd>Esc</kbd>        <span class="shortcut-desc">Go home / close overlay</span></div>
    <div class="shortcuts-section">Flashcards</div>
    <div class="shortcut-row"><kbd>Space</kbd>      <span class="shortcut-desc">Reveal answer</span></div>
    <div class="shortcut-row"><kbd>→ / k</kbd>      <span class="shortcut-desc">Known</span></div>
    <div class="shortcut-row"><kbd>← / u</kbd>      <span class="shortcut-desc">Unknown</span></div>
    <div class="shortcuts-section">Files</div>
    <div class="shortcut-row"><kbd>R</kbd>          <span class="shortcut-desc">Toggle read</span></div>
    <div class="shortcuts-section">Navigation</div>
    <div class="shortcut-row"><kbd>1</kbd>          <span class="shortcut-desc">Files tab</span></div>
    <div class="shortcut-row"><kbd>2</kbd>          <span class="shortcut-desc">Summary tab</span></div>
    <div class="shortcut-row"><kbd>3</kbd>          <span class="shortcut-desc">Study tab</span></div>
    <div class="shortcut-row"><kbd>4</kbd>          <span class="shortcut-desc">Notes tab</span></div>
    <div class="shortcut-row"><kbd>N</kbd>          <span class="shortcut-desc">Toggle file notes</span></div>
    <div class="shortcut-row"><kbd>5</kbd>          <span class="shortcut-desc">Chat tab</span></div>
    <div style="margin-top:20px;text-align:right">
      <button class="tbtn btn-gray" onclick="hideShortcuts()">Close</button>
    </div>
  </div>
</div>

<!-- Confirm modal -->
<div id="confirm-overlay" onclick="hideConfirm()">
  <div class="confirm-box" onclick="event.stopPropagation()">
    <h3 id="confirm-title">Are you sure?</h3>
    <p id="confirm-message">This action cannot be undone.</p>
    <div class="confirm-btns">
      <button class="tbtn btn-gray" onclick="hideConfirm()">Cancel</button>
      <button class="tbtn btn-red" id="confirm-ok">Confirm</button>
    </div>
  </div>
</div>

<!-- Diff modal -->
<div id="diff-overlay" onclick="hideDiffModal()">
  <div id="diff-box" onclick="event.stopPropagation()">
    <div style="display:flex;align-items:center;gap:10px;margin-bottom:14px">
      <h3 style="flex:1;margin:0;font-size:15px;font-weight:700;color:var(--text)">Compare Summaries</h3>
      <button class="tbtn btn-gray" style="padding:3px 10px;font-size:11px" onclick="hideDiffModal()">✕ Close</button>
    </div>
    <div style="display:flex;gap:10px;margin-bottom:12px;flex-wrap:wrap">
      <div style="flex:1;min-width:140px">
        <label style="font-size:11px;color:var(--text3);display:block;margin-bottom:4px">Summary A</label>
        <select id="diff-sel-a" style="width:100%;background:var(--bg3);border:1px solid var(--border);border-radius:5px;color:var(--text);font-size:12px;padding:5px 8px"></select>
      </div>
      <div style="flex:1;min-width:140px">
        <label style="font-size:11px;color:var(--text3);display:block;margin-bottom:4px">Summary B</label>
        <select id="diff-sel-b" style="width:100%;background:var(--bg3);border:1px solid var(--border);border-radius:5px;color:var(--text);font-size:12px;padding:5px 8px"></select>
      </div>
      <div style="display:flex;align-items:flex-end">
        <button class="tbtn btn-blue" style="font-size:12px" onclick="_runDiff()">Compare</button>
      </div>
    </div>
    <div id="diff-view" class="diff-view"></div>
  </div>
</div>

<!-- Command Palette -->
<div id="cmd-overlay" onclick="closeCmdPalette()">
  <div id="cmd-box" onclick="event.stopPropagation()">
    <input id="cmd-input" type="text" placeholder="Search courses or actions…"
      oninput="renderPaletteResults(this.value)"
      onkeydown="handlePaletteKey(event)"
      autocomplete="off" spellcheck="false">
    <div id="cmd-results"></div>
    <div id="cmd-footer">
      <span class="cmd-hint"><kbd>↑↓</kbd> navigieren</span>
      <span class="cmd-hint"><kbd>↵</kbd> open</span>
      <span class="cmd-hint"><kbd>Esc</kbd> close</span>
    </div>
  </div>
</div>

<script>
// ═══════════════════════════════════════════════════════════════════════════
// State
// ═══════════════════════════════════════════════════════════════════════════
let allCourses     = [];
let activeCourse   = null;
let activeTab      = 'home';
let flashState     = { cards: [], index: 0, revealed: false, progress: {}, timerStart: null, timerInterval: null, isGlobal: false };
let notesSaveTimer = null;
let allFilesChecked = true;
let selectionMode = false;
let sidebarFilter  = 'all';
let sidebarSort    = 'name';
let notesPreviewMode = false;

// ═══════════════════════════════════════════════════════════════════════════
// Favorites
// ═══════════════════════════════════════════════════════════════════════════
function applyFontSize() {} // placeholder — font size is fixed for now

// ── Theme ────────────────────────────────────────────────────────────────
function applyTheme(light) {
  document.documentElement.classList.toggle('light', light);
  const btn = document.getElementById('theme-toggle-btn');
  if (btn) btn.textContent = light ? '🌑' : '🌙';
}
function toggleTheme() {
  const isLight = document.documentElement.classList.contains('light');
  localStorage.setItem('theme', isLight ? 'dark' : 'light');
  applyTheme(!isLight);
}
// Apply saved theme immediately (before boot) to avoid flash
applyTheme(localStorage.getItem('theme') === 'light');

// ═══════════════════════════════════════════════════════════════════════════
// User state — persisted server-side so it survives browser/device changes
// ═══════════════════════════════════════════════════════════════════════════
let _usFavs       = [];   // fav_courses
let _usArchived   = [];   // archived_courses
let _usHidden     = [];   // hidden_courses
let _usReadFiles  = {};   // { course_path: [filename, ...] }
let _usPins       = [];   // pins
let _usTodos      = [];   // todos
let _usTodoHideDone = false;

let _usPersistTimer = null;
function _persistUserState() {
  clearTimeout(_usPersistTimer);
  _usPersistTimer = setTimeout(() => {
    fetch('/api/user-state', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({
        fav_courses:    _usFavs,
        archived_courses: _usArchived,
        hidden_courses: _usHidden,
        read_files:     _usReadFiles,
        pins:           _usPins,
        todos:          _usTodos,
        todo_hide_done: _usTodoHideDone,
      })
    });
  }, 300);
}

async function _loadUserState() {
  try {
    const s = await fetch('/api/user-state').then(r => r.json());
    // One-time migration from localStorage if server has no data yet
    const migrateList = (key, fallback) => {
      if (s[key] !== undefined) return s[key];
      try { return JSON.parse(localStorage.getItem(key) || JSON.stringify(fallback)); } catch { return fallback; }
    };
    _usFavs         = migrateList('fav_courses', []);
    _usArchived     = migrateList('archived_courses', []);
    _usHidden       = migrateList('hidden_courses', []);
    _usTodos        = migrateList('todos', []);
    _usPins         = migrateList('pins', []);
    _usTodoHideDone = s.todo_hide_done !== undefined
      ? s.todo_hide_done
      : (localStorage.getItem('todo_hide_done') === '1');
    // Read files: merge all read_files__<course> keys from localStorage
    if (s.read_files !== undefined) {
      _usReadFiles = s.read_files;
    } else {
      _usReadFiles = {};
      for (let i = 0; i < localStorage.length; i++) {
        const k = localStorage.key(i);
        if (k && k.startsWith('read_files__')) {
          const course = k.slice('read_files__'.length);
          try { _usReadFiles[course] = JSON.parse(localStorage.getItem(k) || '[]'); } catch {}
        }
      }
    }
    // Persist migrated data back to server
    _persistUserState();
  } catch (e) {
    console.warn('Could not load user state from server:', e);
  }
}

// ── Read-state helpers ───────────────────────────────────────────────────
function getReadSet(course) {
  return new Set(_usReadFiles[course] || []);
}
function saveReadSet(course, set) {
  _usReadFiles[course] = [...set];
  _persistUserState();
}
function isFileRead(course, filename) { return getReadSet(course).has(filename); }
function setFileRead(course, filename, read) {
  const s = getReadSet(course);
  if (read) s.add(filename); else s.delete(filename);
  saveReadSet(course, s);
  // Update DOM immediately
  document.querySelectorAll(`.file-item[data-filename="${CSS.escape(filename)}"]`).forEach(el => {
    el.classList.toggle('file-read', read);
    const chk = el.querySelector('.file-read-check');
    if (chk) chk.textContent = read ? '✓' : '';
  });
  // Update preview header button
  const btn = document.getElementById('read-toggle-btn');
  if (btn) { btn.textContent = read ? '✓ Read' : 'Mark as read'; btn.classList.toggle('is-read', read); }
}

function getFavorites() { return _usFavs; }
function setFavorites(arr) { _usFavs = arr; _persistUserState(); }
function isFavorite(name) { return _usFavs.includes(name); }
function toggleFavorite(name) {
  _usFavs = _usFavs.includes(name) ? _usFavs.filter(f => f !== name) : [..._usFavs, name];
  _persistUserState();
}

// ───────────────────────── Archive ─────────────────────────
function getArchived() { return _usArchived; }
function setArchived(arr) { _usArchived = arr; _persistUserState(); }
function isArchived(path) { return _usArchived.includes(path); }

// ───────────────────────── Hidden (removed) ─────────────────────────
function getHidden() { return _usHidden; }
function setHidden(arr) { _usHidden = arr; _persistUserState(); }

// ───────────────────────── Context menu ─────────────────────────
let _ctxPath = null;
let _archiveOpen = false;

function showCtxMenu(e, path) {
  e.preventDefault();
  e.stopPropagation();
  _ctxPath = path;
  const archived = isArchived(path);
  document.getElementById('ctx-archive-btn').style.display = archived ? 'none' : '';
  document.getElementById('ctx-restore-btn').style.display = archived ? '' : 'none';
  const menu = document.getElementById('ctx-menu');
  menu.classList.add('visible');
  // Position near cursor, keep inside viewport
  const mw = 170, mh = 90;
  let x = e.clientX, y = e.clientY;
  if (x + mw > window.innerWidth)  x = window.innerWidth  - mw - 8;
  if (y + mh > window.innerHeight) y = window.innerHeight - mh - 8;
  menu.style.left = x + 'px';
  menu.style.top  = y + 'px';
}
function hideCtxMenu() {
  document.getElementById('ctx-menu').classList.remove('visible');
  _ctxPath = null;
}
document.addEventListener('click', hideCtxMenu);
document.addEventListener('keydown', e => { if (e.key === 'Escape') hideCtxMenu(); });

function ctxArchive() {
  if (!_ctxPath) return;
  const arr = getArchived();
  if (!arr.includes(_ctxPath)) arr.push(_ctxPath);
  setArchived(arr);
  if (activeCourse === _ctxPath) goHome();
  _archiveOpen = true; // auto-expand archive after archiving
  hideCtxMenu();
  filterAndRenderSidebar();
}
function ctxRestore() {
  if (!_ctxPath) return;
  setArchived(getArchived().filter(p => p !== _ctxPath));
  hideCtxMenu();
  filterAndRenderSidebar();
}
function ctxRemove() {
  if (!_ctxPath) return;
  if (!confirm(`Permanently remove "${_ctxPath}" from the sidebar?`)) return;
  const h = getHidden();
  if (!h.includes(_ctxPath)) h.push(_ctxPath);
  setHidden(h);
  // also remove from archive if there
  setArchived(getArchived().filter(p => p !== _ctxPath));
  if (activeCourse === _ctxPath) goHome();
  hideCtxMenu();
  filterAndRenderSidebar();
}

// ═══════════════════════════════════════════════════════════════════════════
// Boot
// ═══════════════════════════════════════════════════════════════════════════
// courseTree = raw API response (may include groups)
// allCourses = flat list of all individual courses (for stats/home)
let courseTree = [];
let _lastSyncGlobal = null;

async function _fetchCourses() {
  const res = await fetch('/api/courses').then(r => r.json());
  // support both old array shape and new {tree, last_sync} shape
  if (Array.isArray(res)) { _lastSyncGlobal = null; return res; }
  _lastSyncGlobal = res.last_sync || null;
  return res.tree || [];
}

function flattenTree(tree) {
  const flat = [];
  for (const item of tree) {
    if (item.is_group) flat.push(...item.courses);
    else flat.push(item);
  }
  return flat;
}

async function boot() {
  await _loadUserState();
  courseTree = await _fetchCourses();
  allCourses = flattenTree(courseTree);
  filterAndRenderSidebar();
  renderHome(allCourses);
  initResizeDividers();
  applyFontSize();

  // Init pinboard & to-do list
  pinRender();
  if (_todoLoad().length === 0) _todoSave([{ text: '', done: false }]);
  todoRender();

  // Restore last session state
  const lastCourse = localStorage.getItem('last_course');
  const lastTab    = localStorage.getItem('last_tab') || 'files';
  if (lastCourse && allCourses.find(c => c.path === lastCourse)) {
    activeCourse = lastCourse;
    filterAndRenderSidebar();
    document.getElementById('tabs').style.display = 'flex';
    document.getElementById('course-title-bar').style.display = 'flex';
    document.getElementById('course-title-text').textContent = lastCourse.split('/').pop();
    switchTab(lastTab);
    if (lastTab === 'files') loadFiles();
  }
}

// ═══════════════════════════════════════════════════════════════════════════
// Sidebar
// ═══════════════════════════════════════════════════════════════════════════
function setFilter(f) {
  sidebarFilter = f;
  document.querySelectorAll('.pill').forEach(p => p.classList.toggle('active', p.dataset.filter === f));
  filterAndRenderSidebar();
}

function setSortAndRender(s) {
  sidebarSort = s;
  filterAndRenderSidebar();
}

function filterCourses(courses) {
  const q = document.getElementById('sidebar-search').value.toLowerCase();
  const hidden   = getHidden();
  const archived = getArchived();
  let result = courses.filter(c =>
    !hidden.includes(c.path) && !archived.includes(c.path) &&
    (c.name.toLowerCase().includes(q) || (c.path||'').toLowerCase().includes(q))
  );
  if (sidebarFilter === 'summary')   result = result.filter(c => c.has_summary);
  if (sidebarFilter === 'nosummary') result = result.filter(c => !c.has_summary);
  if (sidebarFilter === 'fav')       result = result.filter(c => isFavorite(c.path));
  if (sidebarSort === 'files')    result.sort((a, b) => b.file_count - a.file_count);
  if (sidebarSort === 'progress') result.sort((a, b) => {
    const pa = a.progress.total ? a.progress.known / a.progress.total : 0;
    const pb = b.progress.total ? b.progress.known / b.progress.total : 0;
    return pb - pa;
  });
  if (sidebarSort === 'recent') result.sort((a, b) => (b.summary_age || 0) - (a.summary_age || 0));
  return result;
}

function filterAndRenderSidebar() {
  const q        = document.getElementById('sidebar-search').value.toLowerCase();
  const hidden   = getHidden();
  const archived = getArchived();

  // Build filtered tree (exclude hidden + archived from main list)
  const filteredTree = [];
  for (const item of courseTree) {
    if (item.is_group) {
      const filtered = filterCourses(item.courses);
      if (filtered.length || item.name.toLowerCase().includes(q)) {
        filteredTree.push({ ...item, courses: filtered.length ? filtered : item.courses.filter(() => !q) });
      }
    } else {
      if (hidden.includes(item.path) || archived.includes(item.path)) continue;
      const matches = item.name.toLowerCase().includes(q);
      const passFilter = (sidebarFilter === 'all') ||
        (sidebarFilter === 'summary'   && item.has_summary) ||
        (sidebarFilter === 'nosummary' && !item.has_summary) ||
        (sidebarFilter === 'fav'       && isFavorite(item.path));
      if (matches && passFilter) filteredTree.push(item);
    }
  }

  renderSidebar(filteredTree);
}

// Track which groups are collapsed (default: open)
const collapsedGroups = new Set(JSON.parse(localStorage.getItem('collapsed_groups') || '[]'));

function toggleGroup(groupName) {
  if (collapsedGroups.has(groupName)) collapsedGroups.delete(groupName);
  else collapsedGroups.add(groupName);
  localStorage.setItem('collapsed_groups', JSON.stringify([...collapsedGroups]));
  filterAndRenderSidebar();
}

function renderSidebar(tree) {
  const el = document.getElementById('course-list');

  function courseHTML(c, indent = false, inArchive = false) {
    const pct = c.progress.total ? Math.round(c.progress.known / c.progress.total * 100) : 0;
    const fav = isFavorite(c.path);
    const lastStudied = c.progress.last_studied
      ? new Date(c.progress.last_studied).toLocaleDateString('de-DE', {day:'2-digit', month:'2-digit'})
      : null;
    const readCount = c.file_count > 0 ? getReadSet(c.path).size : 0;
    const allRead = c.file_count > 0 && readCount >= c.file_count;
    const archiveClass = inArchive ? ' archived' : '';
    return `
    <div class="citem${archiveClass} ${activeCourse === c.path ? 'active' : ''}" data-course="${esc(c.path)}" onclick="selectCourseFromEl(this)"
         oncontextmenu="showCtxMenu(event,'${esc(c.path)}')"
         style="${indent ? 'padding-left:22px' : ''}">
      <div class="citem-dot ${c.has_summary ? 'dot-ok' : 'dot-missing'}"></div>
      <div class="citem-body">
        <div class="citem-name" title="${esc(c.name)}">${esc(c.name)}${c.new_files ? `<span class="new-badge">+${c.new_files}</span>` : ''}${allRead ? `<span style="color:var(--green);font-size:10px;margin-left:4px" title="All files read">✓</span>` : ''}</div>
        <div class="citem-meta">
          <span>${c.file_count} file${c.file_count !== 1 ? 's' : ''}</span>
          ${readCount > 0 && !allRead ? `<span style="color:var(--text3)">${readCount}/${c.file_count} read</span>` : ''}
          ${c.new_files && sidebarSort === 'recent' ? `<span style="color:var(--blue);font-weight:600">+${c.new_files} new</span>` : ''}
          ${c.has_summary ? `<span style="color:var(--green)">✓</span>` : ''}
          ${c.has_notes   ? `<span style="color:var(--purple)">📝</span>` : ''}
          ${lastStudied   ? `<span title="Last studied">🕐 ${lastStudied}</span>` : ''}
        </div>
        ${c.progress.total ? `<div class="progress-mini"><div class="progress-mini-fill" style="width:${pct}%"></div></div>` : ''}
      </div>
      ${inArchive
        ? `<button class="fav-btn" title="Restore" onclick="event.stopPropagation(); setArchived(getArchived().filter(p=>p!=='${esc(c.path)}')); filterAndRenderSidebar()">↩</button>`
        : `<button class="fav-btn ${fav ? 'is-fav' : ''}" title="${fav ? 'Remove favourite' : 'Mark as favourite'}" onclick="event.stopPropagation(); handleFavClick('${esc(c.path)}')">⭐</button>`
      }
    </div>`;
  }

  function groupHTML(item) {
    const collapsed = collapsedGroups.has(item.name);
    const total    = item.courses.length;
    const withSum  = item.courses.filter(c => c.has_summary).length;
    const semClass = item.is_semester ? 'group-header semester-header' : 'group-header';
    const icon     = item.is_semester ? '🎓 ' : '';
    return `
    <div class="${semClass}" onclick="toggleGroup('${esc(item.name)}')">
      <span class="group-chevron">${collapsed ? '▶' : '▼'}</span>
      <span class="group-name">${icon}${esc(item.name)}</span>
      <span class="group-meta">${withSum}/${total}</span>
    </div>
    ${collapsed ? '' : item.courses.map(c => courseHTML(c, true)).join('')}`;
  }

  // Separate favorites from tree
  const favPaths = getFavorites();
  const favCourses = allCourses.filter(c => favPaths.includes(c.path) && !getArchived().includes(c.path) && !getHidden().includes(c.path));

  let html = '';
  if (favCourses.length) {
    html += `<div class="sidebar-section-label">⭐ Favoriten</div>`;
    html += favCourses.map(c => courseHTML(c)).join('');
    html += `<div class="sidebar-section-label" style="margin-top:6px">Courses</div>`;
  } else {
    html += `<div class="sidebar-section-label">Courses</div>`;
  }

  if (!tree.length) {
    html += `<div style="color:var(--text3);font-size:12px;padding:16px 12px">No courses found</div>`;
  } else {
    for (const item of tree) {
      html += item.is_group ? groupHTML(item) : courseHTML(item);
    }
  }

  // Archive section
  const archivedPaths = getArchived();
  const archivedCoursesList = allCourses.filter(c => archivedPaths.includes(c.path));
  if (archivedCoursesList.length) {
    const archCollapsed = !_archiveOpen;
    html += `
    <div class="archive-header" onclick="_archiveOpen=!_archiveOpen; filterAndRenderSidebar()">
      <span class="group-chevron">${archCollapsed ? '▶' : '▼'}</span>
      <span>📦 Archiv</span>
      <span style="margin-left:auto;font-size:10px;opacity:.6">${archivedCoursesList.length}</span>
    </div>
    ${archCollapsed ? '' : archivedCoursesList.map(c => courseHTML(c, false, true)).join('')}`;
  }

  el.innerHTML = html;
}

function selectCourseFromEl(el) { selectCourse(el.dataset.course); }

function handleFavClick(path) {
  toggleFavorite(path);
  filterAndRenderSidebar();
}

// ═══════════════════════════════════════════════════════════════════════════
// Sticky notes pinboard
// ═══════════════════════════════════════════════════════════════════════════
const PIN_COLORS = ['#fef08a','#bbf7d0','#bfdbfe','#fecaca','#f5d0fe','#fed7aa'];

function _pinLoad() { return _usPins; }
function _pinSave(pins) { _usPins = pins; _persistUserState(); }

let _pinDragIdx = null;

function pinRender() {
  const pins = _pinLoad();
  const board = document.getElementById('pinboard');
  if (!board) return;
  board.innerHTML = '';
  pins.forEach((p, i) => {
    const div = document.createElement('div');
    div.className = 'sticky';
    div.style.background = p.color || PIN_COLORS[0];
    div.draggable = true;
    div.dataset.idx = i;
    div.innerHTML = `<button class="sticky-del" onclick="pinDelete(${i})" title="Remove">✕</button>
      <textarea placeholder="Write something…" oninput="pinEdit(${i},this.value)">${p.text || ''}</textarea>`;
    div.addEventListener('dragstart', e => {
      _pinDragIdx = i;
      div.style.opacity = '0.4';
      e.dataTransfer.effectAllowed = 'move';
    });
    div.addEventListener('dragend', () => { div.style.opacity = ''; });
    div.addEventListener('dragover', e => {
      e.preventDefault();
      e.dataTransfer.dropEffect = 'move';
      div.style.outline = '2px dashed rgba(0,0,0,.3)';
    });
    div.addEventListener('dragleave', () => { div.style.outline = ''; });
    div.addEventListener('drop', e => {
      e.preventDefault();
      div.style.outline = '';
      const toIdx = parseInt(div.dataset.idx);
      if (_pinDragIdx === null || _pinDragIdx === toIdx) return;
      const pins2 = _pinLoad();
      const [moved] = pins2.splice(_pinDragIdx, 1);
      pins2.splice(toIdx, 0, moved);
      _pinSave(pins2);
      pinRender();
    });
    board.appendChild(div);
  });
  const add = document.createElement('div');
  add.className = 'pinboard-add';
  add.title = 'Add sticky note';
  add.textContent = '+';
  add.onclick = pinAdd;
  board.appendChild(add);
}

function pinAdd() {
  const pins = _pinLoad();
  const color = PIN_COLORS[pins.length % PIN_COLORS.length];
  pins.push({ text: '', color });
  _pinSave(pins);
  pinRender();
  const board = document.getElementById('pinboard');
  const notes = board.querySelectorAll('.sticky textarea');
  if (notes.length) notes[notes.length - 1].focus();
}

function pinDelete(i) {
  const pins = _pinLoad();
  pins.splice(i, 1);
  _pinSave(pins);
  pinRender();
}

function pinEdit(i, val) {
  const pins = _pinLoad();
  if (pins[i]) { pins[i].text = val; _pinSave(pins); }
}

// ═══════════════════════════════════════════════════════════════════════════
// To-do list
// ═══════════════════════════════════════════════════════════════════════════
function _todoLoad() { return _usTodos; }
function _todoSave(items) { _usTodos = items; _persistUserState(); }

function todoToggleHideDone() {
  _usTodoHideDone = !_usTodoHideDone;
  _persistUserState();
  todoRender();
  if (activeCourse) todoRenderCourse(activeCourse);
}

function _todoUpdateHideBtn() {
  const btn = document.getElementById('todo-hide-done-btn');
  if (btn) btn.textContent = _usTodoHideDone ? 'Show done' : 'Hide done';
}

function _todoDueLabel(due) {
  if (!due) return null;
  const today = new Date(); today.setHours(0,0,0,0);
  const d     = new Date(due + 'T00:00:00');
  const diff  = Math.round((d - today) / 86400000);
  if (diff < 0)  return { cls: 'overdue',  label: `${Math.abs(diff)}d overdue` };
  if (diff === 0) return { cls: 'today',    label: 'Today' };
  if (diff === 1) return { cls: 'upcoming', label: 'Tomorrow' };
  return { cls: 'upcoming', label: `${diff}d` };
}

function _todoItemHTML(item, i, opts = {}) {
  const courseName = item.course ? item.course.split('/').pop() : null;
  const linkPill   = courseName
    ? `<span class="todo-course-pill" onclick="event.stopPropagation();todoUnlinkCourse(${i})" title="Unlink from course">${esc(courseName)} ✕</span>`
    : `<button class="todo-link-btn" onclick="event.stopPropagation();todoShowCoursePicker(${i},this)" title="Link to course">@</button>`;
  const dueInfo = _todoDueLabel(item.due);
  const duePill = opts.courseView ? (dueInfo ? `<span class="todo-due ${dueInfo.cls}">${dueInfo.label}</span>` : '') : `
    <input type="date" class="todo-due-input" value="${item.due || ''}"
      title="Set due date"
      onchange="todoDue(${i},this.value)"
      style="opacity:0;position:absolute;pointer-events:none" id="todo-date-${i}">
    ${dueInfo
      ? `<button class="todo-due ${dueInfo.cls}" onclick="event.stopPropagation();document.getElementById('todo-date-${i}').showPicker?document.getElementById('todo-date-${i}').showPicker():document.getElementById('todo-date-${i}').click()" title="Change due date">${dueInfo.label}</button>`
      : `<button class="todo-due upcoming" onclick="event.stopPropagation();document.getElementById('todo-date-${i}').showPicker?document.getElementById('todo-date-${i}').showPicker():document.getElementById('todo-date-${i}').click()" title="Set due date" style="opacity:0;transition:opacity .15s" onmouseenter="this.style.opacity=1" onmouseleave="this.style.opacity=0">📅</button>`
    }`;
  return `
    <div class="todo-item" data-i="${i}" style="position:relative">
      ${opts.courseView ? '' : '<span class="todo-drag-handle" draggable="false">⠿</span>'}
      <input type="checkbox" class="todo-cb" ${item.done ? 'checked' : ''} onchange="todoToggle(${i},${!!opts.courseView})">
      <textarea class="todo-text${item.done ? ' done' : ''}" rows="1"
        onchange="todoEdit(${i}, this.value)"
        oninput="this.style.height='auto';this.style.height=this.scrollHeight+'px';todoEdit(${i},this.value)"
        onkeydown="${opts.courseView ? '' : `todoKeydown(event,${i})`}"
        ${opts.courseView ? 'readonly' : ''}
      >${esc(item.text)}</textarea>
      ${duePill}
      ${opts.courseView ? '' : linkPill}
      ${opts.courseView ? '' : `<button class="todo-del-btn" onclick="todoDelete(${i})" title="Delete">✕</button>`}
    </div>`;
}

let _todoSortByDue = localStorage.getItem('todo_sort_due') === '1';

function todoToggleSort() {
  _todoSortByDue = !_todoSortByDue;
  localStorage.setItem('todo_sort_due', _todoSortByDue ? '1' : '');
  const btn = document.getElementById('todo-sort-btn');
  if (btn) btn.classList.toggle('active', _todoSortByDue);
  todoRender();
}

function todoDue(i, val) {
  const items = _todoLoad();
  if (!items[i]) return;
  items[i].due = val || null;
  _todoSave(items);
  todoRender();
  if (activeCourse) todoRenderCourse(activeCourse);
}

function _todoSortedItems(items) {
  if (!_todoSortByDue) return items.map((item, i) => ({ item, i }));
  // Sort: overdue first, then by date ascending, then undated at end, done at very end
  return items
    .map((item, i) => ({ item, i }))
    .sort((a, b) => {
      if (a.item.done !== b.item.done) return a.item.done ? 1 : -1;
      const da = a.item.due || '9999-99-99';
      const db = b.item.due || '9999-99-99';
      return da.localeCompare(db);
    });
}

let _todoDragIdx = null;

function _todoDragWire(list) {
  list.querySelectorAll('.todo-item').forEach(row => {
    const handle = row.querySelector('.todo-drag-handle');
    if (!handle) return;
    handle.addEventListener('mousedown', () => { row.draggable = true; });
    row.addEventListener('dragstart', e => {
      _todoDragIdx = parseInt(row.dataset.i);
      row.style.opacity = '0.4';
      e.dataTransfer.effectAllowed = 'move';
    });
    row.addEventListener('dragend', () => {
      row.style.opacity = '';
      row.draggable = false;
      list.querySelectorAll('.todo-item').forEach(r => r.classList.remove('drag-over'));
    });
    row.addEventListener('dragover', e => {
      e.preventDefault();
      list.querySelectorAll('.todo-item').forEach(r => r.classList.remove('drag-over'));
      row.classList.add('drag-over');
    });
    row.addEventListener('dragleave', () => row.classList.remove('drag-over'));
    row.addEventListener('drop', e => {
      e.preventDefault();
      row.classList.remove('drag-over');
      const toIdx = parseInt(row.dataset.i);
      if (_todoDragIdx === null || _todoDragIdx === toIdx) return;
      const items = _todoLoad();
      const [moved] = items.splice(_todoDragIdx, 1);
      items.splice(toIdx, 0, moved);
      _todoSave(items);
      todoRender();
      if (activeCourse) todoRenderCourse(activeCourse);
    });
  });
}

function todoRender() {
  const items = _todoLoad();
  const list = document.getElementById('todo-list');
  if (!list) return;
  list.innerHTML = _todoSortedItems(items)
    .filter(({ item }) => !(_usTodoHideDone && item.done))
    .map(({ item, i }) => _todoItemHTML(item, i))
    .join('');
  list.querySelectorAll('.todo-text').forEach(t => { t.style.height = 'auto'; t.style.height = t.scrollHeight + 'px'; });
  _todoDragWire(list);
  _todoUpdateHideBtn();
  const btn = document.getElementById('todo-sort-btn');
  if (btn) btn.classList.toggle('active', _todoSortByDue);
}

function todoRenderCourse(coursePath) {
  const items = _todoLoad();
  const linked = items
    .map((item, i) => ({ item, i }))
    .filter(({ item }) => item.course === coursePath && !(_usTodoHideDone && item.done));
  const wrap = document.getElementById('course-todos-wrap');
  if (!wrap) return;
  // Show if there are linked todos (even if all hidden, keep header visible with toggle)
  const allLinked = items.filter(item => item.course === coursePath);
  if (!allLinked.length) { wrap.style.display = 'none'; return; }
  wrap.style.display = 'block';
  document.getElementById('course-todos-list').innerHTML =
    linked.map(({ item, i }) => _todoItemHTML(item, i, { courseView: true })).join('');
  wrap.querySelectorAll('.todo-text').forEach(t => { t.style.height = 'auto'; t.style.height = t.scrollHeight + 'px'; });
}

function todoAddItem(text = '', focus = true) {
  const items = _todoLoad();
  items.push({ text, done: false });
  _todoSave(items);
  todoRender();
  if (focus) {
    const inputs = document.querySelectorAll('#todo-list .todo-text');
    if (inputs.length) { const last = inputs[inputs.length - 1]; last.focus(); last.setSelectionRange(last.value.length, last.value.length); }
  }
}

function todoToggle(i, fromCourseView = false) {
  const items = _todoLoad();
  if (!items[i]) return;
  items[i].done = !items[i].done;
  _todoSave(items);
  if (fromCourseView && activeCourse) todoRenderCourse(activeCourse); else todoRender();
  // Also refresh the other view if visible
  if (fromCourseView) todoRender();
  else if (activeCourse) todoRenderCourse(activeCourse);
}

function todoEdit(i, val) {
  const items = _todoLoad();
  if (!items[i]) return;
  items[i].text = val;
  _todoSave(items);
}

function todoDelete(i) {
  const items = _todoLoad();
  items.splice(i, 1);
  _todoSave(items);
  todoRender();
  if (activeCourse) todoRenderCourse(activeCourse);
}

function todoClearDone() {
  _todoSave(_todoLoad().filter(t => !t.done));
  todoRender();
}

function todoKeydown(e, i) {
  if (e.key === 'Enter' && !e.shiftKey) {
    e.preventDefault();
    const items = _todoLoad();
    items.splice(i + 1, 0, { text: '', done: false });
    _todoSave(items);
    todoRender();
    const inputs = document.querySelectorAll('#todo-list .todo-text');
    if (inputs[i + 1]) { inputs[i + 1].focus(); }
  }
  if (e.key === 'Backspace') {
    const items = _todoLoad();
    if (items[i] && items[i].text === '' && items.length > 1) {
      e.preventDefault();
      items.splice(i, 1);
      _todoSave(items);
      todoRender();
      const inputs = document.querySelectorAll('#todo-list .todo-text');
      const prev = inputs[Math.max(0, i - 1)];
      if (prev) { prev.focus(); prev.setSelectionRange(prev.value.length, prev.value.length); }
    }
  }
}

function todoShowCoursePicker(i, btnEl) {
  document.querySelectorAll('.todo-course-picker').forEach(p => p.remove());
  const picker = document.createElement('div');
  picker.className = 'todo-course-picker';
  const courses = allCourses.map(c =>
    `<div class="tcp-item" data-path="${esc(c.path)}" onclick="todoLinkCourse(${i},'${esc(c.path)}');this.closest('.todo-course-picker').remove()">${esc(c.name.split('/').pop() || c.name)}</div>`
  ).join('');
  picker.innerHTML = `<div class="tcp-search-wrap"><input class="tcp-search" placeholder="Search course…" oninput="tcpFilter(this)" onkeydown="tcpKeydown(event,${i})" autofocus></div><div class="tcp-list">${courses}</div>`;
  btnEl.closest('.todo-item').appendChild(picker);
  setTimeout(() => picker.querySelector('.tcp-search')?.focus(), 10);
  setTimeout(() => document.addEventListener('click', function h(e) {
    if (!picker.contains(e.target)) { picker.remove(); document.removeEventListener('click', h); }
  }), 10);
}

function tcpFilter(input) {
  const q = input.value.toLowerCase();
  const picker = input.closest('.todo-course-picker');
  picker.querySelectorAll('.tcp-item').forEach(el => {
    el.style.display = el.textContent.toLowerCase().includes(q) ? '' : 'none';
  });
  // Reset highlight
  picker.querySelectorAll('.tcp-item.tcp-active').forEach(el => el.classList.remove('tcp-active'));
}

function tcpKeydown(e, i) {
  const picker = e.target.closest('.todo-course-picker');
  const visible = [...picker.querySelectorAll('.tcp-item')].filter(el => el.style.display !== 'none');
  const cur = picker.querySelector('.tcp-item.tcp-active');
  let idx = visible.indexOf(cur);
  if (e.key === 'ArrowDown') {
    e.preventDefault();
    const next = visible[Math.min(idx + 1, visible.length - 1)];
    if (cur) cur.classList.remove('tcp-active');
    if (next) { next.classList.add('tcp-active'); next.scrollIntoView({ block: 'nearest' }); }
  } else if (e.key === 'ArrowUp') {
    e.preventDefault();
    const prev = visible[Math.max(idx - 1, 0)];
    if (cur) cur.classList.remove('tcp-active');
    if (prev) { prev.classList.add('tcp-active'); prev.scrollIntoView({ block: 'nearest' }); }
  } else if (e.key === 'Enter') {
    e.preventDefault();
    const active = picker.querySelector('.tcp-item.tcp-active') || visible[0];
    if (active) { todoLinkCourse(i, active.dataset.path); picker.remove(); }
  } else if (e.key === 'Escape') {
    picker.remove();
  }
}

function todoLinkCourse(i, path) {
  const items = _todoLoad();
  if (!items[i]) return;
  items[i].course = path;
  _todoSave(items);
  todoRender();
  if (activeCourse) todoRenderCourse(activeCourse);
}

function todoUnlinkCourse(i) {
  const items = _todoLoad();
  if (!items[i]) return;
  delete items[i].course;
  _todoSave(items);
  todoRender();
  if (activeCourse) todoRenderCourse(activeCourse);
}

// ═══════════════════════════════════════════════════════════════════════════
// Home
// ═══════════════════════════════════════════════════════════════════════════
function renderHome(courses) {
  // Sync tile — prefer per-course _last_sync, fall back to courses.json mtime
  const lastSync = Math.max(0, ...courses.map(c => c.sync_age || 0)) || _lastSyncGlobal || 0;
  const syncTile = document.getElementById('sync-tile');
  if (syncTile) {
    if (lastSync) {
      const now = Math.floor(Date.now() / 1000);
      const diff = now - lastSync;
      let label;
      if (diff < 60)           label = 'just now';
      else if (diff < 3600)    label = `${Math.floor(diff/60)}m ago`;
      else if (diff < 86400)   label = `${Math.floor(diff/3600)}h ago`;
      else                     label = `${Math.floor(diff/86400)}d ago`;

      // Collect all newly downloaded files across courses
      const newItems = [];
      for (const c of courses) {
        if (c.new_file_names && c.new_file_names.length) {
          for (const f of c.new_file_names) {
            newItems.push({ course: c.name, file: f, path: c.path });
          }
        }
      }
      const MAX_SHOWN = 8;
      let filesHtml = '';
      if (newItems.length) {
        const shown = newItems.slice(0, MAX_SHOWN);
        filesHtml = `<div id="sync-tile-files">` +
          shown.map(item => `
            <div class="sync-file" onclick="selectCourse('${esc(item.path)}')" style="cursor:pointer" title="${esc(item.course)}">
              <span class="sync-file-course">${esc(item.course.split('/').pop())}</span>
              <span style="color:var(--border2)">›</span>
              <span class="sync-file-name">${esc(item.file)}</span>
            </div>`).join('') +
          (newItems.length > MAX_SHOWN ? `<div style="font-size:10px;color:var(--text3);margin-top:2px">+${newItems.length - MAX_SHOWN} more</div>` : '') +
          `</div>`;
      }

      syncTile.innerHTML = `<div id="sync-tile-header"><span class="sync-dot"></span>Last sync: ${label}${newItems.length ? ` &middot; <span style="color:var(--blue)">${newItems.length} new file${newItems.length > 1 ? 's' : ''}</span>` : ''}</div>${filesHtml}`;
      syncTile.style.display = '';
    } else {
      syncTile.style.display = 'none';
    }
  }
}

function buildRecommendations(courses) {
  const recs = [];
  const withSummary = courses.filter(c => c.has_summary);

  // 1. Kurse mit neuen Dateien seit letzter Zusammenfassung
  const newFiles = courses.filter(c => c.new_files > 0);
  if (newFiles.length) {
    const c = newFiles.sort((a, b) => b.new_files - a.new_files)[0];
    recs.push({ path: c.path, name: c.name, reason: '📂 New files', meta: `${c.new_files} new file${c.new_files > 1 ? 's' : ''} since last summary`, accent: 'var(--yellow)' });
  }

  // 2. Kurs mit niedrigstem Lernfortschritt (der Lernkarten hat)
  const withProgress = withSummary.filter(c => c.progress.total > 0);
  if (withProgress.length) {
    const c = withProgress.sort((a, b) => {
      const pa = a.progress.known / a.progress.total;
      const pb = b.progress.known / b.progress.total;
      return pa - pb;
    })[0];
    const pct = Math.round(c.progress.known / c.progress.total * 100);
    if (pct < 80) recs.push({ path: c.path, name: c.name, reason: '📉 Behind on studying', meta: `${pct}% known (${c.progress.known}/${c.progress.total} cards)`, accent: 'var(--red)' });
  }

  // 3. Kurs ohne Zusammenfassung aber mit Dateien
  const noSummary = courses.filter(c => !c.has_summary && c.file_count > 0);
  if (noSummary.length) {
    const c = noSummary.sort((a, b) => b.file_count - a.file_count)[0];
    recs.push({ path: c.path, name: c.name, reason: '⚠️ No summary', meta: `${c.file_count} file${c.file_count > 1 ? 's' : ''} available`, accent: 'var(--orange)' });
  }

  // 4. Am längsten nicht gelernt (aus denen mit Lernkarten)
  const withLast = withProgress.filter(c => c.progress.last_studied);
  if (withLast.length) {
    const c = withLast.sort((a, b) => new Date(a.progress.last_studied) - new Date(b.progress.last_studied))[0];
    const days = Math.floor((Date.now() - new Date(c.progress.last_studied)) / 86400000);
    if (days >= 3) recs.push({ path: c.path, name: c.name, reason: '⏰ Not studied recently', meta: `Last studied ${days} day${days > 1 ? 's' : ''} ago`, accent: 'var(--purple)' });
  }

  return recs.slice(0, 3);
}

function goHome() {
  activeCourse = null;
  document.getElementById('tabs').style.display = 'none';
  document.getElementById('course-title-bar').style.display = 'none';
  showPanel('home');
  filterAndRenderSidebar();
}

function toggleSettingsPopover() {
  const pop = document.getElementById('settings-popover');
  const open = pop.style.display === 'none';
  pop.style.display = open ? 'block' : 'none';
  if (open) loadBgServiceCard();
}
document.addEventListener('click', e => {
  const pop = document.getElementById('settings-popover');
  if (pop && pop.style.display !== 'none' &&
      !pop.contains(e.target) && !document.getElementById('settings-btn').contains(e.target))
    pop.style.display = 'none';
});

async function loadBgServiceCard() {
  const instr  = document.getElementById('bg-instructions');
  const toggle = document.getElementById('bg-toggle-input');
  const status = document.getElementById('bg-status-line');
  if (!instr || !toggle) return;
  let data;
  try { data = await (await fetch('/api/background-status')).json(); }
  catch { instr.textContent = 'Could not load status.'; return; }

  toggle.checked = data.enabled;
  if (status) status.innerHTML = data.enabled
    ? '<span style="color:var(--green)">●</span> Auto-start is <strong>on</strong> — starts at every login automatically.'
    : '<span style="color:var(--text3)">●</span> Auto-start is <strong>off</strong> — only runs in the current session.';

  instr.innerHTML = `
    1. Open Terminal<br>
    2. Run: <code onclick="navigator.clipboard.writeText(this.textContent);this.style.outline='1px solid var(--blue)';setTimeout(()=>this.style.outline='',800)" title="Click to copy">${data.start_cmd}</code><br>
    3. Open: <code onclick="navigator.clipboard.writeText('http://localhost:5001');this.style.outline='1px solid var(--blue)';setTimeout(()=>this.style.outline='',800)" title="Click to copy">http://localhost:5001</code>`;
}

async function bgToggle(enabled) {
  const instr = document.getElementById('bg-instructions');
  const status = document.getElementById('bg-status-line');
  if (status) status.innerHTML = '<span style="color:var(--text3)">Updating…</span>';
  try {
    const res = await fetch(enabled ? '/api/background-enable' : '/api/background-disable', { method: 'POST' });
    const data = await res.json();
    if (!data.ok && instr) instr.innerHTML = `<span style="color:var(--red)">Error: ${data.error}</span>`;
    else loadBgServiceCard();
  } catch(e) {
    if (instr) instr.innerHTML = `<span style="color:var(--red)">Error: ${e}</span>`;
  }
}

function goCoursesOverview() {
  activeCourse = null;
  document.getElementById('tabs').style.display = 'none';
  document.getElementById('course-title-bar').style.display = 'none';
  showPanel('courses');
  filterAndRenderSidebar();
  renderCoursesOverview();
}

function renderCoursesOverview() {
  const body = document.getElementById('courses-overview-body');
  if (!body) return;

  // Accent colors cycling per semester
  const ACCENTS = ['var(--blue)','var(--purple)','var(--green)','var(--yellow)','var(--orange)','var(--red)'];

  // Separate semesters from ungrouped
  const semesters = courseTree.filter(item => item.is_group && item.is_semester);
  const ungrouped = courseTree.filter(item => !item.is_group);
  const otherGroups = courseTree.filter(item => item.is_group && !item.is_semester);

  const archived = getArchived();
  const hidden   = getHidden();

  function cardHTML(c) {
    if (archived.includes(c.path) || hidden.includes(c.path)) return '';
    const newLine = c.new_files
      ? `<div style="font-size:11px;color:var(--yellow);margin-top:3px">+${c.new_files} new since last sync</div>`
      : '';
    return `
      <div class="co-card" onclick="selectCourse('${esc(c.path)}')" title="${esc(c.name)}">
        <div class="co-card-accent" style="background:ACCENT"></div>
        <div class="co-card-name">${esc(c.name)}</div>
        ${newLine}
      </div>`;
  }

  let html = '';

  semesters.forEach((sem, si) => {
    const accent = ACCENTS[si % ACCENTS.length];
    const cards = sem.courses.map(c => cardHTML(c).replace('ACCENT', accent)).filter(Boolean).join('');
    if (!cards) return;
    html += `<div class="co-semester">
      <div class="co-semester-title">${esc(sem.name)}</div>
      <div class="co-grid">${cards}</div>
    </div>`;
  });

  otherGroups.forEach((grp, gi) => {
    const accent = ACCENTS[(semesters.length + gi) % ACCENTS.length];
    const cards = grp.courses.map(c => cardHTML(c).replace('ACCENT', accent)).filter(Boolean).join('');
    if (!cards) return;
    html += `<div class="co-semester">
      <div class="co-semester-title">${esc(grp.name)}</div>
      <div class="co-grid">${cards}</div>
    </div>`;
  });

  if (ungrouped.length) {
    const accent = ACCENTS[(semesters.length + otherGroups.length) % ACCENTS.length];
    const cards = ungrouped.map(c => cardHTML(c).replace('ACCENT', accent)).filter(Boolean).join('');
    if (cards) {
      html += `<div class="co-semester">
        <div class="co-semester-title">All courses</div>
        <div class="co-grid">${cards}</div>
      </div>`;
    }
  }

  body.innerHTML = html || '<div style="color:var(--text3);padding:20px">No courses found.</div>';
}

// ═══════════════════════════════════════════════════════════════════════════
// Course selection
// ═══════════════════════════════════════════════════════════════════════════
async function selectCourse(path) {
  activeCourse = path;
  localStorage.setItem('last_course', path);
  _expandedFolders.clear();
  filterAndRenderSidebar();
  document.getElementById('tabs').style.display = 'flex';
  document.getElementById('course-title-bar').style.display = 'flex';
  document.getElementById('course-title-text').textContent = path.split('/').pop();
  // reset selection mode on course change
  selectionMode = false;
  const selBtn = document.getElementById('selection-toggle-btn');
  if (selBtn) { selBtn.style.color = 'var(--text3)'; }
  document.getElementById('file-list')?.classList.remove('selection-mode');
  switchTab('files');
  loadFiles();
  todoRenderCourse(path);
}

// ═══════════════════════════════════════════════════════════════════════════
// Tabs
// ═══════════════════════════════════════════════════════════════════════════
function switchTab(tab) {
  activeTab = tab;
  if (activeCourse) localStorage.setItem('last_tab', tab);
  document.querySelectorAll('.tab').forEach(t => t.classList.toggle('active', t.dataset.tab === tab));
  showPanel(tab);
  if (tab === 'files') _updateSummarizeBtn();

  if (tab === 'files')   loadFiles();
  if (tab === 'info')    loadInfo();
  if (tab === 'summary') loadSummary();
  if (tab === 'learn')   loadFlashcards();
  if (tab === 'notes')   loadNotes();
  if (tab === 'chat')    loadChat();
}

function showPanel(name) {
  document.querySelectorAll('.panel').forEach(p => p.classList.toggle('active', p.id === `panel-${name}`));
}

// ═══════════════════════════════════════════════════════════════════════════
// Files tab
// ═══════════════════════════════════════════════════════════════════════════
let fileSort = localStorage.getItem('file_sort') || 'name';
let activePreviewFile = null;
let _allLoadedFiles = [];
let _allLoadedMeta  = {};
let _oldFiles       = new Set(); // paths that are old course material (from Alle Kurse)

function filterFileList(q) {
  const clearBtn = document.getElementById('file-search-clear');
  if (clearBtn) clearBtn.style.display = q ? '' : 'none';
  _renderFileListWithFilter(q);
}

function clearFileSearch() {
  const inp = document.getElementById('file-search');
  if (inp) { inp.value = ''; inp.dispatchEvent(new Event('input')); inp.focus(); }
}

function _renderFileListWithFilter(q) {
  const el = document.getElementById('file-list');
  if (!el) return;
  const course = allCourses.find(c => c.path === activeCourse);
  const summaryAge = course?.sync_age || course?.summary_age || 0;
  const newCount = course?.new_files || 0;

  const banner = newCount > 0
    ? `<div style="margin-bottom:8px;padding:7px 10px;background:rgba(234,179,8,.1);border:1px solid rgba(234,179,8,.3);border-radius:6px;font-size:11px;color:var(--yellow)">
        ⚠ ${newCount} new file${newCount>1?'s':''} since last sync
       </div>`
    : '';

  const lower = (q || '').toLowerCase();
  if (lower) {
    // Flat filtered list — skip tree, show matches directly
    const matches = _allLoadedFiles.filter(f => f.toLowerCase().includes(lower));
    if (!matches.length) {
      el.innerHTML = banner + `<div style="color:var(--text3);font-size:12px;padding:8px">No files matching "${esc(q)}"</div>`;
      return;
    }
    el.innerHTML = banner + matches.map(f => {
      const m = _allLoadedMeta[f] || {};
      const isNew = summaryAge && m.mtime && m.mtime > summaryAge;
      const isOld = _oldFiles.has(f);
      const read = isFileRead(activeCourse, f);
      return `
        <div class="file-item${isNew ? ' new-file' : ''}${read ? ' file-read' : ''}" data-filename="${esc(f)}"
             style="padding-left:8px${isOld ? ';opacity:.7' : ''}" onclick="previewFileFromEl(this)">
          <input type="checkbox" name="file" value="${esc(f)}" checked onclick="event.stopPropagation()" onchange="if(selectionMode)updateSelectionToolbar()">
          <span class="file-icon">${fileIcon(f)}</span>
          <span class="file-name" title="${esc(f)}">${esc(f)}</span>
          ${isOld ? '<span class="new-badge" style="flex-shrink:0;background:var(--bg4);color:var(--text3)">old</span>' : ''}
          ${isNew ? '<span class="new-badge" style="flex-shrink:0">New</span>' : ''}
          <span class="file-read-check">${read ? '✓' : ''}</span>
        </div>`;
    }).join('');
  } else {
    const currentFiles = _allLoadedFiles.filter(f => !_oldFiles.has(f));
    let html = renderFileTree(buildFileTree(currentFiles), currentFiles, summaryAge, _allLoadedMeta, 0);

    // Old course material section — collapsed by default
    if (_oldFiles.size > 0) {
      const oldArr = [..._oldFiles];
      const isOpen = _expandedFolders.has('__old_material__');
      html += `
      <div class="folder-item old-material-header${isOpen ? '' : ' collapsed'}" style="padding-left:8px;margin-top:10px;border-top:1px solid var(--border);padding-top:8px"
           onclick="toggleFileFolder('__old_material__', this)">
        <span class="folder-chevron">▼</span>
        <span class="folder-icon">📦</span>
        <span class="folder-name" style="color:var(--text3)">Old Course Material</span>
        <span class="folder-count" style="color:var(--text3)">${oldArr.length}</span>
      </div>
      <div class="folder-contents${isOpen ? '' : ' collapsed'}" data-folder="__old_material__" style="opacity:.7">
        ${renderFileTree(buildFileTree(oldArr), oldArr, summaryAge, _allLoadedMeta, 0)}
      </div>`;
    }

    el.innerHTML = banner + html;
  }
}

function setFileSort(val) {
  fileSort = val;
  localStorage.setItem('file_sort', val);
  const sel = document.getElementById('file-sort-select');
  if (sel) sel.value = val;
  loadFiles();
}

async function loadFiles() {
  const [data, meta] = await Promise.all([
    fetch(`/api/files/${enc(activeCourse)}`).then(r => r.json()),
    fetch(`/api/file-meta/${enc(activeCourse)}`).then(r => r.json()).catch(() => []),
  ]);
  const el = document.getElementById('file-list');

  // Support both new {files, old_files} format and legacy plain array
  const files    = Array.isArray(data) ? data : (data.files    || []);
  const oldFiles = Array.isArray(data) ? []   : (data.old_files || []);

  if (!files.length && !oldFiles.length) {
    el.innerHTML = '<div style="color:var(--text3);font-size:12px;padding:8px">No files</div>';
    return;
  }
  const metaMap = Object.fromEntries((meta || []).map(m => [m.name, m]));

  // Apply sort
  const sel = document.getElementById('file-sort-select');
  if (sel) sel.value = fileSort;
  if (fileSort === 'date') {
    files.sort((a, b)    => (metaMap[b]?.mtime || 0) - (metaMap[a]?.mtime || 0));
    oldFiles.sort((a, b) => (metaMap[b]?.mtime || 0) - (metaMap[a]?.mtime || 0));
  }

  // Store for filtering (search covers both sets)
  _allLoadedFiles = [...files, ...oldFiles];
  _allLoadedMeta  = metaMap;
  _oldFiles       = new Set(oldFiles);

  // Clear search box on course change
  const searchInp = document.getElementById('file-search');
  if (searchInp) searchInp.value = '';
  const clearBtn = document.getElementById('file-search-clear');
  if (clearBtn) clearBtn.style.display = 'none';

  _renderFileListWithFilter('');

  // Re-open last viewed file
  const lastFile = localStorage.getItem('last_file__' + activeCourse);
  if (lastFile && _allLoadedFiles.includes(lastFile)) {
    previewFile(lastFile);
  }
}

// Build a nested tree from flat relative paths
function buildFileTree(files) {
  const root = { dirs: {}, files: [] };
  for (const f of files) {
    const parts = f.split('/');
    let node = root;
    for (let i = 0; i < parts.length - 1; i++) {
      const p = parts[i];
      if (!node.dirs[p]) node.dirs[p] = { dirs: {}, files: [], path: parts.slice(0, i + 1).join('/') };
      node = node.dirs[p];
    }
    node.files.push(f);
  }
  return root;
}

function countTreeFiles(node) {
  return node.files.length + Object.values(node.dirs).reduce((s, d) => s + countTreeFiles(d), 0);
}

// Track which folders are explicitly expanded (default: all collapsed)
const _expandedFolders = new Set();

function hasNewFilesInTree(node, metaMap, summaryAge) {
  if (!summaryAge) return false;
  for (const f of node.files) {
    const m = metaMap[f] || {};
    if (m.mtime && m.mtime > summaryAge) return true;
  }
  for (const child of Object.values(node.dirs)) {
    if (hasNewFilesInTree(child, metaMap, summaryAge)) return true;
  }
  return false;
}

function renderFileTree(node, allFiles, summaryAge, metaMap, depth) {
  let html = '';
  const indent = depth * 14;

  // Files at this level
  for (const f of node.files) {
    const m = metaMap[f] || {};
    const isNew = summaryAge && m.mtime && m.mtime > summaryAge;
    const displayName = f.split('/').pop();
    const dateStr = m.mtime ? formatFileDate(m.mtime) : '';
    const read = isFileRead(activeCourse, f);
    html += `
    <div class="file-item${isNew ? ' new-file' : ''}${read ? ' file-read' : ''}" data-filename="${esc(f)}"
         style="padding-left:${8 + indent}px" onclick="previewFileFromEl(this)">
      <input type="checkbox" name="file" value="${esc(f)}" checked onclick="event.stopPropagation()" onchange="if(selectionMode)updateSelectionToolbar()">
      <span class="file-icon">${fileIcon(f)}</span>
      <span class="file-name" title="${esc(f)}">${esc(displayName)}</span>
      ${isNew ? '<span class="new-badge" style="flex-shrink:0">New</span>' : ''}
      ${dateStr ? `<span class="file-date">${dateStr}</span>` : ''}
      <span class="file-read-check">${read ? '✓' : ''}</span>
    </div>`;
  }

  // Subfolders
  for (const [name, child] of Object.entries(node.dirs).sort(([a],[b]) => a.localeCompare(b))) {
    const folderPath = child.path;
    const collapsed = !_expandedFolders.has(folderPath);
    const count = countTreeFiles(child);
    const folderIsNew = hasNewFilesInTree(child, metaMap, summaryAge);
    html += `
    <div class="folder-item${collapsed ? ' collapsed' : ''}${folderIsNew ? ' has-new' : ''}" style="padding-left:${8 + indent}px"
         onclick="toggleFileFolder('${esc(folderPath)}', this)">
      <span class="folder-chevron">▼</span>
      <span class="folder-icon">📁</span>
      <span class="folder-name">${esc(name)}</span>
      ${folderIsNew ? '<span class="folder-new-badge">New</span>' : ''}
      <span class="folder-count">${count}</span>
    </div>
    <div class="folder-contents${collapsed ? ' collapsed' : ''}" data-folder="${esc(folderPath)}">
      ${renderFileTree(child, allFiles, summaryAge, metaMap, depth + 1)}
    </div>`;
  }
  return html;
}

function toggleFileFolder(path, headerEl) {
  const expanded = _expandedFolders.has(path);
  if (expanded) _expandedFolders.delete(path); else _expandedFolders.add(path);
  headerEl.classList.toggle('collapsed', expanded);
  const contents = document.querySelector(`.folder-contents[data-folder="${CSS.escape(path)}"]`);
  if (contents) contents.classList.toggle('collapsed', expanded);
}

function previewFileFromEl(el) { previewFile(el.dataset.filename); }

function formatFileDate(mtime) {
  const d = new Date(mtime * 1000);
  const now = new Date();
  const diffDays = Math.floor((now - d) / 86400000);
  if (diffDays === 0) return 'Heute';
  if (diffDays === 1) return 'Gestern';
  if (diffDays < 7) return d.toLocaleDateString('de-DE', { weekday: 'short' });
  return d.toLocaleDateString('de-DE', { day: '2-digit', month: '2-digit', year: '2-digit' });
}

function fileIcon(name) {
  const ext = name.split('.').pop().toLowerCase();
  return {pdf:'📕', doc:'📘', docx:'📘', ppt:'📙', pptx:'📙', txt:'📃', md:'📝',
          jpg:'🖼', jpeg:'🖼', png:'🖼', gif:'🖼', svg:'🖼',
          zip:'🗜', rar:'🗜', '7z':'🗜',
          mp4:'🎬', mov:'🎬', avi:'🎬', mp3:'🎵', wav:'🎵',
          py:'💻', js:'💻', ts:'💻', java:'💻', cpp:'💻', c:'💻'}[ext] || '📄';
}

const PREVIEWABLE_EXT  = new Set(['pdf','doc','docx','ppt','pptx','txt','md','py','js','ts','java','cpp','c','h','css','html','json','xml','csv']);
const IMAGE_EXT        = new Set(['jpg','jpeg','png','gif','svg','webp','bmp','ico']);
const AUDIO_EXT        = new Set(['mp3','wav','ogg','m4a','aac','flac']);
const VIDEO_EXT        = new Set(['mp4','webm','mov','avi','mkv']);

function _downloadPlaceholder(filename, rawUrl, errorMsg) {
  const name = filename.split('/').pop();
  return `<div class="preview-placeholder">
    <div class="icon">${fileIcon(filename)}</div>
    <div style="margin-bottom:6px;color:var(--text2);font-weight:500">${esc(name)}</div>
    ${errorMsg ? `<div style="font-size:11px;color:var(--red);margin-bottom:12px">${esc(errorMsg)}</div>`
               : `<div style="font-size:12px;color:var(--text3);margin-bottom:16px">No preview available</div>`}
    <a href="${rawUrl}" download
       style="background:var(--blue);color:#fff;padding:9px 22px;border-radius:8px;text-decoration:none;font-size:13px;font-weight:600">
      ⬇ Download
    </a>
  </div>`;
}

let _readTimer = null;
let _readTimerFile = null;

function _clearReadTimer() {
  if (_readTimer) { clearInterval(_readTimer); _readTimer = null; }
  _readTimerFile = null;
}

function _startReadTimer(filename) {
  _clearReadTimer();
  if (isFileRead(activeCourse, filename)) return;
  _readTimerFile = filename;
  const DURATION = 20; // seconds
  let elapsed = 0;
  _readTimer = setInterval(() => {
    elapsed++;
    if (elapsed >= DURATION) {
      _clearReadTimer();
      if (_readTimerFile === filename) { setFileRead(activeCourse, filename, true); filterAndRenderSidebar(); }
    }
  }, 1000);
}

function _saveScrollPos() {
  if (!activeCourse || !activePreviewFile) return;
  const body = document.getElementById('preview-body');
  if (body) localStorage.setItem('last_scroll__' + activeCourse + '__' + activePreviewFile, body.scrollTop);
}

async function previewFile(filename) {
  _clearReadTimer();
  // Save scroll position of the file we're leaving
  _saveScrollPos();
  // Tear down any active PDF so the ResizeObserver can't re-render it into a non-PDF view
  if (_pdfRO) { _pdfRO.disconnect(); _pdfRO = null; }
  if (_pdfDoc) { _pdfDoc.destroy(); _pdfDoc = null; }
  pdfFindClose();
  activePreviewFile = filename;
  localStorage.setItem('last_file__' + activeCourse, filename);
  _loadFileNote(filename); // load notes for this file if panel is open
  document.querySelectorAll('.file-item').forEach(el => {
    el.classList.toggle('active', el.dataset.filename === filename);
  });

  const ext = filename.split('.').pop().toLowerCase();
  const header = document.getElementById('preview-header');
  const body   = document.getElementById('preview-body');
  const read = isFileRead(activeCourse, filename);

  const isPdf = ext === 'pdf';
  document.getElementById('pdf-find-topbar-btn').style.display = isPdf ? '' : 'none';
  header.innerHTML = `
    <span class="preview-header-name">${esc(filename.split('/').pop())}</span>
    ${isPdf ? `
      <span id="pdf-page-ind" style="font-size:11px;color:var(--text3);flex-shrink:0"></span>
      <div class="zoom-controls">
        <button class="zoom-btn" onclick="changePdfZoom(-20)" title="Verkleinern">−</button>
        <span class="zoom-label" id="zoom-label">Auto</span>
        <button class="zoom-btn" onclick="changePdfZoom(20)" title="Zoom in">+</button>
      </div>` : ''}
    <button id="read-toggle-btn" class="read-toggle-btn${read ? ' is-read' : ''}"
      onclick="toggleReadFile('${esc(filename)}')">${read ? '✓ Read' : 'Mark as read'}</button>
    <a href="/api/file-raw/${enc(activeCourse)}/${enc(filename)}" download title="Download"
       style="color:var(--text3);font-size:13px;text-decoration:none" onclick="event.stopPropagation()">⬇</a>
    <button id="fnotes-toggle-btn" onclick="toggleFileNotes()" title="File notes (N)"
      style="background:none;border:none;cursor:pointer;font-size:13px;padding:2px 5px;border-radius:4px;color:var(--text3);transition:color var(--transition),background var(--transition)"
      onmouseover="this.style.background='var(--bg4)'" onmouseout="this.style.background='none'">✏️</button>
    ${isPdf ? `<button id="fullscreen-btn" onclick="togglePreviewFullscreen()" title="Vollbild (F)">⛶</button>` : ''}`;
  body.innerHTML = '<div class="preview-placeholder"><div class="icon">⏳</div><div>Loading…</div></div>';
  body.className = 'preview-body';

  const rawUrl = `/api/file-raw/${enc(activeCourse)}/${enc(filename)}`;

  if (isPdf) {
    _pdfScale = null;
    _pdfUrl = rawUrl;
    body.className = 'preview-body pdf-wrap';
    body.innerHTML = '';
    await _loadPdfJs(body);
    _setupPdfResizeObserver(body);
    _setupPdfWheelZoom(body);
    _restoreScrollPos(filename);
  } else if (IMAGE_EXT.has(ext)) {
    body.className = 'preview-body';
    body.innerHTML = `<div style="display:flex;flex-direction:column;align-items:center;gap:12px;padding:20px">
      <img src="${rawUrl}" alt="${esc(filename.split('/').pop())}"
           style="max-width:100%;max-height:70vh;border-radius:6px;box-shadow:var(--shadow)">
      <a href="${rawUrl}" download style="font-size:12px;color:var(--text3)">⬇ Download</a>
    </div>`;
    _restoreScrollPos(filename);
  } else if (AUDIO_EXT.has(ext)) {
    body.className = 'preview-body';
    body.innerHTML = `<div class="preview-placeholder">
      <div class="icon">🎵</div>
      <div style="margin-bottom:16px;color:var(--text2)">${esc(filename.split('/').pop())}</div>
      <audio controls style="width:100%;max-width:400px">
        <source src="${rawUrl}">
        Your browser does not support audio playback.
      </audio>
      <a href="${rawUrl}" download style="margin-top:12px;font-size:12px;color:var(--text3)">⬇ Download</a>
    </div>`;
  } else if (VIDEO_EXT.has(ext)) {
    body.className = 'preview-body';
    body.innerHTML = `<div style="display:flex;flex-direction:column;align-items:center;gap:12px;padding:20px">
      <div id="video-wrap" style="width:100%"></div>
      <a href="${rawUrl}" download style="font-size:12px;color:var(--text3)">⬇ Download</a>
    </div>`;
    // Create video element programmatically so the browser does not auto-mute it
    const vid = document.createElement('video');
    vid.controls = true;
    vid.muted = false;
    vid.style.cssText = 'max-width:100%;max-height:70vh;border-radius:6px;box-shadow:var(--shadow);display:block;margin:0 auto';
    vid.src = rawUrl;
    document.getElementById('video-wrap').appendChild(vid);
  } else if (PREVIEWABLE_EXT.has(ext)) {
    const data = await fetch(`/api/file-text/${enc(activeCourse)}/${enc(filename)}`).then(r => r.json());
    const text = data.text || '';
    if (text && !text.startsWith('Fehler')) {
      body.className = 'preview-body';
      body.innerHTML = `<div style="margin-bottom:12px;text-align:right">
        <a href="${rawUrl}" download style="font-size:11px;color:var(--text3);text-decoration:none">⬇ Download</a>
      </div>${esc(text)}`;
    } else {
      body.className = 'preview-body';
      body.innerHTML = _downloadPlaceholder(filename, rawUrl, text || '');
    }
    _restoreScrollPos(filename);
  } else {
    body.className = 'preview-body';
    body.innerHTML = _downloadPlaceholder(filename, rawUrl, '');
    _restoreScrollPos(filename);
  }
  _startReadTimer(filename);
}

function markAllFilesRead(read) {
  const items = document.querySelectorAll('.file-item[data-filename]');
  items.forEach(el => setFileRead(activeCourse, el.dataset.filename, read));
  if (read) _clearReadTimer();
  else if (activePreviewFile) _startReadTimer(activePreviewFile);
  filterAndRenderSidebar(); // update completion indicator
}

function downloadCourseZip(withNotes = false) {
  window.location.href = `/api/download-zip/${enc(activeCourse)}${withNotes ? '?notes=1' : ''}`;
}

function toggleReadFile(filename) {
  const nowRead = !isFileRead(activeCourse, filename);
  setFileRead(activeCourse, filename, nowRead);
  if (nowRead) _clearReadTimer();
  else _startReadTimer(filename);
  filterAndRenderSidebar(); // update completion indicator
}

function _restoreScrollPos(filename) {
  const key = 'last_scroll__' + activeCourse + '__' + filename;
  const saved = parseFloat(localStorage.getItem(key) || '0');
  if (!saved) return;
  requestAnimationFrame(() => {
    const body = document.getElementById('preview-body');
    if (body) body.scrollTop = saved;
  });
}

// Continuously save scroll position while user is scrolling
document.addEventListener('DOMContentLoaded', () => {
  const body = document.getElementById('preview-body');
  if (body) body.addEventListener('scroll', () => _saveScrollPos(), { passive: true });
});

// ── PDF.js state ─────────────────────────────────────────────────────────
let _pdfUrl      = '';
let _pdfScale    = null;  // null = auto-fit to container width
let _pdfDoc      = null;
let _pdfRO       = null, _pdfROTimer = null;
let _pdfGen      = 0;    // generation counter — incremented on every new render request

// Returns the scale that fills the container width exactly.
function _pdfAutoScale(container) {
  if (!_pdfDoc) return 1;
  // Use cached viewport from first page (already fetched as _pdfPage1Vp)
  const w = container.clientWidth - 16;
  return w > 0 ? w / _pdfPage1Width : 1;
}
let _pdfPage1Width = 0; // natural width of page 1 at scale=1

async function _loadPdfJs(container) {
  if (_pdfDoc) { _pdfDoc.destroy(); _pdfDoc = null; }
  try {
    _pdfDoc = await pdfjsLib.getDocument(_pdfUrl).promise;
  } catch(e) {
    container.innerHTML = `<div style="color:var(--red);padding:20px">Error: ${e.message}</div>`;
    return;
  }
  // Cache page-1 natural width so _pdfAutoScale() doesn't need to be async.
  const p1 = await _pdfDoc.getPage(1);
  _pdfPage1Width = p1.getViewport({ scale: 1 }).width;

  // Wait for the browser to finish laying out the container so clientWidth is real.
  await new Promise(r => requestAnimationFrame(r));
  await new Promise(r => requestAnimationFrame(r)); // two frames to be sure

  await _renderPdf(container);
}

async function _renderPdf(container) {
  if (!_pdfDoc) return;
  const gen = ++_pdfGen; // any concurrent render with a smaller gen is stale

  const containerWidth = container.clientWidth;
  if (!containerWidth) return;
  const scale = _pdfScale !== null ? _pdfScale : _pdfAutoScale(container);

  // Collect all page objects first (fast, no rendering yet)
  const pages = [];
  for (let i = 1; i <= _pdfDoc.numPages; i++) {
    pages.push(await _pdfDoc.getPage(i));
  }
  if (gen !== _pdfGen) return; // superseded

  // Render all pages into an off-screen DocumentFragment first, then
  // atomically swap — canvas rendering works without being in the DOM,
  // so the old content stays visible until everything is ready.
  const fragment = document.createDocumentFragment();
  const dpr = window.devicePixelRatio || 1;

  for (const page of pages) {
    if (gen !== _pdfGen) return; // superseded mid-render
    const viewport = page.getViewport({ scale: scale * dpr }); // render at physical pixels
    const cssW = (viewport.width  / dpr) + 'px';
    const cssH = (viewport.height / dpr) + 'px';

    // Wrapper — relative so text layer can be absolutely positioned on top
    const wrapper = document.createElement('div');
    wrapper.className = 'pdf-page-wrapper';
    wrapper.style.width  = cssW;
    wrapper.style.height = cssH;

    const canvas     = document.createElement('canvas');
    canvas.className = 'pdf-page-canvas';
    canvas.width     = viewport.width;  // physical pixel size
    canvas.height    = viewport.height;
    canvas.style.width  = cssW;
    canvas.style.height = cssH;
    wrapper.appendChild(canvas);

    // Text layer for selection
    const textLayer = document.createElement('div');
    textLayer.className = 'textLayer';
    textLayer.style.width  = cssW;
    textLayer.style.height = cssH;
    wrapper.appendChild(textLayer);

    fragment.appendChild(wrapper);

    // Render canvas while it's still off-screen (works fine without being in the DOM)
    await page.render({ canvasContext: canvas.getContext('2d'), viewport }).promise;
    if (gen !== _pdfGen) return;

    // Text layer at CSS scale so positions match CSS pixels
    const cssViewport = page.getViewport({ scale });
    try {
      const textContent = await page.getTextContent();
      if (gen !== _pdfGen) return;
      if (pdfjsLib.renderTextLayer) {
        pdfjsLib.renderTextLayer({
          textContentSource: textContent,
          container: textLayer,
          viewport: cssViewport,
          textDivs: [],
        });
      }
    } catch(_) { /* text layer is best-effort */ }
  }

  if (gen !== _pdfGen) return;

  // Atomic swap — old pages stay visible until all new pages are fully painted
  container.innerHTML = '';
  container.appendChild(fragment);

  _updateZoomLabel(container, scale);
  _setupPdfPageIndicator(container);
  // Re-apply any active find highlights after re-render
  if (_pdfFindTerm) requestAnimationFrame(() => pdfFindRun(_pdfFindTerm));
}

let _pdfPageObserver = null;
let _pdfPageRatios   = new Map();
function _setupPdfPageIndicator(container) {
  if (_pdfPageObserver) _pdfPageObserver.disconnect();
  _pdfPageRatios.clear();
  const ind = document.getElementById('pdf-page-ind');
  if (!ind || !_pdfDoc) return;
  const canvases = [...container.querySelectorAll('.pdf-page-canvas')];
  ind.textContent = `1 / ${canvases.length}`;
  _pdfPageObserver = new IntersectionObserver(entries => {
    // Update the persistent ratio map for every changed entry
    entries.forEach(e => _pdfPageRatios.set(e.target, e.intersectionRatio));
    // Pick the canvas with the highest visible ratio across ALL observed canvases
    let bestCanvas = null, bestRatio = -1;
    _pdfPageRatios.forEach((ratio, canvas) => {
      if (ratio > bestRatio) { bestRatio = ratio; bestCanvas = canvas; }
    });
    if (bestCanvas) ind.textContent = `${canvases.indexOf(bestCanvas) + 1} / ${canvases.length}`;
  }, { root: container, threshold: Array.from({length: 11}, (_, i) => i * 0.1) });
  canvases.forEach(c => { _pdfPageRatios.set(c, 0); _pdfPageObserver.observe(c); });
}

function _updateZoomLabel(container, scale) {
  const label = document.getElementById('zoom-label');
  if (!label) return;
  if (_pdfScale === null) {
    label.textContent = 'Auto';
  } else {
    const autoScale = _pdfAutoScale(container);
    label.textContent = Math.round((scale / autoScale) * 100) + '%';
  }
}

function applyPdfZoom() {
  const body = document.getElementById('preview-body');
  if (!body || !_pdfDoc) return;
  _renderPdf(body);
}

function changePdfZoom(delta) {
  const body = document.getElementById('preview-body');
  if (!body || !_pdfDoc) return;
  // Resolve current scale (auto → actual number) before applying delta
  const current = _pdfScale !== null ? _pdfScale : _pdfAutoScale(body);
  _pdfScale = Math.min(5, Math.max(0.1, current * (1 + delta / 100)));
  _renderPdf(body);
}

function _setupPdfResizeObserver(body) {
  if (_pdfRO) _pdfRO.disconnect();
  _pdfRO = new ResizeObserver(() => {
    if (_pdfScale !== null) return; // manual zoom active — don't clobber it
    clearTimeout(_pdfROTimer);
    _pdfROTimer = setTimeout(() => {
      const b = document.getElementById('preview-body');
      if (b && _pdfDoc) _renderPdf(b);
    }, 150);
  });
  _pdfRO.observe(body);
}

function _setupPdfWheelZoom(body) {
  body.addEventListener('wheel', e => {
    if (!_pdfDoc || !e.ctrlKey) return; // only intercept pinch-to-zoom (ctrlKey = true on macOS trackpad pinch)
    e.preventDefault();
    changePdfZoom(e.deltaY < 0 ? 15 : -15);
  }, { passive: false });
}

// ── PDF in-page search ────────────────────────────────────────────────────────
let _pdfFindMatches = [];
let _pdfFindIdx     = -1;
let _pdfFindTerm    = '';

function togglePdfFind() {
  const bar = document.getElementById('pdf-find-bar');
  if (!bar) return;
  if (bar.style.display === 'flex') { pdfFindClose(); }
  else { bar.style.display = 'flex'; document.getElementById('pdf-find-input').focus(); }
}

function pdfFindClose() {
  const bar = document.getElementById('pdf-find-bar');
  if (bar) bar.style.display = 'none';
  _pdfFindClear();
  const inp = document.getElementById('pdf-find-input');
  if (inp) inp.value = '';
  _pdfFindUpdateCount();
}

function _pdfFindClear() {
  document.querySelectorAll('#preview-body .pdf-hl').forEach(el => {
    el.classList.remove('pdf-hl', 'pdf-hl-cur');
  });
  _pdfFindMatches = []; _pdfFindIdx = -1; _pdfFindTerm = '';
}

function pdfFindRun(term) {
  document.querySelectorAll('#preview-body .pdf-hl').forEach(el => {
    el.classList.remove('pdf-hl', 'pdf-hl-cur');
  });
  _pdfFindMatches = [];
  _pdfFindTerm = term.toLowerCase();
  if (!term) { _pdfFindUpdateCount(); return; }

  document.querySelectorAll('#preview-body .textLayer span').forEach(span => {
    if (span.textContent.toLowerCase().includes(_pdfFindTerm)) {
      span.classList.add('pdf-hl');
      _pdfFindMatches.push(span);
    }
  });

  _pdfFindIdx = _pdfFindMatches.length > 0 ? 0 : -1;
  _pdfFindGoto(_pdfFindIdx);
}

function pdfFindNext() {
  if (!_pdfFindMatches.length) return;
  _pdfFindIdx = (_pdfFindIdx + 1) % _pdfFindMatches.length;
  _pdfFindGoto(_pdfFindIdx);
}

function pdfFindPrev() {
  if (!_pdfFindMatches.length) return;
  _pdfFindIdx = (_pdfFindIdx - 1 + _pdfFindMatches.length) % _pdfFindMatches.length;
  _pdfFindGoto(_pdfFindIdx);
}

function _pdfFindGoto(idx) {
  document.querySelectorAll('#preview-body .pdf-hl-cur').forEach(el => el.classList.remove('pdf-hl-cur'));
  if (idx >= 0 && _pdfFindMatches[idx]) {
    _pdfFindMatches[idx].classList.add('pdf-hl-cur');
    _pdfFindMatches[idx].scrollIntoView({ block: 'center', behavior: 'smooth' });
  }
  _pdfFindUpdateCount();
}

function _pdfFindUpdateCount() {
  const el = document.getElementById('pdf-find-count');
  if (!el) return;
  if (!_pdfFindTerm) { el.textContent = ''; return; }
  if (_pdfFindMatches.length === 0) {
    el.textContent = 'No matches'; el.style.color = 'var(--red)';
  } else {
    el.textContent = `${_pdfFindIdx + 1} / ${_pdfFindMatches.length}`; el.style.color = 'var(--text3)';
  }
}

let _isFullscreen = false;

function togglePreviewFullscreen() {
  const area = document.getElementById('preview-area');
  _isFullscreen = !_isFullscreen;
  area.classList.toggle('fullscreen', _isFullscreen);

  const btn = document.getElementById('fullscreen-btn');
  if (btn) btn.textContent = _isFullscreen ? '✕' : '⛶';

  // Re-render PDF to fit the new size
  setTimeout(() => { _pdfScale = null; applyPdfZoom(); }, 50);
}

// Esc closes fake fullscreen
document.addEventListener('keydown', e => {
  if (e.key === 'Escape' && _isFullscreen) {
    e.stopImmediatePropagation();
    togglePreviewFullscreen();
  }
}, true);

function startSummaryFlow() {
  if (!selectionMode) toggleSelectionMode();
  toast('Select files, then click Summarize.', '');
}

function toggleSelectionMode() {
  selectionMode = !selectionMode;
  const list    = document.getElementById('file-list');
  const btn     = document.getElementById('selection-toggle-btn');
  const toolbar = document.getElementById('selection-toolbar');
  const bulkRow = document.getElementById('bulk-read-row');
  list.classList.toggle('selection-mode', selectionMode);
  if (selectionMode) {
    btn.style.color = 'var(--blue)';
    btn.style.borderColor = 'var(--blue)';
    btn.textContent = 'Done';
    allFilesChecked = false;
    document.querySelectorAll('input[name="file"]').forEach(cb => cb.checked = false);
    toolbar.classList.add('visible');
    if (bulkRow) bulkRow.style.display = 'none';
    updateSelectionToolbar();
    setSummaryLength(localStorage.getItem('summary_length') || 'short');
  } else {
    btn.style.color = 'var(--text3)';
    btn.style.borderColor = 'var(--border)';
    btn.textContent = 'Select';
    document.querySelectorAll('input[name="file"]').forEach(cb => cb.checked = false);
    toolbar.classList.remove('visible');
    if (bulkRow) bulkRow.style.display = 'flex';
  }
}

function toggleAllFiles() {
  allFilesChecked = !allFilesChecked;
  document.querySelectorAll('input[name="file"]').forEach(cb => cb.checked = allFilesChecked);
}

function selectAllFiles() {
  allFilesChecked = true;
  document.querySelectorAll('input[name="file"]').forEach(cb => cb.checked = true);
  updateSelectionToolbar();
}

function selectNoFiles() {
  allFilesChecked = false;
  document.querySelectorAll('input[name="file"]').forEach(cb => cb.checked = false);
  updateSelectionToolbar();
}

function getSelectedFiles() {
  if (!selectionMode) return [...document.querySelectorAll('input[name="file"]')].map(cb => cb.value);
  return [...document.querySelectorAll('input[name="file"]:checked')].map(cb => cb.value);
}

// ═══════════════════════════════════════════════════════════════════════════
// Info tab
// ═══════════════════════════════════════════════════════════════════════════
let _customInfoFields = []; // [{label, value}]

async function loadInfo() {
  const el = document.getElementById('info-body');
  el.innerHTML = '<div style="color:var(--text3);font-size:13px;padding:16px">Loading…</div>';

  const course = allCourses.find(c => c.path === activeCourse);
  let info = {};
  try { info = await fetch(`/api/course-info/${enc(activeCourse)}`).then(r => r.json()); } catch(_) {}
  try { _customInfoFields = await fetch(`/api/custom-info/${enc(activeCourse)}`).then(r => r.json()); } catch(_) { _customInfoFields = []; }

  _renderInfo(course, info);
}

function _renderInfo(course, info) {
  const el = document.getElementById('info-body');
  const row = (label, value) => value
    ? `<div class="info-label">${esc(label)}</div><div class="info-value">${esc(String(value))}</div>`
    : '';
  const lecturers = Array.isArray(info?.lecturers) ? info.lecturers.join(', ') : (info?.lecturers || '');
  const hasScraped = info && Object.values(info).some(v => v && (Array.isArray(v) ? v.length : true));

  const scrapedCard = hasScraped ? `
    <div class="info-card">
      <div class="info-title">${esc(info.title || course?.name || activeCourse.split('/').pop())}</div>
      ${info.subtitle ? `<div class="info-subtitle">${esc(info.subtitle)}</div>` : ''}
      <div class="info-grid">
        ${row('Lecturer', lecturers)}
        ${row('Type', info.type)}
        ${row('Semester', info.semester)}
        ${row('ECTS', info.ects)}
        ${row('SWS', info.sws)}
        ${row('Participants', info.participants)}
        ${row('Location', info.location)}
        ${row('Course no.', info.course_no)}
        ${row('Institution', info.institution)}
        ${row('Language', info.language)}
      </div>
      ${info.description ? `<div class="info-desc">${esc(info.description)}</div>` : ''}
    </div>` : '';

  el.innerHTML = scrapedCard + `
    <div class="info-card" id="custom-info-card">
      <div style="font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:.06em;color:var(--text3);margin-bottom:12px">Course details</div>
      <div id="custom-info-rows"></div>
      <button onclick="customInfoAddRow()" style="margin-top:8px;background:none;border:1px dashed var(--border2);border-radius:6px;padding:4px 12px;font-size:12px;color:var(--text3);cursor:pointer;transition:border-color .15s,color .15s" onmouseover="this.style.color='var(--text)'" onmouseout="this.style.color='var(--text3)'">+ Add field</button>
    </div>`;

  _renderCustomRows();
}

function _renderCustomRows() {
  const wrap = document.getElementById('custom-info-rows');
  if (!wrap) return;
  wrap.innerHTML = _customInfoFields.map((f, i) => `
    <div class="custom-info-row">
      <input type="text" value="${esc(f.label)}" placeholder="Category"
        oninput="_customInfoEdit(${i},'label',this.value)" onblur="customInfoSave()">
      <textarea placeholder="Content" rows="1"
        oninput="_customInfoEdit(${i},'value',this.value);this.style.height='auto';this.style.height=this.scrollHeight+'px';customInfoSave()"
      >${esc(f.value)}</textarea>
      <button class="custom-info-del" onclick="customInfoDelRow(${i})" title="Remove">✕</button>
    </div>`).join('');
  // Auto-size textareas
  wrap.querySelectorAll('textarea').forEach(ta => {
    ta.style.height = 'auto';
    ta.style.height = ta.scrollHeight + 'px';
  });
}

function _customInfoEdit(i, key, val) {
  if (_customInfoFields[i]) _customInfoFields[i][key] = val;
}

function customInfoAddRow() {
  _customInfoFields.push({ label: '', value: '' });
  _renderCustomRows();
  // Focus the new label input
  const rows = document.querySelectorAll('.custom-info-row input');
  if (rows.length) rows[rows.length - 1].focus();
}

function customInfoDelRow(i) {
  _customInfoFields.splice(i, 1);
  _renderCustomRows();
  customInfoSave();
}

let _customInfoSaveTimer = null;
function customInfoSave() {
  clearTimeout(_customInfoSaveTimer);
  _customInfoSaveTimer = setTimeout(async () => {
    await fetch(`/api/custom-info/${enc(activeCourse)}`, {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify(_customInfoFields)
    });
  }, 600);
}

// ═══════════════════════════════════════════════════════════════════════════
// Summary tab
// ═══════════════════════════════════════════════════════════════════════════
let _summaryEditMode = false;
let _summaryActiveFile = null;

function summaryLabel(name) {
  if (name === '_zusammenfassung.md') return 'Summary 1';
  const m = name.match(/^_zusammenfassung_(\d+)\.md$/);
  if (m) return `Summary ${m[1]}`;
  // fallback for renamed/custom filenames
  return name.replace('_zusammenfassung_', '').replace('.md', '');
}

async function loadSummary(filename = null) {
  _summaryEditMode = false;
  const [summaries, data] = await Promise.all([
    fetch(`/api/summaries/${enc(activeCourse)}`).then(r => r.json()),
    fetch(`/api/summary/${enc(activeCourse)}${filename ? '?file='+enc(filename) : ''}`).then(r => r.json()),
  ]);
  _summaryActiveFile = data.file || null;
  const el = document.getElementById('summary-body');

  if (!data.html) {
    el.innerHTML = `
      <div class="empty-state">
        <div class="icon">📄</div>
        <h3>No summary yet</h3>
        <p style="max-width:340px;margin:0 auto 20px;line-height:1.7;color:var(--text3)">
          To create a summary, go to <strong style="color:var(--text2)">Files</strong>, click
          <strong style="color:var(--text2)">Select</strong> to enter selection mode,
          tick the files you want included, then hit
          <strong style="color:var(--text2)">Create Summary</strong>.
        </p>
        <button class="tbtn btn-blue" onclick="switchTab('files')">Go to Files →</button>
      </div>`;
    return;
  }

  const pillsHtml = summaries.map(s => {
    const label  = summaryLabel(s.name);
    const active = s.name === _summaryActiveFile;
    return `<span class="summary-file-pill ${active?'active':''}" onclick="loadSummary('${esc(s.name)}')">
      ${esc(label)}
      <button class="summary-pill-rename" onclick="event.stopPropagation();_summaryRename('${esc(s.name)}','${esc(label)}')" title="Rename">✎</button>
      <button class="summary-pill-del"    onclick="event.stopPropagation();_summaryDelete('${esc(s.name)}')" title="Delete">✕</button>
    </span>`;
  }).join('');

  el.innerHTML = `
    <div class="summary-layout" id="summary-layout">
      <div class="summary-content-wrap">
        <div class="summary-toolbar">
          <div style="flex:1;min-width:0;display:flex;gap:6px;flex-wrap:wrap">${pillsHtml}</div>
          ${summaries.length >= 2 ? `<button class="tbtn btn-gray" onclick="_openDiffModal()" title="Compare summaries">⇄ Compare</button>` : ''}
          <button class="tbtn btn-gray" id="summary-edit-btn" onclick="_summaryToggleEdit()" title="Edit markdown">✎ Edit</button>
          <button class="tbtn btn-gray" onclick="copyToClipboard(summaryMD)" title="Copy markdown">📋 Copy</button>
          <a class="tbtn btn-gray" href="/api/summary-raw/${enc(activeCourse)}?file=${enc(_summaryActiveFile)}" download style="text-decoration:none">⬇ Download</a>
        </div>
        <div id="summary-view-wrap">
          <div class="md-content" id="summary-md-content">${data.html}</div>
        </div>
        <div id="summary-edit-wrap" style="display:none;flex-direction:column;gap:8px">
          <textarea id="summary-editor">${esc(data.md)}</textarea>
          <div style="display:flex;gap:8px">
            <button class="tbtn btn-blue" onclick="_summarySave()">Save</button>
            <button class="tbtn btn-gray" onclick="_summaryToggleEdit()">Cancel</button>
          </div>
        </div>
      </div>
    </div>`;
  window.summaryMD = data.md;
  _buildAndInjectToC();
  renderLatexIn('summary-md-content');
}

function _summaryToggleEdit() {
  _summaryEditMode = !_summaryEditMode;
  document.getElementById('summary-view-wrap').style.display = _summaryEditMode ? 'none' : '';
  document.getElementById('summary-edit-wrap').style.display = _summaryEditMode ? 'flex' : 'none';
  document.getElementById('summary-edit-btn').textContent = _summaryEditMode ? '👁 View' : '✎ Edit';
}

async function _summarySave() {
  const content = document.getElementById('summary-editor').value;
  await fetch(`/api/summary-save/${enc(activeCourse)}`, {
    method: 'POST', headers: {'Content-Type':'application/json'},
    body: JSON.stringify({ file: _summaryActiveFile, content }),
  });
  toast('Saved', 'ok');
  loadSummary(_summaryActiveFile);
}

function _summaryRename(filename, currentLabel) {
  const newLabel = prompt('Rename summary:', currentLabel);
  if (!newLabel || newLabel === currentLabel) return;
  fetch(`/api/summary-rename/${enc(activeCourse)}`, {
    method: 'POST', headers: {'Content-Type':'application/json'},
    body: JSON.stringify({ file: filename, label: newLabel }),
  }).then(r => r.json()).then(d => {
    if (d.ok) { toast('Renamed', 'ok'); loadSummary(d.file); }
    else toast(d.error || 'Rename failed', 'err');
  });
}

function _summaryDelete(filename) {
  showConfirm('Delete summary?', `"${filename}" will be permanently deleted.`, async () => {
    await fetch(`/api/summary-delete/${enc(activeCourse)}`, {
      method: 'POST', headers: {'Content-Type':'application/json'},
      body: JSON.stringify({ file: filename }),
    });
    toast('Deleted', 'ok');
    loadSummary();
  });
}

let _diffSummaries = [];
function _openDiffModal() {
  fetch(`/api/summaries/${enc(activeCourse)}`).then(r => r.json()).then(summaries => {
    _diffSummaries = summaries || [];
    const selA = document.getElementById('diff-sel-a');
    const selB = document.getElementById('diff-sel-b');
    selA.innerHTML = _diffSummaries.map((s,i) => `<option value="${esc(s.name)}" ${i===0?'selected':''}>${esc(summaryLabel(s.name))}</option>`).join('');
    selB.innerHTML = _diffSummaries.map((s,i) => `<option value="${esc(s.name)}" ${i===1?'selected':''}>${esc(summaryLabel(s.name))}</option>`).join('');
    document.getElementById('diff-view').innerHTML = '<div style="padding:20px;text-align:center;color:var(--text3)">Select summaries and click Compare.</div>';
    document.getElementById('diff-overlay').classList.add('open');
  });
}
function hideDiffModal() { document.getElementById('diff-overlay').classList.remove('open'); }

function _computeDiff(a, b) {
  const linesA = a.split('\\n');
  const linesB = b.split('\\n');
  const m = linesA.length, n = linesB.length;
  // LCS table
  const dp = Array.from({length: m+1}, () => new Array(n+1).fill(0));
  for (let i = m-1; i >= 0; i--)
    for (let j = n-1; j >= 0; j--)
      dp[i][j] = linesA[i] === linesB[j] ? dp[i+1][j+1]+1 : Math.max(dp[i+1][j], dp[i][j+1]);
  const result = [];
  let i = 0, j = 0;
  while (i < m || j < n) {
    if (i < m && j < n && linesA[i] === linesB[j]) { result.push({type:'same', line: linesA[i]}); i++; j++; }
    else if (j < n && (i >= m || dp[i][j+1] >= dp[i+1][j])) { result.push({type:'add', line: linesB[j]}); j++; }
    else { result.push({type:'remove', line: linesA[i]}); i++; }
  }
  return result;
}

async function _runDiff() {
  const nameA = document.getElementById('diff-sel-a').value;
  const nameB = document.getElementById('diff-sel-b').value;
  const view = document.getElementById('diff-view');
  view.innerHTML = '<div style="padding:20px;text-align:center;color:var(--text3)">Loading…</div>';
  try {
    const [resA, resB] = await Promise.all([
      fetch(`/api/summary/${enc(activeCourse)}?file=${enc(nameA)}`).then(r => r.json()),
      fetch(`/api/summary/${enc(activeCourse)}?file=${enc(nameB)}`).then(r => r.json()),
    ]);
    const diff = _computeDiff(resA.md || '', resB.md || '');
    view.innerHTML = diff.map(d => {
      const cls = d.type === 'add' ? 'diff-line-add' : d.type === 'remove' ? 'diff-line-remove' : 'diff-line-same';
      const prefix = d.type === 'add' ? '+ ' : d.type === 'remove' ? '- ' : '  ';
      return `<div class="diff-line ${cls}">${prefix}${esc(d.line)}</div>`;
    }).join('') || '<div style="padding:20px;text-align:center;color:var(--text3)">No differences found.</div>';
  } catch(e) {
    view.innerHTML = `<div style="padding:20px;text-align:center;color:var(--red)">Error: ${esc(e.message)}</div>`;
  }
}

function _updateSummarizeBtn() { /* no-op: always create new numbered summary */ }

function setSummaryLength(len) {
  localStorage.setItem('summary_length', len);
  document.getElementById('sum-len-short')?.classList.toggle('btn-blue', len === 'short');
  document.getElementById('sum-len-short')?.classList.toggle('btn-gray',  len !== 'short');
  document.getElementById('sum-len-long')?.classList.toggle('btn-blue',  len === 'long');
  document.getElementById('sum-len-long')?.classList.toggle('btn-gray',   len !== 'long');
}

async function generateSummary() {
  const files  = getSelectedFiles();
  const limit  = parseInt(document.getElementById('limit-input')?.value) || 50;
  const lang   = localStorage.getItem('summary_lang') || 'en';
  const length = localStorage.getItem('summary_length') || 'short';
  setLoading(true);
  logShow(`Generating summary for "${activeCourse}"…\n`);

  const res  = await fetch('/api/summarize', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({ course: activeCourse, limit, force: true, files, lang, new_file: true, length })
  });
  const data = await res.json();
  logAppend(data.log || '');
  setLoading(false);

  if (data.success) {
    logAppend('\n✅ Done!');
    courseTree = await _fetchCourses();
    allCourses = flattenTree(courseTree);
    filterAndRenderSidebar();
    _updateSummarizeBtn();
    switchTab('summary');
    toast('Summary created!', 'ok');
  } else {
    // Show specific, actionable error messages
    const log = (data.log || '').toLowerCase();
    if (log.includes('anthropic_api_key missing') || log.includes('api_key missing')) {
      logAppend(_setupErrorMsg('no_key'));
    } else if (log.includes('ungültig') || log.includes('authenticationerror') || log.includes('invalid') && log.includes('key')) {
      logAppend(_setupErrorMsg('bad_key'));
    } else if (log.includes('connection') || log.includes('verbindung')) {
      logAppend(_setupErrorMsg('no_connection'));
    }
    toast('Error creating summary.', 'err');
  }
}

function _setupErrorMsg(type) {
  const msgs = {
    no_key: `

╔══════════════════════════════════════════════════════════════╗
  ❌  No API key configured
══════════════════════════════════════════════════════════════════
  An API key is required for AI summaries.
  Add one of the following keys to your .env file:

  Option A – Anthropic (recommended, free trial credit):
    ANTHROPIC_API_KEY=sk-ant-...
    → https://console.anthropic.com  →  "API Keys"

  Option B – OpenAI (GPT):
    OPENAI_API_KEY=sk-...
    → https://platform.openai.com/api-keys

  Option C – Other (Groq, Mistral, Ollama …):
    OPENAI_API_KEY=your-key
    OPENAI_BASE_URL=https://api.groq.com/openai/v1
    OPENAI_MODEL=llama-3.3-70b-versatile
╚══════════════════════════════════════════════════════════════╝`,
    bad_key: `

╔══════════════════════════════════════════════════════════════╗
  ❌  API key invalid or expired
══════════════════════════════════════════════════════════════════
  The ANTHROPIC_API_KEY in your .env was rejected.

  Check the following:
  1. Open .env and verify the key was copied correctly
     (no spaces, complete, starts with sk-ant-)
  2. Check https://console.anthropic.com that the key is
     still active and has remaining credit
╚══════════════════════════════════════════════════════════════╝`,
    no_connection: `

╔══════════════════════════════════════════════════════════════╗
  ❌  No connection to Anthropic API
══════════════════════════════════════════════════════════════════
  Check your internet connection and try again.
╚══════════════════════════════════════════════════════════════╝`,
  };
  return msgs[type] || '';
}


async function copyToClipboard(text) {
  try {
    await navigator.clipboard.writeText(text);
    toast('Copied to clipboard!', 'ok');
  } catch {
    toast('Copy not possible.', 'err');
  }
}

// ═══════════════════════════════════════════════════════════════════════════
// SRS Study System
// ═══════════════════════════════════════════════════════════════════════════
let _srsAllCards = [], _srsData = {}, _srsSession = null;

function _srsToday() { return new Date().toISOString().split('T')[0]; }

function _srsCardState(entry) {
  if (!entry || !entry.reps) return 'new';
  const today = _srsToday();
  if (entry.state === 'mastered') return 'mastered';
  if (entry.due < today)  return 'overdue';
  if (entry.due === today) return 'due';
  if (entry.interval <= 7) return 'learning';
  return 'review';
}

function srsSchedule(entry, rating) {
  let { interval = 1, ease = 2.5, reps = 0, lapses = 0 } = entry || {};
  let state;
  if (rating === 1) {
    lapses++; ease = Math.max(1.3, ease - 0.2); interval = 1; reps = 0; state = 'learning';
  } else {
    reps++;
    if      (reps === 1) interval = rating === 4 ? 4 : 1;
    else if (reps === 2) interval = rating === 4 ? 7 : 3;
    else if (rating === 2) { interval = Math.max(1, Math.round(interval * 1.2)); ease = Math.max(1.3, ease - 0.15); }
    else if (rating === 3) interval = Math.round(interval * ease);
    else                   { interval = Math.round(interval * ease * 1.3); ease = Math.min(3.0, ease + 0.15); }
    state = interval >= 21 ? 'mastered' : 'review';
  }
  const due = new Date(); due.setDate(due.getDate() + interval);
  return { interval, ease, reps, lapses, state, due: due.toISOString().split('T')[0] };
}

async function loadFlashcards() {
  const [summaryCards, customCards, srsRaw] = await Promise.all([
    fetch(`/api/flashcards/${enc(activeCourse)}`).then(r => r.json()),
    fetch(`/api/custom-cards/${enc(activeCourse)}`).then(r => r.json()),
    fetch(`/api/srs/${enc(activeCourse)}`).then(r => r.json()),
  ]);
  const normCustom = customCards.map((c, i) => ({
    id: `custom_${i}_${c.q.slice(0,16).replace(/\s/g,'_')}`,
    section: 'Custom', q: c.q, a: c.a,
  }));
  _srsAllCards = [...summaryCards, ...normCustom];
  _srsData = srsRaw || {};

  if (!_srsAllCards.length) {
    document.getElementById('learn-body').innerHTML = `
      <div class="empty-state">
        <div class="icon">🧠</div>
        <h3>No flashcards yet</h3>
        <p>Add your own cards or generate them from the AI summary.</p>
        <button class="tbtn btn-blue" onclick="showManageCards()">+ Add cards</button>
      </div>`;
    return;
  }
  renderStudyOverview();
}

function renderStudyOverview() {
  let nDue = 0, nNew = 0, nLearning = 0, nMastered = 0;
  for (const c of _srsAllCards) {
    const s = _srsCardState(_srsData[c.id]);
    if (s === 'new')                    nNew++;
    else if (s === 'due' || s === 'overdue') nDue++;
    else if (s === 'mastered')          nMastered++;
    else                                nLearning++;
  }
  const studyCount = nDue + Math.min(nNew, 20);
  document.getElementById('learn-body').innerHTML = `
    <div class="srs-overview">
      <div class="srs-overview-header">
        <h2>Study</h2>
        <button class="tbtn btn-gray" style="font-size:11px" onclick="showManageCards()">✎ Cards</button>
      </div>
      <div class="srs-stats">
        <div class="srs-stat srs-stat-due">
          <div class="srs-stat-num">${nDue}</div>
          <div class="srs-stat-label">Due</div>
        </div>
        <div class="srs-stat srs-stat-new">
          <div class="srs-stat-num">${nNew}</div>
          <div class="srs-stat-label">New</div>
        </div>
        <div class="srs-stat srs-stat-learning">
          <div class="srs-stat-num">${nLearning}</div>
          <div class="srs-stat-label">Learning</div>
        </div>
        <div class="srs-stat srs-stat-mastered">
          <div class="srs-stat-num">${nMastered}</div>
          <div class="srs-stat-label">Mastered</div>
        </div>
      </div>
      <div class="srs-action-btns">
        ${studyCount > 0
          ? `<button class="tbtn btn-blue" style="flex:1" onclick="_srsStart('due')">▶ Study due + new (${studyCount})</button>`
          : `<button class="tbtn btn-gray" style="flex:1;opacity:.6" disabled>✓ All caught up for today!</button>`}
        <button class="tbtn btn-gray" onclick="_srsStart('all')">All (${_srsAllCards.length})</button>
      </div>
      ${nMastered > 0 ? `<p style="font-size:12px;color:var(--text3);margin:0 0 12px">🏆 ${nMastered} card${nMastered>1?'s':''} mastered — well done!</p>` : ''}
    </div>`;
}

function _srsStart(mode) {
  let queue;
  if (mode === 'due') {
    const due = _srsAllCards.filter(c => { const s = _srsCardState(_srsData[c.id]); return s === 'due' || s === 'overdue'; });
    const nw  = _srsAllCards.filter(c => _srsCardState(_srsData[c.id]) === 'new').slice(0, 20);
    queue = shuffleArr([...due, ...nw]);
  } else {
    queue = shuffleArr([..._srsAllCards]);
  }
  if (!queue.length) { renderStudyOverview(); return; }
  _srsSession = {
    queue, current: 0, revealed: false,
    ratings: {1:0,2:0,3:0,4:0},
    hardCards: [],
    startTime: Date.now(),
    timerInterval: setInterval(() => {
      const el = document.getElementById('srs-timer');
      if (el) el.textContent = formatTime(Date.now() - _srsSession.startTime);
    }, 1000),
  };
  _srsRenderCard();
}

const _TYPE_LABELS = { recall:'Definition', mechanism:'Mechanism', contrast:'Contrast', application:'Application', cloze:'Fill in the blank' };

function _cardTypeBadge(type) {
  const label = _TYPE_LABELS[type] || 'Definition';
  return `<span class="card-type-badge ct-${type||'recall'}">${label}</span>`;
}

function _clozeQ(q) {
  return q.split('___').map(mdToHtml).join('<span class="cloze-blank"></span>');
}

function _clozeFilled(q, a) {
  return q.split('___').map(mdToHtml).join(`<span class="cloze-fill">${mdToHtml(a)}</span>`);
}

function _srsRenderCard() {
  const { queue, current } = _srsSession;
  if (current >= queue.length) { _srsDone(); return; }
  const card  = queue[current];
  const entry = _srsData[card.id];
  const state = _srsCardState(entry);
  const pct   = Math.round(current / queue.length * 100);
  const badgeClass = { new:'srs-badge-new', due:'srs-badge-due', overdue:'srs-badge-overdue', learning:'srs-badge-learning', review:'srs-badge-review', mastered:'srs-badge-mastered' }[state] || 'srs-badge-new';
  const intervalHint = entry?.interval ? (entry.interval === 1 ? 'interval: 1d' : `interval: ${entry.interval}d`) : 'first time';
  const isCloze = card.type === 'cloze';
  const questionHtml = isCloze ? _clozeQ(card.q) : mdToHtml(card.q);

  document.getElementById('learn-body').innerHTML = `
    <div class="flash-layout">
      <div class="flash-header">
        <span class="flash-header-title">${esc((activeCourse||'').split('/').pop())}</span>
        <span id="srs-timer" class="flash-timer">0:00</span>
        <button class="tbtn btn-gray" style="font-size:11px;padding:4px 10px" onclick="_srsAbort()">✕ Exit</button>
      </div>
      <div class="flash-progress-bar"><div class="flash-progress-fill" style="width:${pct}%"></div></div>
      <div class="flash-meta">
        <span>${current+1} / ${queue.length}</span>
        <span class="srs-badge ${badgeClass}">${state}</span>
        <span style="font-size:11px;color:var(--text3)">${intervalHint}</span>
      </div>
      <div class="flash-card" id="srs-card-el">
        <div style="display:flex;align-items:center;gap:8px;margin-bottom:6px">
          <div class="flash-section" style="margin:0">${esc(card.section||'')}</div>
          ${_cardTypeBadge(card.type)}
        </div>
        <div class="flash-question" id="srs-question">${questionHtml}</div>
        <div class="flash-answer" id="srs-answer" style="display:none">${isCloze ? `<span style="font-size:12px;color:var(--text3)">Answer: </span><strong>${mdToHtml(card.a)}</strong>` : mdToHtml(card.a)}</div>
      </div>
      <div id="srs-btns">
        <div style="text-align:center">
          <button class="flash-btn fb-reveal" onclick="_srsReveal()">Show answer</button>
        </div>
        <div class="flash-kbd-hint"><kbd>Space</kbd> Show answer</div>
      </div>
    </div>`;

  const t = document.getElementById('srs-timer');
  if (t) t.textContent = formatTime(Date.now() - _srsSession.startTime);
  _srsSession.revealed = false;
  renderLatexIn('srs-question');
}

function _srsReveal() {
  _srsSession.revealed = true;
  const card = _srsSession.queue[_srsSession.current];
  const isCloze = card.type === 'cloze';
  if (isCloze) {
    document.getElementById('srs-question').innerHTML = _clozeFilled(card.q, card.a);
    renderLatexIn('srs-question');
  }
  document.getElementById('srs-answer').style.display = 'block';
  document.getElementById('srs-btns').innerHTML = `
    <div class="rating-grid">
      <button class="rb rb-again" onclick="_srsRate(1)"><span class="rb-label">Again</span><span class="rb-hint">1</span></button>
      <button class="rb rb-hard"  onclick="_srsRate(2)"><span class="rb-label">Hard</span><span class="rb-hint">2</span></button>
      <button class="rb rb-good"  onclick="_srsRate(3)"><span class="rb-label">Good</span><span class="rb-hint">3</span></button>
      <button class="rb rb-easy"  onclick="_srsRate(4)"><span class="rb-label">Easy</span><span class="rb-hint">4</span></button>
    </div>
    <div class="flash-kbd-hint"><kbd>1</kbd> Again &nbsp; <kbd>2</kbd> Hard &nbsp; <kbd>3</kbd> Good &nbsp; <kbd>4</kbd> Easy</div>`;
  renderLatexIn('srs-answer');
}

function _srsRate(rating) {
  const card = _srsSession.queue[_srsSession.current];
  _srsData[card.id] = srsSchedule(_srsData[card.id], rating);
  _srsSession.ratings[rating]++;
  if (rating === 1) _srsSession.hardCards.push(card.q);
  _srsSession.current++;
  _srsSession.revealed = false;
  const el = document.getElementById('srs-card-el');
  const cls = (rating <= 2) ? 'flash-card-unknown' : 'flash-card-known';
  if (el) { el.classList.add(cls); setTimeout(_srsRenderCard, 120); } else _srsRenderCard();
  _srsSave();
}

function _srsSave() {
  fetch(`/api/srs/${enc(activeCourse)}`, {
    method: 'POST', headers: {'Content-Type':'application/json'},
    body: JSON.stringify(_srsData),
  });
}

function _srsDone() {
  clearInterval(_srsSession.timerInterval);
  const elapsed = formatTime(Date.now() - _srsSession.startTime);
  const { ratings, queue, hardCards } = _srsSession;
  const tomorrow = new Date(); tomorrow.setDate(tomorrow.getDate() + 1);
  const tStr = tomorrow.toISOString().split('T')[0];
  const dueTomorrow = _srsAllCards.filter(c => _srsData[c.id]?.due === tStr).length;
  const rColor = {1:'var(--red)',2:'var(--orange)',3:'var(--blue)',4:'var(--green)'};
  const rLabel = {1:'Again',2:'Hard',3:'Good',4:'Easy'};

  // Streak logic
  const streakKey = `srs_streak_${activeCourse}`;
  const today = new Date().toISOString().split('T')[0];
  const yesterday = new Date(Date.now() - 86400000).toISOString().split('T')[0];
  let streakData = { date: null, count: 0 };
  try { streakData = JSON.parse(localStorage.getItem(streakKey) || '{}'); } catch(e) {}
  if (streakData.date === today) { /* keep */ }
  else if (streakData.date === yesterday) streakData = { date: today, count: (streakData.count || 0) + 1 };
  else streakData = { date: today, count: 1 };
  localStorage.setItem(streakKey, JSON.stringify(streakData));
  const streakCount = streakData.count || 1;
  const streakHtml = streakCount > 1
    ? `<div class="srs-streak">🔥 ${streakCount} day streak</div>`
    : `<div class="srs-streak srs-streak-day1">Day 1</div>`;

  // Hard cards
  const uniqueHard = [...new Set(hardCards)];
  const hardHtml = uniqueHard.length ? `
    <details class="srs-hard-cards">
      <summary>⚠ ${uniqueHard.length} card${uniqueHard.length>1?'s':''} to review</summary>
      <ul>${uniqueHard.map(q => `<li>${esc(q.length > 80 ? q.slice(0,80)+'…' : q)}</li>`).join('')}</ul>
    </details>` : '';

  document.getElementById('learn-body').innerHTML = `
    <div class="srs-summary">
      <div style="font-size:52px;margin-bottom:12px">🎉</div>
      <h2>Session complete!</h2>
      ${streakHtml}
      <p class="srs-summary-sub">${queue.length} card${queue.length>1?'s':''} reviewed &middot; ⏱ ${elapsed}</p>
      <div class="srs-rating-breakdown">
        ${[1,2,3,4].map(r=>`
          <div class="srb-item">
            <div class="srb-num" style="color:${rColor[r]}">${ratings[r]}</div>
            <div class="srb-label">${rLabel[r]}</div>
          </div>`).join('')}
      </div>
      ${hardHtml}
      ${dueTomorrow > 0 ? `
      <div class="srs-next-due">
        <div class="srs-next-due-title">Up next</div>
        ${dueTomorrow} card${dueTomorrow>1?'s':''} due tomorrow
      </div>` : ''}
      <div style="display:flex;gap:10px;justify-content:center;flex-wrap:wrap">
        <button class="tbtn btn-blue" onclick="renderStudyOverview()">← Overview</button>
        <button class="tbtn btn-gray" onclick="_srsStart('due')">Study again</button>
      </div>
    </div>`;
}

function _srsAbort() {
  clearInterval(_srsSession?.timerInterval);
  _srsSession = null;
  renderStudyOverview();
}

function resetProgress() {
  showConfirm('Reset SRS progress?',
    `All study data for "${activeCourse}" will be deleted.`,
    async () => {
      await fetch(`/api/srs/${enc(activeCourse)}`, { method: 'POST', headers: {'Content-Type':'application/json'}, body: '{}' });
      _srsData = {};
      toast('Progress reset', 'ok');
      renderStudyOverview();
    }
  );
}

// ── Manage cards ────────────────────────────────────────────────────────────
let _mcSummaries = [], _mcSelectedSummary = null;

async function showManageCards() {
  const [summaryCards, customCards, summaries] = await Promise.all([
    fetch(`/api/flashcards/${enc(activeCourse)}`).then(r => r.json()),
    fetch(`/api/custom-cards/${enc(activeCourse)}`).then(r => r.json()),
    fetch(`/api/summaries/${enc(activeCourse)}`).then(r => r.json()),
  ]);
  _mcCustomCards = [...customCards];
  _mcSummaries   = summaries || [];
  if (!_mcSelectedSummary && _mcSummaries.length) _mcSelectedSummary = _mcSummaries[0].name;
  _renderManageCards(summaryCards);
}

let _mcCustomCards = [];

function _renderManageCards(summaryCards) {
  const courseInfo = allCourses.find(c => c.path === activeCourse);
  const hasSummary = courseInfo?.has_summary || false;
  const customList = _mcCustomCards.map((c, i) => `
    <div class="mc-card" id="mc-card-${i}">
      <div class="mc-card-body">
        ${_cardTypeBadge(c.type)}
        <div class="mc-card-q">${c.type === 'cloze' ? _clozeQ(c.q) : mdToHtml(c.q)}</div>
        <div class="mc-card-a">${mdToHtml(c.a)}</div>
      </div>
      <div style="display:flex;flex-direction:column;gap:4px;flex-shrink:0">
        <button class="mc-card-del" onclick="_mcEdit(${i})" title="Edit">✎</button>
        <button class="mc-card-del" onclick="_mcDelete(${i})" title="Delete">✕</button>
      </div>
    </div>`).join('');

  const totalCards = summaryCards.length + _mcCustomCards.length;

  document.getElementById('learn-body').innerHTML = `
    <div class="manage-cards-wrap">
      <div class="manage-cards-header">
        <h2>Flashcards</h2>
        ${totalCards ? `<button class="tbtn btn-blue" onclick="renderStudyOverview()">▶ Study (${totalCards})</button>` : ''}
      </div>

      ${hasSummary ? `
      <div class="mc-gen-section">
        <div class="mc-gen-info">
          <div class="mc-gen-icon">✨</div>
          <div>
            <div class="mc-gen-title">Generate from summary</div>
            <div class="mc-gen-desc">
              <select id="mc-gen-file" style="background:var(--bg3);border:1px solid var(--border);border-radius:5px;color:var(--text);font-size:11px;padding:3px 6px;margin-top:4px;max-width:220px"
                onchange="_mcSelectedSummary=this.value">
                ${_mcSummaries.map(s => `<option value="${esc(s.name)}" ${s.name===_mcSelectedSummary?'selected':''}>${esc(summaryLabel(s.name))}</option>`).join('')}
              </select>
            </div>
          </div>
        </div>
        <div style="display:flex;align-items:center;gap:8px;flex-shrink:0">
          <input id="mc-gen-count" type="number" value="10" min="3" max="30"
            style="width:52px;background:var(--bg3);border:1px solid var(--border);border-radius:6px;color:var(--text);font-size:12px;padding:4px 6px;text-align:center"
            title="Number of cards">
          <button class="tbtn btn-blue" id="mc-gen-btn" onclick="_mcGenerate()">Generate</button>
        </div>
      </div>
` : `
      <div class="mc-gen-section mc-gen-section-disabled">
        <div class="mc-gen-icon">✨</div>
        <span>Create a summary first to generate cards with AI.</span>
      </div>`}

      <div class="mc-add-form">
        <textarea id="mc-q" placeholder="Question…" rows="2"></textarea>
        <textarea id="mc-a" placeholder="Answer…" rows="2"></textarea>
        <div class="mc-add-row">
          <button class="tbtn btn-blue" onclick="_mcAdd()">+ Add card</button>
        </div>
      </div>

      ${_mcCustomCards.length ? `
        <div class="mc-section-label">Your cards (${_mcCustomCards.length})</div>
        <div class="mc-list">${customList}</div>` : ''}

      ${hasSummary ? `
        <div class="mc-section-label">From summary (${summaryCards.length})</div>
        <div class="mc-list">${summaryCards.map(c => `
          <div class="mc-card">
            <div class="mc-card-body">
              ${_cardTypeBadge(c.type)}
              <div class="mc-card-q">${c.type === 'cloze' ? _clozeQ(c.q) : mdToHtml(c.q)}</div>
              <div class="mc-card-a">${mdToHtml(c.a)}</div>
            </div>
          </div>`).join('')}
        </div>` : ''}
    </div>`;
  renderLatexIn('learn-body');
}

function _mcAdd() {
  const q = document.getElementById('mc-q').value.trim();
  const a = document.getElementById('mc-a').value.trim();
  if (!q || !a) return;
  _mcCustomCards.push({ q, a });
  _mcSave();
  showManageCards();
}

function _mcDelete(i) {
  _mcCustomCards.splice(i, 1);
  _mcSave();
  showManageCards();
}

function _mcEdit(i) {
  const card = _mcCustomCards[i];
  const el = document.getElementById(`mc-card-${i}`);
  if (!el) return;
  el.innerHTML = `
    <div class="mc-card-edit" style="flex:1">
      <textarea id="mc-edit-q-${i}" rows="2" style="width:100%;resize:vertical;background:var(--bg3);border:1px solid var(--border);border-radius:var(--radius);color:var(--text);font-size:13px;padding:6px 8px;font-family:inherit;box-sizing:border-box;margin-bottom:6px">${esc(card.q)}</textarea>
      <textarea id="mc-edit-a-${i}" rows="2" style="width:100%;resize:vertical;background:var(--bg3);border:1px solid var(--border);border-radius:var(--radius);color:var(--text);font-size:13px;padding:6px 8px;font-family:inherit;box-sizing:border-box;margin-bottom:8px">${esc(card.a)}</textarea>
      <div style="display:flex;gap:8px">
        <button class="tbtn btn-blue" style="font-size:11px;padding:4px 12px" onclick="_mcEditSave(${i})">Save</button>
        <button class="tbtn btn-gray" style="font-size:11px;padding:4px 12px" onclick="showManageCards()">Cancel</button>
      </div>
    </div>`;
}

function _mcEditSave(i) {
  const q = document.getElementById(`mc-edit-q-${i}`)?.value.trim();
  const a = document.getElementById(`mc-edit-a-${i}`)?.value.trim();
  if (!q || !a) return;
  _mcCustomCards[i].q = q;
  _mcCustomCards[i].a = a;
  _mcSave();
  showManageCards();
}

function _mcSave() {
  fetch(`/api/custom-cards/${enc(activeCourse)}`, {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify(_mcCustomCards),
  });
}

async function _mcDeleteSummary(filename) {
  showConfirm('Delete summary?', `"${filename}" will be permanently deleted.`, async () => {
    await fetch(`/api/summary-delete/${enc(activeCourse)}`, {
      method: 'POST', headers: {'Content-Type':'application/json'},
      body: JSON.stringify({ file: filename }),
    });
    toast('Summary deleted', 'ok');
    if (_mcSelectedSummary === filename) _mcSelectedSummary = null;
    showManageCards();
  });
}

async function _mcGenerate() {
  const btn = document.getElementById('mc-gen-btn');
  if (!btn) return;
  btn.disabled = true;
  btn.textContent = '⏳ Generating…';
  const count = parseInt(document.getElementById('mc-gen-count')?.value) || 10;
  const file  = document.getElementById('mc-gen-file')?.value || _mcSelectedSummary || null;
  logShow(`Generating ${count} flashcards for "${activeCourse}"…\n`);
  try {
    const res  = await fetch(`/api/generate-cards/${enc(activeCourse)}`, {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ count, file }),
    });
    const data = await res.json();
    if (!res.ok) {
      logAppend(`\n${data.error || '❌ Generation failed — unknown error.'}\n`);
      btn.disabled = false;
      btn.textContent = 'Generate';
      return;
    }
    const generated = Array.isArray(data) ? data : [];
    const existing = new Set(_mcCustomCards.map(c => c.q.toLowerCase()));
    const added = generated.filter(c => !existing.has(c.q.toLowerCase()));
    _mcCustomCards.push(...added);
    _mcSave();
    logAppend(`✅ Done — ${added.length} card${added.length !== 1 ? 's' : ''} added.\n`);
    showManageCards();
  } catch(e) {
    logAppend(`\n❌ Network error: ${e.message}\n`);
    btn.disabled = false;
    btn.textContent = 'Generate';
  }
}

// ═══════════════════════════════════════════════════════════════════════════
// Flashcards — global (all courses)
// ═══════════════════════════════════════════════════════════════════════════

function shuffleArr(arr) {
  for (let i = arr.length - 1; i > 0; i--) {
    const j = Math.floor(Math.random() * (i + 1));
    [arr[i], arr[j]] = [arr[j], arr[i]];
  }
  return arr;
}

function formatTime(ms) {
  const s = Math.floor(ms / 1000);
  const m = Math.floor(s / 60);
  return `${m}:${String(s % 60).padStart(2, '0')}`;
}

// ═══════════════════════════════════════════════════════════════════════════
// Keyboard shortcuts
// ═══════════════════════════════════════════════════════════════════════════
document.addEventListener('keydown', e => {
  // Overlay close
  if (e.key === 'Escape') {
    if (document.getElementById('cmd-overlay').classList.contains('open')) { closeCmdPalette(); return; }
    if (document.getElementById('shortcuts-overlay').classList.contains('open')) { hideShortcuts(); return; }
    if (document.getElementById('confirm-overlay').classList.contains('open')) { hideConfirm(); return; }
    goHome();
    return;
  }

  if (e.key === '?' && e.target.tagName !== 'INPUT' && e.target.tagName !== 'TEXTAREA') {
    showShortcuts(); return;
  }

  if (e.key === 'b' && e.target.tagName !== 'INPUT' && e.target.tagName !== 'TEXTAREA') {
    toggleSidebar(); return;
  }

  if (e.key === 'f' && e.target.tagName !== 'INPUT' && e.target.tagName !== 'TEXTAREA') {
    if (document.getElementById('fullscreen-btn')) { togglePreviewFullscreen(); return; }
  }

  // Ctrl/Cmd+F → PDF find bar (when PDF is active)
  if ((e.ctrlKey || e.metaKey) && e.key === 'f' && _pdfDoc) {
    e.preventDefault();
    togglePdfFind();
    return;
  }

  // Ctrl/Cmd+K → command palette
  if ((e.ctrlKey || e.metaKey) && e.key === 'k') {
    e.preventDefault();
    document.getElementById('cmd-overlay').classList.contains('open') ? closeCmdPalette() : openCmdPalette();
    return;
  }

  // Tab shortcuts (1-4) when a course is active
  if (activeCourse && !e.ctrlKey && !e.metaKey && e.target.tagName !== 'INPUT' && e.target.tagName !== 'TEXTAREA') {
    if (e.key === '1') { switchTab('files'); return; }
    if (e.key === '2') { switchTab('summary'); return; }
    if (e.key === '3') { switchTab('learn'); return; }
    if (e.key === '4') { switchTab('notes'); return; }
    if (e.key === '5') { switchTab('chat'); return; }
    // R = toggle read on active preview file
    if (e.key === 'r' && activeTab === 'files' && activePreviewFile) {
      toggleReadFile(activePreviewFile); return;
    }
    // N = toggle file notes panel
    if (e.key === 'n' && activeTab === 'files') {
      toggleFileNotes(); return;
    }
  }

  // SRS Study shortcuts
  if (activeTab !== 'learn') return;
  if (e.target.tagName === 'INPUT' || e.target.tagName === 'TEXTAREA') return;
  if (!_srsSession) return;

  if (e.key === ' ' || e.code === 'Space') {
    e.preventDefault();
    if (!_srsSession.revealed && _srsSession.current < _srsSession.queue.length) _srsReveal();
  } else if ((e.key === '1' || e.key === '2' || e.key === '3' || e.key === '4') && _srsSession.revealed) {
    e.preventDefault();
    _srsRate(parseInt(e.key));
  }
});

// ═══════════════════════════════════════════════════════════════════════════
// Notes
// ═══════════════════════════════════════════════════════════════════════════
async function loadNotes() {
  const data = await fetch(`/api/notes/${enc(activeCourse)}`).then(r => r.json());
  const editor = document.getElementById('notes-editor');
  editor.value = data.text || '';
  document.getElementById('notes-saved').style.display = 'none';
  // Reset preview mode
  if (notesPreviewMode) toggleNotesPreview();
}

async function saveNotes() {
  const text = document.getElementById('notes-editor').value;
  await fetch(`/api/notes/${enc(activeCourse)}`, {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({ text })
  });
  const el = document.getElementById('notes-saved');
  el.textContent = '✓ Gespeichert';
  el.style.color = 'var(--green)';
  el.style.display = 'inline';
  setTimeout(() => { el.style.display = 'none'; }, 2000);
  allCourses = allCourses.map(c => c.path === activeCourse ? {...c, has_notes: true} : c);
}

// ═══════════════════════════════════════════════════════════════════════════
// File-level notes (sliding right panel)
// ═══════════════════════════════════════════════════════════════════════════
let _fileNotesOpen   = false;
let _fileNotesFile   = null;   // currently loaded filename
let _fileNotesSaveTimer = null;

function toggleFileList() {
  const col = document.getElementById('files-list-col');
  const tab = document.getElementById('filelist-toggle-tab');
  const collapsed = col.classList.toggle('collapsed');
  if (tab) tab.textContent = collapsed ? '▶' : '◀';
}

function toggleFileNotes() {
  _fileNotesOpen = !_fileNotesOpen;
  const col = document.getElementById('files-notes-col');
  if (_fileNotesOpen) {
    // Restore saved width (if user had resized it) before re-opening
    const saved = localStorage.getItem('notes_col_width');
    if (saved) col.style.width = saved + 'px';
  } else {
    // Clear inline width so the CSS 'width: 0' (no .open class) can take effect
    col.style.width = '';
  }
  col.classList.toggle('open', _fileNotesOpen);
  document.getElementById('divider-notes').style.display = _fileNotesOpen ? '' : 'none';
  if (_fileNotesOpen && activePreviewFile) {
    _loadFileNote(activePreviewFile);
  }
  // update toggle button color to indicate state
  const btn = document.getElementById('fnotes-toggle-btn');
  if (btn) btn.style.color = _fileNotesOpen ? 'var(--blue)' : 'var(--text3)';
}

const FNOTES_DEFAULT = '## Questions\n\n- [ ] \n\n## Notes\n\n';

// Called whenever a new file is selected in the preview
async function _loadFileNote(filename) {
  if (!_fileNotesOpen) return;
  _fileNotesFile = filename;
  const editor = document.getElementById('fnotes-editor');
  const title  = document.getElementById('fnotes-title');
  title.textContent = filename.split('/').pop();
  editor.value = '';
  document.getElementById('fnotes-saved').textContent = '';
  try {
    const data = await fetch(`/api/file-note/${enc(activeCourse)}/${enc(filename)}`).then(r => r.json());
    editor.value = data.text || FNOTES_DEFAULT;
  } catch(_) { editor.value = FNOTES_DEFAULT; }
  if (_fnotesPreviewMode) _renderFnotes();
}

let _fnotesPreviewMode = false;

function toggleFnotesMode() {
  _fnotesPreviewMode = !_fnotesPreviewMode;
  const editor   = document.getElementById('fnotes-editor');
  const rendered = document.getElementById('fnotes-rendered');
  const btn      = document.getElementById('fnotes-mode-btn');
  if (_fnotesPreviewMode) {
    _renderFnotes();
    editor.style.display   = 'none';
    rendered.style.display = 'block';
    btn.title = 'Bearbeiten';
    btn.textContent = '✏️';
  } else {
    editor.style.display   = 'block';
    rendered.style.display = 'none';
    btn.title = 'Vorschau';
    btn.textContent = '👁';
  }
}

async function _renderFnotes() {
  const md = document.getElementById('fnotes-editor').value;
  const rendered = document.getElementById('fnotes-rendered');
  try {
    const data = await fetch('/api/notes-preview', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({ text: md })
    }).then(r => r.json());
    rendered.innerHTML = data.html || '';
    // Enable and wire up task-list checkboxes (markdown2 renders them as disabled)
    rendered.querySelectorAll('input.task-list-item-checkbox').forEach((cb, i) => {
      cb.removeAttribute('disabled');
      cb.addEventListener('change', () => _toggleFnoteCheckbox(i, cb.checked));
    });
  } catch(_) {
    rendered.innerHTML = '<pre style="color:var(--text3);white-space:pre-wrap">' + md.replace(/&/g,'&amp;').replace(/</g,'&lt;') + '</pre>';
  }
}

function _toggleFnoteCheckbox(index, checked) {
  const editor = document.getElementById('fnotes-editor');
  let md = editor.value;
  let count = 0;
  md = md.replace(/^(\s*-\s*\[)([ x])(\])/gm, (match, pre, state, post) => {
    if (count === index) {
      count++;
      return pre + (checked ? 'x' : ' ') + post;
    }
    count++;
    return match;
  });
  editor.value = md;
  clearTimeout(_fileNotesSaveTimer);
  _fileNotesSaveTimer = setTimeout(_saveFileNote, 600);
}

async function _saveFileNote() {
  if (!_fileNotesFile) return;
  const text = document.getElementById('fnotes-editor').value;
  await fetch(`/api/file-note/${enc(activeCourse)}/${enc(_fileNotesFile)}`, {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({ text }),
  });
  const el = document.getElementById('fnotes-saved');
  el.textContent = '✓';
  el.style.color = 'var(--green)';
  setTimeout(() => { if (el) el.textContent = ''; }, 1800);
}

function downloadFileNote() {
  if (!_fileNotesFile) return;
  window.location.href = `/api/file-note-download/${enc(activeCourse)}/${enc(_fileNotesFile)}`;
}

// Auto-save on input (debounced)
document.addEventListener('DOMContentLoaded', () => {
  const ed = document.getElementById('fnotes-editor');
  ed.addEventListener('input', () => {
    clearTimeout(_fileNotesSaveTimer);
    _fileNotesSaveTimer = setTimeout(_saveFileNote, 600);
  });
  // Tab / Shift+Tab → indent / outdent by 2 spaces
  ed.addEventListener('keydown', e => {
    if (e.key === 'Tab') {
      e.preventDefault();
      const start = ed.selectionStart, end = ed.selectionEnd;
      if (e.shiftKey) {
        const lineStart = ed.value.lastIndexOf('\n', start - 1) + 1;
        if (ed.value.slice(lineStart, lineStart + 2) === '  ') {
          ed.value = ed.value.slice(0, lineStart) + ed.value.slice(lineStart + 2);
          ed.selectionStart = ed.selectionEnd = Math.max(lineStart, start - 2);
        }
      } else {
        ed.value = ed.value.slice(0, start) + '  ' + ed.value.slice(end);
        ed.selectionStart = ed.selectionEnd = start + 2;
      }
      clearTimeout(_fileNotesSaveTimer);
      _fileNotesSaveTimer = setTimeout(_saveFileNote, 600);
      return;
    }
  });
  // Enter on list lines → continue list; empty list line → break out
  ed.addEventListener('keydown', e => {
    if (e.key !== 'Enter') return;
    const start = ed.selectionStart;
    const lineStart = ed.value.lastIndexOf('\n', start - 1) + 1;
    const line = ed.value.slice(lineStart, start);
    let insert = null;
    const cbMatch  = line.match(/^(\s*)- \[[ x]\] (.*)/);
    const ulMatch  = line.match(/^(\s*)([-*+]) (.*)/);
    const olMatch  = line.match(/^(\s*)(\d+)\. (.*)/);
    if (cbMatch) {
      // checkbox: empty item → break out, else new checkbox
      insert = cbMatch[2].trim() === '' ? '\n' : '\n' + cbMatch[1] + '- [ ] ';
    } else if (ulMatch) {
      insert = ulMatch[3].trim() === '' ? '\n' : '\n' + ulMatch[1] + ulMatch[2] + ' ';
    } else if (olMatch) {
      insert = olMatch[3].trim() === '' ? '\n' : '\n' + olMatch[1] + (parseInt(olMatch[2]) + 1) + '. ';
    }
    if (insert !== null) {
      e.preventDefault();
      ed.value = ed.value.slice(0, start) + insert + ed.value.slice(ed.selectionEnd);
      const pos = start + insert.length;
      ed.selectionStart = ed.selectionEnd = pos;
      clearTimeout(_fileNotesSaveTimer);
      _fileNotesSaveTimer = setTimeout(_saveFileNote, 600);
    }
  });
});

function toggleNotesPreview() {
  notesPreviewMode = !notesPreviewMode;
  const editor  = document.getElementById('notes-editor');
  const preview = document.getElementById('notes-preview');
  const btn     = document.getElementById('notes-preview-btn');

  if (notesPreviewMode) {
    // Render markdown preview
    const md = editor.value;
    // Simple client-side markdown: use server for rendering
    fetch('/api/notes-preview', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({ text: md })
    }).then(r => r.json()).then(d => {
      preview.innerHTML = d.html || '';
    }).catch(() => {
      // Fallback: basic rendering
      preview.innerHTML = '<pre style="color:var(--text2);white-space:pre-wrap">' + esc(md) + '</pre>';
    });
    editor.style.display = 'none';
    preview.style.display = 'block';
    btn.textContent = 'Bearbeiten';
  } else {
    editor.style.display = 'block';
    preview.style.display = 'none';
    btn.textContent = 'Vorschau';
  }
}

// ═══════════════════════════════════════════════════════════════════════════
// Chat / Ask AI
// ═══════════════════════════════════════════════════════════════════════════
let chatHistory = [];   // [{role, content}]
let chatStreaming = false;
let chatContextFile = '';  // filename of currently selected file context ('' = summary)
let _chatHistoryData = []; // loaded saved conversations, indexed by position

const CHAT_SUGGESTIONS = [
  'What are the most important concepts?',
  'Explain the key term simply.',
  'What exam questions might come up?',
  'Summarize the most important points in 3 bullets.',
  'What connects to this topic?',
];

function _chatNew() {
  if (chatStreaming) return;
  chatHistory = [];
  chatContextFile = '';
  // Reset context selector
  const sel = document.getElementById('chat-context-select');
  if (sel) sel.value = '';
  document.getElementById('chat-file-pill-row').innerHTML = '';
  // Clear messages and show welcome bubble again
  const course = allCourses.find(c => c.path === activeCourse);
  const courseName = (activeCourse || '').split('/').pop();
  const msgs = document.getElementById('chat-messages');
  msgs.innerHTML = `<div class="chat-msg">
    <div class="chat-avatar">🤖</div>
    <div class="chat-bubble" id="chat-welcome-bubble">
      Hi! I'm using the <strong>summary</strong> of <strong>${esc(courseName)}</strong> as context. Select a specific file below to ask about it directly.
    </div>
  </div>`;
  // Re-add suggestions if they were removed
  if (course?.has_summary && !document.getElementById('chat-suggestions')) {
    const suggDiv = document.createElement('div');
    suggDiv.className = 'chat-suggestions';
    suggDiv.id = 'chat-suggestions';
    suggDiv.innerHTML = CHAT_SUGGESTIONS.map(s => `<button class="chat-suggestion" onclick="sendSuggestion('${esc(s)}')">${esc(s)}</button>`).join('');
    msgs.after(suggDiv);
  }
}

async function loadChat() {
  chatHistory = [];
  chatContextFile = '';
  const course = allCourses.find(c => c.path === activeCourse);
  const hasSummary = course?.has_summary;
  document.getElementById('chat-body').innerHTML = `
    <div class="chat-layout" style="position:relative">
      <div class="chat-toolbar">
        <button class="tbtn btn-gray" style="font-size:11px" onclick="_chatNew()">＋ New</button>
        <button class="tbtn btn-gray" style="font-size:11px" onclick="_chatSave()">💾 Save</button>
        <button class="tbtn btn-gray" style="font-size:11px" onclick="_chatDownload()">⬇ Download</button>
        <button class="tbtn btn-gray" style="font-size:11px" onclick="_chatHistoryToggle()">🕘 History</button>
      </div>
      <div class="chat-messages" id="chat-messages">
        <div class="chat-msg">
          <div class="chat-avatar">🤖</div>
          <div class="chat-bubble" id="chat-welcome-bubble">
            ${hasSummary
              ? `Hi! I'm using the <strong>summary</strong> of <strong>${esc(activeCourse.split('/').pop())}</strong> as context. Select a specific file below to ask about it directly.`
              : `No summary exists for this course yet. Create one under "Files" first so I can help you.`}
          </div>
        </div>
      </div>
      ${hasSummary ? `<div class="chat-suggestions" id="chat-suggestions">
        ${CHAT_SUGGESTIONS.map(s => `<button class="chat-suggestion" onclick="sendSuggestion('${esc(s)}')">${esc(s)}</button>`).join('')}
      </div>` : ''}
      <div class="chat-context-bar" id="chat-context-bar">
        <span class="chat-context-label">Context:</span>
        <select class="chat-context-select" id="chat-context-select" onchange="_chatContextChange(this.value)">
          <option value="">Summary (default)</option>
        </select>
      </div>
      <div id="chat-file-pill-row"></div>
      <div class="chat-input-row">
        <textarea id="chat-input" placeholder="Ask a question… (Enter to send, Shift+Enter for newline)"
          ${hasSummary ? '' : 'disabled'}
          onkeydown="handleChatKey(event)" oninput="this.style.height='auto';this.style.height=this.scrollHeight+'px'"></textarea>
        <button class="chat-input-btn" id="chat-send-btn" onclick="sendChat()" ${hasSummary ? '' : 'disabled'}>Send</button>
      </div>
      <div class="chat-history-panel" id="chat-history-panel">
        <div class="chat-history-header">
          <h3>Saved conversations</h3>
          <button class="tbtn btn-gray" style="font-size:11px;padding:3px 8px" onclick="_chatHistoryToggle()">✕</button>
        </div>
        <div class="chat-history-list" id="chat-history-list">Loading…</div>
      </div>
    </div>`;
  // Populate file dropdown
  try {
    const files = await fetch(`/api/file-meta/${enc(activeCourse)}`).then(r => r.json());
    const sel = document.getElementById('chat-context-select');
    if (sel && files.length) {
      for (const f of files) {
        const opt = document.createElement('option');
        opt.value = f.name;
        opt.textContent = f.name;
        sel.appendChild(opt);
      }
    }
  } catch {}
}

function _chatContextChange(value) {
  chatContextFile = value;
  const pillRow = document.getElementById('chat-file-pill-row');
  const welcome = document.getElementById('chat-welcome-bubble');
  const courseName = (activeCourse || '').split('/').pop();
  if (value) {
    pillRow.innerHTML = `<div class="chat-file-pill">
      <span>📄 ${esc(value)}</span>
      <button onclick="_chatContextChange(''); document.getElementById('chat-context-select').value=''" title="Clear">✕</button>
    </div>`;
    if (welcome) welcome.innerHTML = `Asking about <strong>${esc(value)}</strong>. I'll read the raw file content as context.`;
  } else {
    pillRow.innerHTML = '';
    if (welcome) welcome.innerHTML = `Hi! I'm using the <strong>summary</strong> of <strong>${esc(courseName)}</strong> as context. Select a specific file below to ask about it directly.`;
  }
}

function _chatMdExport() {
  if (!chatHistory.length) return null;
  const course = (activeCourse || '').split('/').pop();
  const date   = new Date().toLocaleString();
  const lines  = [`# Chat — ${course}\n\n*${date}*\n`];
  for (const m of chatHistory)
    lines.push(`**${m.role === 'user' ? 'You' : 'AI'}:** ${m.content}\n`);
  return lines.join('\n');
}

function _chatAutoTitle() {
  const first = chatHistory.find(m => m.role === 'user');
  if (!first) return 'Conversation';
  // Strip LaTeX, markdown bold/italic/code, and collapse whitespace
  let t = first.content
    .replace(/\$\$[\s\S]+?\$\$|\$[^$\n]+?\$/g, '')  // LaTeX
    .replace(/\*\*(.+?)\*\*|__(.+?)__/g, '$1$2')      // bold
    .replace(/\*(.+?)\*|_(.+?)_/g, '$1$2')            // italic
    .replace(/`[^`]+`/g, '')                           // inline code
    .replace(/\s+/g, ' ').trim();
  if (!t) t = 'Conversation';
  return t.slice(0, 60) + (t.length > 60 ? '…' : '');
}

async function _chatSave() {
  if (!chatHistory.length) { toast('Nothing to save yet', ''); return; }
  await fetch(`/api/chat-history/${enc(activeCourse)}`, {
    method: 'POST',
    headers: {'Content-Type':'application/json'},
    body: JSON.stringify({ title: _chatAutoTitle(), messages: chatHistory }),
  });
  toast('Conversation saved', 'ok');
}

function _chatDownload() {
  const md = _chatMdExport();
  if (!md) { toast('Nothing to download yet', ''); return; }
  const blob = new Blob([md], { type: 'text/markdown' });
  const a = document.createElement('a');
  a.href = URL.createObjectURL(blob);
  a.download = `chat_${(activeCourse||'').split('/').pop().replace(/\s+/g,'_')}_${Date.now()}.md`;
  a.click();
}

function _chatHistoryToggle() {
  const panel = document.getElementById('chat-history-panel');
  const isOpen = panel.classList.toggle('open');
  if (isOpen) _chatHistoryLoad();
}

async function _chatHistoryLoad() {
  const list = document.getElementById('chat-history-list');
  _chatHistoryData = await fetch(`/api/chat-history/${enc(activeCourse)}`).then(r => r.json());
  if (!_chatHistoryData.length) { list.innerHTML = '<p style="padding:12px;font-size:12px;color:var(--text3)">No saved conversations yet.</p>'; return; }
  list.innerHTML = _chatHistoryData.map((c, i) => `
    <div class="chat-history-item" data-idx="${i}">
      <span class="chat-history-item-label" onclick="_chatHistoryLoad_open(${i})">${esc(c.title)}</span>
      <button class="chat-history-rename-btn" onclick="_chatHistoryRename(${i}, event)" title="Rename">✏️</button>
      <button class="chat-history-rename-btn" onclick="_chatHistoryDelete(${i}, event)" title="Delete" style="color:var(--red)">🗑</button>
    </div>`).join('');
}

function _chatHistoryRename(idx, e) {
  e.stopPropagation();
  const conv = _chatHistoryData[idx];
  const item = document.querySelector(`.chat-history-item[data-idx="${idx}"]`);
  const label = item.querySelector('.chat-history-item-label');
  const btn   = item.querySelector('.chat-history-rename-btn');
  // Replace label with input
  const input = document.createElement('input');
  input.className = 'chat-history-rename-input';
  input.value = conv.title;
  label.replaceWith(input);
  btn.style.display = 'none';
  input.focus();
  input.select();
  const save = async () => {
    const newTitle = input.value.trim() || conv.title;
    conv.title = newTitle;
    await fetch(`/api/chat-history/${enc(activeCourse)}/rename`, {
      method: 'PATCH',
      headers: {'Content-Type':'application/json'},
      body: JSON.stringify({ id: conv.id, title: newTitle }),
    });
    // Restore label
    const newLabel = document.createElement('span');
    newLabel.className = 'chat-history-item-label';
    newLabel.textContent = newTitle;
    newLabel.onclick = () => _chatHistoryLoad_open(idx);
    input.replaceWith(newLabel);
    btn.style.display = '';
  };
  input.addEventListener('blur', save);
  input.addEventListener('keydown', e => {
    if (e.key === 'Enter') { e.preventDefault(); input.blur(); }
    if (e.key === 'Escape') { input.value = conv.title; input.blur(); }
  });
}

async function _chatHistoryDelete(idx, e) {
  e.stopPropagation();
  const conv = _chatHistoryData[idx];
  await fetch(`/api/chat-history/${enc(activeCourse)}/delete`, {
    method: 'DELETE',
    headers: {'Content-Type':'application/json'},
    body: JSON.stringify({ id: conv.id }),
  });
  _chatHistoryData.splice(idx, 1);
  // Re-render list
  _chatHistoryLoad();
}

function _chatHistoryLoad_open(idx) {
  const conv = _chatHistoryData[idx];
  chatHistory = [...conv.messages];
  const msgs = document.getElementById('chat-messages');
  msgs.innerHTML = '';
  for (const m of conv.messages) appendChatMsg(m.role, m.content);
  document.getElementById('chat-history-panel').classList.remove('open');
  msgs.scrollTop = msgs.scrollHeight;
  // Re-enable input so the user can continue the conversation
  const input = document.getElementById('chat-input');
  const sendBtn = document.getElementById('chat-send-btn');
  if (input) { input.disabled = false; input.focus(); }
  if (sendBtn) sendBtn.disabled = false;
}

function _chatCopyBubble(btn) {
  const bubble = btn.closest('.chat-msg').querySelector('.chat-bubble');
  navigator.clipboard.writeText(bubble.innerText).then(() => {
    btn.textContent = '✓';
    setTimeout(() => btn.textContent = '⎘', 1200);
  });
}

function sendSuggestion(text) {
  document.getElementById('chat-input').value = text;
  document.getElementById('chat-suggestions')?.remove();
  sendChat();
}

function handleChatKey(e) {
  if (e.key === 'Enter' && !e.shiftKey) { e.preventDefault(); sendChat(); }
}

async function sendChat() {
  if (chatStreaming) return;
  const input = document.getElementById('chat-input');
  const question = input.value.trim();
  if (!question) return;

  input.value = '';
  input.style.height = 'auto';
  document.getElementById('chat-suggestions')?.remove();

  // Add user message
  appendChatMsg('user', question);
  chatHistory.push({ role: 'user', content: question });

  // Add empty AI bubble with streaming cursor
  const aiId = 'ai-msg-' + Date.now();
  const messagesEl = document.getElementById('chat-messages');
  const bubble = document.createElement('div');
  bubble.className = 'chat-msg';
  bubble.innerHTML = `<div class="chat-avatar">🤖</div><div class="chat-bubble streaming" id="${aiId}"></div>`;
  messagesEl.appendChild(bubble);
  messagesEl.scrollTop = messagesEl.scrollHeight;

  chatStreaming = true;
  document.getElementById('chat-send-btn').disabled = true;
  input.disabled = true;

  let fullAnswer = '';
  try {
    const body = { question, history: chatHistory.slice(0, -1) };
    if (chatContextFile) body.file = chatContextFile;
    const resp = await fetch(`/api/chat/${enc(activeCourse)}`, {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify(body)
    });

    const reader = resp.body.getReader();
    const decoder = new TextDecoder();
    const bubbleEl = document.getElementById(aiId);

    while (true) {
      const { done, value } = await reader.read();
      if (done) break;
      const chunk = decoder.decode(value);
      for (const line of chunk.split('\n')) {
        if (!line.startsWith('data: ')) continue;
        const data = line.slice(6);
        if (data === '[DONE]') break;
        try {
          const { delta } = JSON.parse(data);
          fullAnswer += delta;
          bubbleEl.textContent = fullAnswer;
          messagesEl.scrollTop = messagesEl.scrollHeight;
        } catch {}
      }
    }

    const bubbleEl2 = document.getElementById(aiId);
    if (bubbleEl2) {
      bubbleEl2.classList.remove('streaming');
      bubbleEl2.style.position = 'relative';
      bubbleEl2.innerHTML = chatMdToHtml(fullAnswer);
      // Add copy button
      const copyBtn = document.createElement('button');
      copyBtn.textContent = '⎘';
      copyBtn.title = 'Copy';
      copyBtn.style.cssText = 'position:absolute;top:8px;right:8px;background:var(--bg3);border:1px solid var(--border);border-radius:5px;color:var(--text3);cursor:pointer;font-size:11px;padding:2px 7px;opacity:0;transition:opacity .15s';
      copyBtn.onclick = () => _chatCopyBubble(copyBtn);
      bubbleEl2.appendChild(copyBtn);
      bubbleEl2.closest('.chat-msg')?.addEventListener('mouseenter', () => copyBtn.style.opacity = '1');
      bubbleEl2.closest('.chat-msg')?.addEventListener('mouseleave', () => copyBtn.style.opacity = '0');
      renderLatexIn(bubbleEl2);
    }
    chatHistory.push({ role: 'assistant', content: fullAnswer });
  } catch (err) {
    const bubbleEl = document.getElementById(aiId);
    if (bubbleEl) { bubbleEl.classList.remove('streaming'); bubbleEl.textContent = 'Error: ' + err.message; }
  }

  chatStreaming = false;
  document.getElementById('chat-send-btn').disabled = false;
  input.disabled = false;
  input.focus();
}

function chatMdToHtml(text) {
  // Stash LaTeX
  const stash = [];
  text = text.replace(/\$\$[\s\S]+?\$\$|\$[^$\n]+?\$/g, m => { stash.push(m); return `\x00${stash.length-1}\x00`; });
  // HTML-escape
  text = text.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
  // Headings (must come before bold to avoid ## being eaten)
  text = text.replace(/^### (.+)$/gm, '<h4 style="font-size:13px;font-weight:700;color:var(--text);margin:14px 0 4px">$1</h4>');
  text = text.replace(/^## (.+)$/gm,  '<h3 style="font-size:15px;font-weight:700;color:#93c5fd;margin:16px 0 6px">$1</h3>');
  text = text.replace(/^# (.+)$/gm,   '<h2 style="font-size:17px;font-weight:800;color:var(--text);margin:18px 0 8px">$1</h2>');
  // Inline code
  text = text.replace(/`([^`]+)`/g, '<code style="background:var(--bg3);padding:1px 5px;border-radius:4px;font-size:12px;color:#6ee7b7">$1</code>');
  // Bold / italic
  text = text.replace(/\*\*(.+?)\*\*/g, '<strong>$1</strong>');
  text = text.replace(/\*([^*\n]+?)\*/g, '<em>$1</em>');
  // Bullet lists
  text = text.replace(/^[ \t]*[-*] (.+)$/gm, '<li>$1</li>');
  text = text.replace(/(<li>[\s\S]*?<\/li>)(?!\s*<li>)/g, '<ul style="margin:6px 0 10px 20px;line-height:1.7">$1</ul>');
  // Paragraphs
  text = text.split(/\n{2,}/).map(block => {
    block = block.trim();
    if (!block) return '';
    if (/^<(h[2-4]|ul|li)/.test(block)) return block;
    block = block.replace(/\n/g, '<br>');
    return `<p>${block}</p>`;
  }).filter(Boolean).join('');
  // Restore LaTeX
  text = text.replace(/\x00(\d+)\x00/g, (_, i) => stash[+i]);
  return text;
}

function appendChatMsg(role, text) {
  const el = document.getElementById('chat-messages');
  const div = document.createElement('div');
  div.className = `chat-msg ${role}`;
  const bubble = document.createElement('div');
  bubble.className = 'chat-bubble';
  bubble.style.position = 'relative';
  bubble.innerHTML = role === 'user' ? esc(text) : chatMdToHtml(text);
  if (role !== 'user') {
    const copyBtn = document.createElement('button');
    copyBtn.textContent = '⎘';
    copyBtn.title = 'Copy';
    copyBtn.style.cssText = 'position:absolute;top:6px;right:6px;background:none;border:none;color:var(--text3);cursor:pointer;font-size:13px;padding:2px 4px;opacity:0;transition:opacity .15s';
    copyBtn.onclick = () => _chatCopyBubble(copyBtn);
    bubble.appendChild(copyBtn);
    div.addEventListener('mouseenter', () => copyBtn.style.opacity = '1');
    div.addEventListener('mouseleave', () => copyBtn.style.opacity = '0');
  }
  div.innerHTML = `<div class="chat-avatar">${role === 'user' ? '🧑' : '🤖'}</div>`;
  div.appendChild(bubble);
  el.appendChild(div);
  if (role !== 'user') renderLatexIn(bubble);
  el.scrollTop = el.scrollHeight;
}

// Tab / Shift+Tab indent in course notes-editor
document.addEventListener('keydown', e => {
  if (e.target.id !== 'notes-editor' || e.key !== 'Tab') return;
  e.preventDefault();
  const ed = e.target, start = ed.selectionStart, end = ed.selectionEnd;
  if (e.shiftKey) {
    const lineStart = ed.value.lastIndexOf('\n', start - 1) + 1;
    if (ed.value.slice(lineStart, lineStart + 2) === '  ') {
      ed.value = ed.value.slice(0, lineStart) + ed.value.slice(lineStart + 2);
      ed.selectionStart = ed.selectionEnd = Math.max(lineStart, start - 2);
    }
  } else {
    ed.value = ed.value.slice(0, start) + '  ' + ed.value.slice(end);
    ed.selectionStart = ed.selectionEnd = start + 2;
  }
  clearTimeout(notesSaveTimer);
  notesSaveTimer = setTimeout(saveNotes, 3000);
});

// Auto-save notes + unsaved indicator
document.addEventListener('input', e => {
  if (e.target.id === 'notes-editor') {
    const el = document.getElementById('notes-saved');
    el.textContent = '● Unsaved';
    el.style.color = 'var(--orange)';
    el.style.display = 'inline';
    clearTimeout(notesSaveTimer);
    notesSaveTimer = setTimeout(saveNotes, 3000);
  }
});

// ═══════════════════════════════════════════════════════════════════════════
// Global search
// ═══════════════════════════════════════════════════════════════════════════
let searchTimer = null;
function handleGlobalSearch(e) {
  const q = e.target.value.trim();
  clearTimeout(searchTimer);
  if (!q) { if (activeCourse) switchTab(activeTab); else goHome(); return; }
  searchTimer = setTimeout(() => doSearch(q), 300);
}

function _renderSearchResults(results, q, el, showFileSearchBtn = true) {
  const escaped = q.replace(/[.*+?^${}()|[\]\\]/g, '\\$&');
  const SOURCE_LABEL = {
    summary: 'Summary', notes: 'Notes', both: 'Summary + Notes', file: 'File content',
  };
  const SOURCE_CLASS = { summary: 'summary', notes: 'notes', both: 'both', file: 'summary' };
  const rows = results.map(r => {
    const snippet = r.snippet.replace(new RegExp(escaped, 'gi'), m => `<mark>${m}</mark>`);
    const tab = r.source === 'notes' ? 'notes' : r.source === 'file' ? 'files' : 'summary';
    const fileTag = r.file ? `<span style="font-size:10px;color:var(--text3);margin-left:4px">${esc(r.file)}</span>` : '';
    return `
      <div class="search-result" data-course="${esc(r.course)}" onclick="selectCourseFromEl(this);switchTab('${tab}')">
        <div class="search-result-course">
          ${esc(r.name || r.course)}${fileTag}
          ${r.count > 1 ? `<span class="search-result-count">${r.count} matches</span>` : ''}
          ${r.source ? `<span class="search-source search-source-${SOURCE_CLASS[r.source]||'summary'}">${SOURCE_LABEL[r.source]||r.source}</span>` : ''}
        </div>
        <div class="search-result-snippet">…${snippet}…</div>
      </div>`;
  }).join('');
  const fileBtn = showFileSearchBtn
    ? `<button class="tbtn btn-gray" style="font-size:11px;padding:4px 10px" onclick="doFileSearch(${JSON.stringify(q)})">📄 Search in files</button>`
    : `<span style="font-size:11px;color:var(--text3)">incl. file contents</span>`;
  el.innerHTML = `
    <div class="search-results-header" style="display:flex;align-items:center;gap:10px">
      <span>${results.length} result${results.length !== 1 ? 's' : ''} for „${esc(q)}"</span>
      <div style="flex:1"></div>
      ${fileBtn}
    </div>${rows || `<div class="search-empty">No results in summaries/notes</div>`}`;
}

async function doSearch(q) {
  showPanel('search');
  document.getElementById('tabs').style.display = 'none';
  const el = document.getElementById('search-results-body');
  el.innerHTML = '<div class="search-empty">Suche…</div>';
  const results = await fetch(`/api/search?q=${encodeURIComponent(q)}`).then(r => r.json());
  _renderSearchResults(results, q, el);
}

async function doFileSearch(q) {
  const el = document.getElementById('search-results-body');
  el.innerHTML = '<div class="search-empty">Searching file contents (may take a moment)…</div>';
  const [sumResults, fileResults] = await Promise.all([
    fetch(`/api/search?q=${encodeURIComponent(q)}`).then(r => r.json()),
    fetch(`/api/search-files?q=${encodeURIComponent(q)}`).then(r => r.json()),
  ]);
  const seen = new Set(sumResults.map(r => r.course));
  const combined = [...sumResults, ...fileResults.filter(r => !seen.has(r.course))];
  _renderSearchResults(combined, q, el, false);
}

// ═══════════════════════════════════════════════════════════════════════════
// Scraper
// ═══════════════════════════════════════════════════════════════════════════
let _syncPolling = false;

async function _pollSyncStatus(onDone) {
  if (_syncPolling) return;
  _syncPolling = true;
  let lastLen = 0;
  while (_syncPolling) {
    await new Promise(r => setTimeout(r, 1000));
    let data;
    try { data = await (await fetch('/api/sync-status')).json(); }
    catch { break; }
    if (data.log && data.log.length > lastLen) {
      logAppend(data.log.slice(lastLen));
      lastLen = data.log.length;
    }
    if (!data.running) {
      _syncPolling = false;
      onDone(data);
      return;
    }
  }
  _syncPolling = false;
}

async function runScraper() {
  setLoading(true);
  logShow('Syncing new files from Stud.IP…\n');
  const data = await (await fetch('/api/scrape', { method: 'POST' })).json();
  if (!data.started) {
    logAppend(data.error || 'Could not start.');
    setLoading(false);
    return;
  }
  _pollSyncStatus(async result => {
    setLoading(false);
    if (result.success) {
      logAppend('\n✅ Done!');
      courseTree = await _fetchCourses();
      allCourses = flattenTree(courseTree);
      filterAndRenderSidebar();
      if (activeCourse) loadFiles();
      toast('Files updated!', 'ok');
    } else {
      toast('Scraping error.', 'err');
    }
  });
}

async function syncCourse() {
  if (!activeCourse) return;
  const btn = document.getElementById('course-sync-btn');
  if (btn) { btn.disabled = true; btn.textContent = '⏳ Sync…'; }
  setLoading(true);
  logShow(`Sync: ${activeCourse}…\n`);
  const data = await (await fetch('/api/sync-course', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({course: activeCourse}),
  })).json();
  if (!data.started) {
    logAppend(data.error || 'Could not start.');
    setLoading(false);
    if (btn) { btn.disabled = false; btn.textContent = '↓ Sync'; }
    toast(data.error?.includes('not in registry') ? 'Run a full sync first.' : 'Sync failed.', 'err');
    return;
  }
  _pollSyncStatus(async result => {
    setLoading(false);
    if (btn) { btn.disabled = false; btn.textContent = '↓ Sync'; }
    if (result.success) {
      logAppend('\n✅ Done!');
      courseTree = await _fetchCourses();
      allCourses = flattenTree(courseTree);
      filterAndRenderSidebar();
      loadFiles();
      toast('Course updated!', 'ok');
    } else {
      toast('Sync failed.', 'err');
    }
  });
}

// ═══════════════════════════════════════════════════════════════════════════
// Resizable dividers
// ═══════════════════════════════════════════════════════════════════════════
function toggleSidebar() {
  const sb  = document.getElementById('sidebar');
  const bd  = document.getElementById('sidebar-backdrop');
  const btn = document.getElementById('sidebar-toggle');
  if (window.innerWidth <= 768) {
    const open = sb.classList.toggle('mobile-open');
    bd.classList.toggle('visible', open);
    btn.textContent = open ? '✕' : '☰';
  } else {
    const collapsed = sb.classList.toggle('collapsed');
    localStorage.setItem('sidebar_collapsed', collapsed ? '1' : '');
    btn.textContent = collapsed ? '▶' : '☰';
  }
}

function initResizeDividers() {
  // Restore collapsed state
  if (localStorage.getItem('sidebar_collapsed') === '1') {
    const sb = document.getElementById('sidebar');
    sb.classList.add('collapsed');
    document.getElementById('sidebar-toggle').textContent = '▶';
  }
  const savedSidebar = localStorage.getItem('sidebar_width');
  if (savedSidebar) document.getElementById('sidebar').style.width = savedSidebar + 'px';
  const savedFilesCol = localStorage.getItem('files_col_width');
  if (savedFilesCol) document.getElementById('files-list-col').style.width = savedFilesCol + 'px';
  const savedNotesCol = localStorage.getItem('notes_col_width');
  if (savedNotesCol) document.getElementById('files-notes-col').style.width = savedNotesCol + 'px';

  setupColResize(document.getElementById('divider-sidebar'),  document.getElementById('sidebar'),        140, 480, 'sidebar_width');
  setupColResize(document.getElementById('divider-files'),    document.getElementById('files-list-col'), 120, 500, 'files_col_width');
  setupColResize(document.getElementById('divider-notes'),    document.getElementById('files-notes-col'), 180, 600, 'notes_col_width', true);
}

function setupColResize(divider, targetEl, minW, maxW, storageKey, reverse = false) {
  let startX, startW;
  divider.addEventListener('mousedown', e => {
    e.preventDefault();
    startX = e.clientX;
    startW = targetEl.getBoundingClientRect().width;
    divider.classList.add('dragging');
    targetEl.style.transition = 'none'; // disable CSS transition while dragging
    // Prevent iframes from swallowing mouse events during drag
    document.querySelectorAll('iframe').forEach(f => f.classList.add('no-pointer'));

    function onMove(e) {
      const delta = reverse ? -(e.clientX - startX) : (e.clientX - startX);
      const newW = Math.min(maxW, Math.max(minW, startW + delta));
      targetEl.style.width = newW + 'px';
      if (storageKey === 'sidebar_width') targetEl.style.minWidth = newW + 'px';
    }
    function onUp() {
      divider.classList.remove('dragging');
      targetEl.style.transition = ''; // restore CSS transition
      document.querySelectorAll('iframe').forEach(f => f.classList.remove('no-pointer'));
      localStorage.setItem(storageKey, parseInt(targetEl.style.width));
      document.removeEventListener('mousemove', onMove);
      document.removeEventListener('mouseup', onUp);
    }
    document.addEventListener('mousemove', onMove);
    document.addEventListener('mouseup', onUp);
  });
}

// ═══════════════════════════════════════════════════════════════════════════
// Toast notifications
// ═══════════════════════════════════════════════════════════════════════════
function toast(msg, type = '') {
  const container = document.getElementById('toast-container');
  const el = document.createElement('div');
  el.className = `toast${type ? ' toast-' + type : ''}`;
  el.innerHTML = `${type === 'ok' ? '✓' : type === 'err' ? '✗' : 'ℹ'} ${esc(msg)}`;
  container.appendChild(el);
  setTimeout(() => { el.style.opacity = '0'; el.style.transition = 'opacity .3s'; setTimeout(() => el.remove(), 300); }, 2500);
}

// ═══════════════════════════════════════════════════════════════════════════
// Shortcuts overlay
// ═══════════════════════════════════════════════════════════════════════════
function showShortcuts() { document.getElementById('shortcuts-overlay').classList.add('open'); }
function hideShortcuts() { document.getElementById('shortcuts-overlay').classList.remove('open'); }

// ═══════════════════════════════════════════════════════════════════════════
// Confirm modal
// ═══════════════════════════════════════════════════════════════════════════
let confirmCallback = null;
function showConfirm(title, message, callback) {
  document.getElementById('confirm-title').textContent   = title;
  document.getElementById('confirm-message').textContent = message;
  confirmCallback = callback;
  document.getElementById('confirm-overlay').classList.add('open');
  document.getElementById('confirm-ok').onclick = () => { hideConfirm(); callback(); };
}
function hideConfirm() { document.getElementById('confirm-overlay').classList.remove('open'); }

// ═══════════════════════════════════════════════════════════════════════════
// Helpers
// ═══════════════════════════════════════════════════════════════════════════
function setLoading(on) {
  document.getElementById('tab-spinner').style.display = on ? 'inline-block' : 'none';
  document.querySelectorAll('.tbtn, .flash-btn').forEach(b => b.disabled = on);
}

function logShow(text) {
  const box = document.getElementById('log-box');
  box.style.display = 'flex';
  document.getElementById('log-content').textContent = text;
  box.scrollTop = box.scrollHeight;
}
function logAppend(text) {
  const el = document.getElementById('log-content');
  el.textContent += text;
  document.getElementById('log-box').scrollTop = 9999;
}

function esc(s) {
  return String(s).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');
}
function enc(s) { return encodeURIComponent(s); }

/** Render card text: protect LaTeX, HTML-escape, apply **bold** / *italic* markdown. */
function mdToHtml(text) {
  text = String(text);
  // 1. Stash LaTeX so it doesn't get escaped or processed
  const stash = [];
  text = text.replace(/\$\$[\s\S]+?\$\$|\$[^$\n]+?\$/g, m => {
    stash.push(m);
    return `\x00${stash.length - 1}\x00`;
  });
  // 2. HTML-escape the rest
  text = text.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
  // 3. Markdown: bold then italic
  text = text.replace(/\*\*(.+?)\*\*/g, '<strong>$1</strong>');
  text = text.replace(/\*([^*\n]+?)\*/g, '<em>$1</em>');
  // 4. Restore LaTeX
  text = text.replace(/\x00(\d+)\x00/g, (_, i) => stash[+i]);
  return text;
}

function renderLatexIn(elId) {
  const el = typeof elId === 'string' ? document.getElementById(elId) : elId;
  if (!el) return;
  if (!window.renderMathInElement) {
    setTimeout(() => renderLatexIn(el), 150);
    return;
  }
  renderMathInElement(el, {
    delimiters: [
      { left: '$$', right: '$$', display: true  },
      { left: '$',  right: '$',  display: false },
      { left: '\\[', right: '\\]', display: true  },
      { left: '\\(', right: '\\)', display: false },
    ],
    throwOnError: false,
  });
}

// ═══════════════════════════════════════════════════════════════════════════
// Command Palette
// ═══════════════════════════════════════════════════════════════════════════
let _cmdIndex = 0;
let _cmdItems = [];

function openCmdPalette() {
  document.getElementById('cmd-overlay').classList.add('open');
  const input = document.getElementById('cmd-input');
  input.value = '';
  renderPaletteResults('');
  setTimeout(() => input.focus(), 10);
}

function closeCmdPalette() {
  document.getElementById('cmd-overlay').classList.remove('open');
}

function _buildPaletteItems(q) {
  const items = [];
  const ql = q.toLowerCase();

  const actions = [
    { icon: '🏠', label: 'Go home',                 meta: 'Navigation',   fn: () => goHome() },
    { icon: '↓',  label: 'Sync new files',          meta: 'Sync',         fn: () => runScraper() },
    { icon: '🧠', label: 'Study all cards',         meta: 'Learn',        fn: () => startGlobalLearn() },
    { icon: '⌨️', label: 'Show shortcuts',          meta: 'Help',         fn: () => showShortcuts() },
    { icon: '🌙', label: 'Toggle light/dark',       meta: 'Settings',     fn: () => toggleTheme() },
    { icon: '🔍', label: 'Search summaries',        meta: 'Search',       fn: () => { document.getElementById('search-global').focus(); } },
  ];

  if (activeCourse) {
    const cname = activeCourse.split('/').pop();
    actions.push(
      { icon: '📁', label: `Files — ${cname}`,                   meta: 'Tab',    fn: () => switchTab('files') },
      { icon: '📄', label: `Summary — ${cname}`,                 meta: 'Tab',    fn: () => switchTab('summary') },
      { icon: '🧠', label: `Study — ${cname}`,                   meta: 'Tab',    fn: () => switchTab('learn') },
      { icon: '✏️', label: `Notes — ${cname}`,                   meta: 'Tab',    fn: () => switchTab('notes') },
      { icon: '💬', label: `Chat — ${cname}`,                    meta: 'Tab',    fn: () => switchTab('chat') },
      { icon: '✨', label: `Create summary — ${cname}`,          meta: 'Action', fn: () => { switchTab('files'); setTimeout(() => generateSummary(), 100); } },
      { icon: '↓',  label: `Sync course — ${cname}`,             meta: 'Sync',   fn: () => syncCourse() },
    );
  }

  const filteredActions = ql
    ? actions.filter(a => a.label.toLowerCase().includes(ql) || (a.meta || '').toLowerCase().includes(ql))
    : actions;

  filteredActions.forEach(a => items.push({ type: 'action', ...a }));

  const courses = allCourses
    .filter(c => !ql || c.name.toLowerCase().includes(ql) || c.path.toLowerCase().includes(ql))
    .slice(0, ql ? 20 : 8);

  courses.forEach(c => items.push({
    type: 'course',
    icon: c.has_summary ? '📚' : '📁',
    label: c.name,
    meta: c.path.includes('/') ? c.path.split('/')[0] : '',
    fn: () => selectCourse(c.path),
  }));

  return items;
}

function renderPaletteResults(q) {
  _cmdItems = _buildPaletteItems(q);
  _cmdIndex = 0;

  const el = document.getElementById('cmd-results');
  if (!_cmdItems.length) {
    el.innerHTML = '<div id="cmd-empty">No results</div>';
    return;
  }

  const actions = _cmdItems.filter(i => i.type === 'action');
  const courses = _cmdItems.filter(i => i.type === 'course');

  let html = '';
  if (actions.length) {
    html += `<div class="cmd-section">Aktionen</div>`;
    actions.forEach(item => {
      const idx = _cmdItems.indexOf(item);
      html += `<div class="cmd-item${idx === _cmdIndex ? ' cmd-selected' : ''}" data-idx="${idx}"
        onmouseenter="setCmdIndex(${idx})" onclick="executePaletteItem(${idx})">
        <span class="cmd-item-icon">${item.icon}</span>
        <span class="cmd-item-label">${esc(item.label)}</span>
        ${item.meta ? `<span class="cmd-item-meta">${esc(item.meta)}</span>` : ''}
      </div>`;
    });
  }
  if (courses.length) {
    html += `<div class="cmd-section">Kurse</div>`;
    courses.forEach(item => {
      const idx = _cmdItems.indexOf(item);
      html += `<div class="cmd-item${idx === _cmdIndex ? ' cmd-selected' : ''}" data-idx="${idx}"
        onmouseenter="setCmdIndex(${idx})" onclick="executePaletteItem(${idx})">
        <span class="cmd-item-icon">${item.icon}</span>
        <span class="cmd-item-label">${esc(item.label)}</span>
        ${item.meta ? `<span class="cmd-item-meta">${esc(item.meta)}</span>` : ''}
      </div>`;
    });
  }
  el.innerHTML = html;
}

function setCmdIndex(i) {
  _cmdIndex = i;
  document.querySelectorAll('#cmd-results .cmd-item').forEach(el => {
    el.classList.toggle('cmd-selected', parseInt(el.dataset.idx) === i);
  });
}

function handlePaletteKey(e) {
  if (e.key === 'Escape')     { closeCmdPalette(); return; }
  if (e.key === 'ArrowDown')  { e.preventDefault(); setCmdIndex(Math.min(_cmdIndex + 1, _cmdItems.length - 1)); _scrollPaletteItem(); return; }
  if (e.key === 'ArrowUp')    { e.preventDefault(); setCmdIndex(Math.max(_cmdIndex - 1, 0)); _scrollPaletteItem(); return; }
  if (e.key === 'Enter')      { e.preventDefault(); executePaletteItem(_cmdIndex); return; }
}

function _scrollPaletteItem() {
  const el = document.querySelector(`#cmd-results .cmd-item[data-idx="${_cmdIndex}"]`);
  if (el) el.scrollIntoView({ block: 'nearest' });
}

function executePaletteItem(idx) {
  const item = _cmdItems[idx];
  if (!item) return;
  closeCmdPalette();
  setTimeout(() => item.fn(), 50);
}

// ═══════════════════════════════════════════════════════════════════════════
// Summary Table of Contents
// ═══════════════════════════════════════════════════════════════════════════
let _tocObserver = null;

function _buildAndInjectToC() {
  const layout    = document.getElementById('summary-layout');
  const contentEl = document.getElementById('summary-md-content');
  if (!layout || !contentEl) return;

  const headings = [...contentEl.querySelectorAll('h2, h3')];
  if (headings.length < 1) return;

  let html = '<div class="summary-toc"><div class="summary-toc-title">Contents</div>';
  headings.forEach((h, i) => {
    const id = `sh-${i}`;
    h.id = id;
    const cls = h.tagName === 'H3' ? 'toc-h3' : '';
    // Strip "File N: " prefix from h2 labels so they're short and clean
    let label = h.textContent.trim().replace(/^File\s+\d+:\s*/i, '').replace(/^\[|\]$/g, '');
    html += `<a class="toc-link ${cls}" data-hid="${id}" href="#"
      onclick="event.preventDefault();_scrollToSummaryHeading('${id}')"
      title="${esc(h.textContent.trim())}">${esc(label)}</a>`;
  });
  html += '</div>';

  layout.insertAdjacentHTML('afterbegin', html);
  _setupTocObserver();
}

function _scrollToSummaryHeading(id) {
  const el    = document.getElementById(id);
  const panel = document.getElementById('panel-summary');
  if (el && panel) {
    const offset = el.getBoundingClientRect().top - panel.getBoundingClientRect().top + panel.scrollTop - 60;
    panel.scrollTo({ top: offset, behavior: 'smooth' });
  }
}

function _setupTocObserver() {
  if (_tocObserver) _tocObserver.disconnect();
  const panel = document.getElementById('panel-summary');
  if (!panel) return;
  const headings = panel.querySelectorAll('[id^="sh-"]');
  if (!headings.length) return;

  _tocObserver = new IntersectionObserver(entries => {
    entries.forEach(entry => {
      const link = panel.querySelector(`.toc-link[data-hid="${entry.target.id}"]`);
      if (link) link.classList.toggle('toc-active', entry.isIntersecting);
    });
  }, { root: panel, rootMargin: '-5% 0px -75% 0px', threshold: 0 });

  headings.forEach(h => _tocObserver.observe(h));
}

// ═══════════════════════════════════════════════════════════════════════════
// Selection toolbar helpers
// ═══════════════════════════════════════════════════════════════════════════
function updateSelectionToolbar() {
  const count  = document.querySelectorAll('input[name="file"]:checked').length;
  const countEl = document.getElementById('sel-count');
  const sumBtn  = document.getElementById('sel-summarize-btn');
  if (countEl) countEl.textContent = `${count} file${count !== 1 ? 's' : ''} selected`;
  if (sumBtn)  sumBtn.textContent  = `Summarize selected (${count})`;
}

function markSelectedFilesRead(read) {
  document.querySelectorAll('input[name="file"]:checked').forEach(cb => {
    setFileRead(activeCourse, cb.value, read);
  });
  if (read) _clearReadTimer();
  filterAndRenderSidebar();
}

// ── Welcome modal ─────────────────────────────────────────────────────────
function initWelcome() {
  if (!localStorage.getItem('welcome_seen')) {
    document.getElementById('welcome-overlay').classList.remove('hidden');
  }
}
function closeWelcome() {
  localStorage.setItem('welcome_seen', '1');
  const el = document.getElementById('welcome-overlay');
  el.style.opacity = '0';
  el.style.transition = 'opacity .2s ease';
  setTimeout(() => el.classList.add('hidden'), 200);
}
document.addEventListener('keydown', e => {
  if (e.key === 'Escape') closeWelcome();
});
document.getElementById('welcome-overlay').addEventListener('click', e => {
  if (e.target === e.currentTarget) closeWelcome();
});

boot();
initWelcome();
</script>
</body>
</html>"""

@app.route("/api/chat/<path:course_name>", methods=["POST"])
def api_chat(course_name):
    from anthropic import Anthropic
    from flask import Response, stream_with_context
    question  = request.json.get("question", "").strip()
    history   = request.json.get("history", [])   # [{role, content}, …]
    file_name = request.json.get("file", "").strip()
    if not question:
        return jsonify({"error": "No question provided"}), 400

    course_dir = COURSES_DIR / course_name
    notes_path = course_dir / NOTES_FILENAME
    notes      = notes_path.read_text(encoding="utf-8")[:3000] if notes_path.exists() else ""

    if file_name:
        # Use raw file text as context
        from summarize import extract_text as _extract_text
        file_path = course_dir / file_name
        raw_text  = _extract_text(file_path)[:15000] if file_path.exists() else ""
        system = f"""You are a concise study assistant for the course "{course_name.split('/')[-1]}". \
Always respond in English, keep answers brief and learning-focused.

The user is asking about the specific file: {file_name}

<file_content>
{raw_text or "Could not extract text from this file."}
</file_content>
{f"<notes>{notes}</notes>" if notes else ""}

Answer questions based solely on the file content above. \
If something is not covered in the file, say so."""
    else:
        summary_path = get_latest_summary(course_dir)
        context = summary_path.read_text(encoding="utf-8")[:10000] if summary_path and summary_path.exists() else ""
        system = f"""You are a concise study assistant for the course "{course_name.split('/')[-1]}". \
Always respond in English, keep answers brief and learning-focused.

<course_summary>
{context or "No summary available."}
</course_summary>
{f"<notes>{notes}</notes>" if notes else ""}

Answer questions based solely on these materials. \
If something is not covered, say so."""

    messages = history + [{"role": "user", "content": question}]

    def generate():
        client = Anthropic()
        with client.messages.stream(
            model="claude-haiku-4-5-20251001",
            max_tokens=800,
            system=system,
            messages=messages,
        ) as stream:
            for text in stream.text_stream:
                yield f"data: {json.dumps({'delta': text})}\n\n"
        yield "data: [DONE]\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

@app.route("/api/file-meta/<path:course_name>")
def api_file_meta(course_name):
    d = COURSES_DIR / course_name
    if not d.is_dir():
        return jsonify([])
    result = []
    for fname in list_files(d):
        p = d / fname
        if p.exists():
            result.append({"name": fname, "size": p.stat().st_size, "mtime": int(p.stat().st_mtime)})
    return jsonify(result)

@app.route("/api/notes-preview", methods=["POST"])
def api_notes_preview():
    text = request.json.get("text", "")
    html = markdown2.markdown(text, extras=["fenced-code-blocks", "tables", "task_list"])
    return jsonify({"html": html})

@app.route("/")
def index():
    return render_template_string(HTML)

if __name__ == "__main__":
    print("Dashboard → http://localhost:5001")
    app.run(debug=False, port=5001)
