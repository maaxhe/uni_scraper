"""
Summarizer: liest heruntergeladene Stud.IP-Dateien, erstellt pro Kurs eine
einzige Zusammenfassungs-Markdown-Datei mit Trainingsfragen via Claude API.

Usage:
    python summarize.py                                     # alle Kurse
    python summarize.py --course italienisch                # Kurs per Teilname
    python summarize.py --dir /abs/path/to/course          # direkter Pfad (Dashboard nutzt das)
    python summarize.py --course italienisch --limit 5
    python summarize.py --course italienisch --force        # neu generieren
    python summarize.py --lang en                          # Englisch (Standard)
    python summarize.py --lang de                          # Deutsch
"""

import argparse
import logging
import os
import re
import sys
from datetime import date
from pathlib import Path

import anthropic
import fitz  # PyMuPDF
from docx import Document
from dotenv import load_dotenv

load_dotenv()

COURSES_DIR = Path(os.environ.get("COURSES_DIR", "/Users/maxmacbookpro/Documents/Uni/Cognitive Science [Course]/Courses"))
SUPPORTED_EXTENSIONS = {".pdf", ".doc", ".docx", ".txt", ".md", ".pptx"}
MAX_CHARS = 40_000
MODEL_ANTHROPIC = os.environ.get("ANTHROPIC_MODEL", "claude-opus-4-6")
MODEL_OPENAI    = os.environ.get("OPENAI_MODEL", "gpt-4o-mini")
OUTPUT_FILENAME = "_zusammenfassung.md"

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# File reading
# ---------------------------------------------------------------------------

def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            with fitz.open(str(path)) as doc:
                return "\n".join(page.get_text() for page in doc)
        elif suffix in {".doc", ".docx"}:
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        elif suffix in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="replace")
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts)
    except Exception as exc:
        log.warning("Konnte Text nicht lesen aus %s: %s", path.name, exc)
    return ""


# ---------------------------------------------------------------------------
# Claude API
# ---------------------------------------------------------------------------

PROMPTS = {
    "en": {
        "intro": (
            "You are a concise study assistant for university students.\n\n"
            "Course: {course_name}\n\n"
            "I'm giving you {n} file(s). For each file produce a study note using bullet points. "
            "Group bullets into clearly separated sub-sections with blank lines between them "
            "so the content breathes and is easy to scan. Each logical topic or concept cluster "
            "should be its own paragraph/group — do NOT lump everything into one continuous list.\n\n"
            "{length_instruction}\n\n"
            "HIGHLIGHTING: Use **bold** for the most important terms, concepts, and keywords "
            "throughout every section — not just in Key Concepts. Bold 2–4 words per bullet "
            "where they are genuinely significant. Do not over-bold.\n\n"
            "IMPORTANT: Always use LaTeX notation for ALL mathematical expressions, "
            "equations, symbols, and formulas. "
            "Inline math: $...$ (e.g. $E = mc^2$, $\\alpha$, $\\nabla f$). "
            "Display/block math: $$...$$ for standalone equations. "
            "Never write math as plain text."
        ),
        "section": "## [Filename]",
        "summary_label": "**Summary** (3–5 bullets, one idea each)",
        "concepts_label": "**Key Concepts**",
        "questions_label": "**Training Questions**",
        "answers_label": "**Answers**",
        "header": "# Summary: {name}\n\n*Generated on {date}*  \n*{n} of {total} file(s) summarised*\n\n---\n\n",
        "truncated": "\n\n[... text truncated ...]",
        "file_label": "### File {i}: {name}",
    },
    "de": {
        "intro": (
            "Du bist ein prägnanter Lernassistent für Studierende.\n\n"
            "Kurs: {course_name}\n\n"
            "Ich gebe dir {n} Datei(en). Erstelle für jede Datei eine Lernnotiz in Stichpunkten. "
            "Gruppiere die Bullets in klar getrennte Unterabschnitte mit Leerzeilen dazwischen, "
            "damit der Inhalt Luft hat und leicht zu scannen ist. Jedes logische Thema oder "
            "Konzept-Cluster soll eine eigene Gruppe/Absatz bilden — "
            "alles in eine einzige Liste zu packen ist nicht erwünscht.\n\n"
            "{length_instruction}\n\n"
            "HERVORHEBUNGEN: Verwende **Fettschrift** für die wichtigsten Begriffe, Konzepte "
            "und Schlüsselwörter in allen Abschnitten — nicht nur in den Kernkonzepten. "
            "2–4 Wörter pro Bullet fett setzen, wo sie wirklich bedeutsam sind. Nicht übertreiben.\n\n"
            "WICHTIG: Verwende für ALLE mathematischen Ausdrücke, Gleichungen, Symbole "
            "und Formeln immer LaTeX-Notation. "
            "Inline-Math: $...$ (z.B. $E = mc^2$, $\\alpha$, $\\nabla f$). "
            "Block-Math: $$...$$ für eigenständige Gleichungen. "
            "Schreibe Mathematik niemals als reinen Text."
        ),
        "section": "## [Dateiname]",
        "summary_label": "**Zusammenfassung** (3–5 Bullets, je eine Kernaussage)",
        "concepts_label": "**Kernkonzepte**",
        "questions_label": "**Trainingsfragen**",
        "answers_label": "**Antworten**",
        "header": "# Zusammenfassung: {name}\n\n*Generiert am {date}*  \n*{n} von {total} Datei(en) zusammengefasst*\n\n---\n\n",
        "truncated": "\n\n[... Text gekürzt ...]",
        "file_label": "### Datei {i}: {name}",
    },
}


_LENGTH_INSTRUCTIONS = {
    "en": {
        "short": "LENGTH: Aim for roughly 300 words per file. Be very concise — only the most essential points.",
        "long":  "LENGTH: Aim for roughly 1000 words per file. Be thorough — cover all key ideas, details, and nuances.",
    },
    "de": {
        "short": "LÄNGE: Ziele auf ca. 300 Wörter pro Datei. Sehr knapp — nur die wichtigsten Punkte.",
        "long":  "LÄNGE: Ziele auf ca. 1000 Wörter pro Datei. Ausführlich — alle wichtigen Ideen, Details und Nuancen.",
    },
}

def summarize_files(client, client_type: str, course_name: str, files: list[dict], lang: str = "en", length: str = "short") -> str:
    p = PROMPTS.get(lang, PROMPTS["en"])
    length_instructions = _LENGTH_INSTRUCTIONS.get(lang, _LENGTH_INSTRUCTIONS["en"])
    length_instruction  = length_instructions.get(length, length_instructions["short"])

    files_block = ""
    for i, f in enumerate(files, 1):
        text = f["text"][:MAX_CHARS]
        if len(f["text"]) > MAX_CHARS:
            text += p["truncated"]
        files_block += f"\n\n{p['file_label'].format(i=i, name=f['name'])}\n\n{text}"

    prompt = f"""{p['intro'].format(course_name=course_name, n=len(files), length_instruction=length_instruction)}

Use this exact structure for each file (all bullet points, no prose):

{p['section']}

{p['summary_label']}
- ...
- ...

{p['concepts_label']}
- **Term**: one-line definition
- **Term**: one-line definition

{p['questions_label']}
1. Question?
2. Question?
3. Question?

{p['answers_label']}
1. Answer (one sentence)
2. Answer (one sentence)
3. Answer (one sentence)

---

Files:{files_block}"""

    if client_type == "anthropic":
        try:
            response = client.messages.create(
                model=MODEL_ANTHROPIC,
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}],
            )
        except anthropic.AuthenticationError:
            log.error("❌ API key invalid — create a new key at https://console.anthropic.com")
            sys.exit(1)
        except anthropic.PermissionDeniedError:
            log.error("❌ API key lacks permission — check your Anthropic account and plan.")
            sys.exit(1)
        except anthropic.RateLimitError:
            log.error("❌ Rate limit hit — too many requests. Wait a moment and try again.")
            sys.exit(1)
        except anthropic.APIStatusError as e:
            if e.status_code in (402, 403):
                log.error("❌ API quota exhausted or billing issue (HTTP %d). "
                          "Check your usage at https://console.anthropic.com", e.status_code)
            else:
                log.error("❌ Anthropic API error (HTTP %d): %s", e.status_code, e.message)
            sys.exit(1)
        except anthropic.APIConnectionError:
            log.error("❌ Cannot reach Anthropic API — check your internet connection.")
            sys.exit(1)
        except anthropic.APITimeoutError:
            log.error("❌ Request timed out — the API took too long to respond. Try again.")
            sys.exit(1)
        return response.content[0].text
    else:  # openai-compatible (OpenAI, Groq, Mistral, Ollama, …)
        try:
            response = client.chat.completions.create(
                model=MODEL_OPENAI,
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}],
            )
        except Exception as e:
            err = str(e).lower()
            if "401" in err or "authentication" in err or "invalid" in err and "key" in err:
                log.error("❌ API key invalid — check OPENAI_API_KEY in your .env file.")
            elif "429" in err or "rate limit" in err:
                log.error("❌ Rate limit hit — too many requests. Wait a moment and try again.")
            elif "402" in err or "quota" in err or "billing" in err or "insufficient_quota" in err:
                log.error("❌ API quota exhausted or billing issue. Check your account.")
            elif "connection" in err or "timeout" in err:
                log.error("❌ Cannot reach API — check your internet connection.")
            else:
                log.error("❌ API error: %s", e)
            sys.exit(1)
        return response.choices[0].message.content


# ---------------------------------------------------------------------------
# Q&A split
# ---------------------------------------------------------------------------

_QA_BLOCK_RE = re.compile(
    r'\*\*(?:Training Questions|Trainingsfragen|Answers|Antworten)\*\*\n(?:\d+\..+\n?)*',
    re.MULTILINE,
)

def _split_qa(text: str) -> tuple[str, str]:
    """Split AI output into (clean_summary, qa_only).

    The qa_only string preserves ## section headers so parse_flashcards()
    can still associate questions with their file/section context.
    """
    chunks = re.split(r'(^## .+$)', text, flags=re.MULTILINE)
    clean_parts: list[str] = []
    qa_parts:    list[str] = []

    i = 0
    if chunks:
        clean_parts.append(chunks[0])   # content before first ##
        i = 1

    while i < len(chunks) - 1:
        header = chunks[i]          # '## File 1: name.pdf'
        body   = chunks[i + 1]

        # -- qa file: keep header + Q&A blocks only --
        q_match = re.search(
            r'\*\*(Training Questions|Trainingsfragen)\*\*\n((?:\d+\..+\n?)+)', body)
        a_match = re.search(
            r'\*\*(Answers|Antworten)\*\*\n((?:\d+\..+\n?)+)', body)
        if q_match:
            qa_block = f"{header}\n\n**{q_match.group(1)}**\n{q_match.group(2).rstrip()}"
            if a_match:
                qa_block += f"\n\n**{a_match.group(1)}**\n{a_match.group(2).rstrip()}"
            qa_parts.append(qa_block)

        # -- clean summary: remove Q&A blocks --
        clean_body = _QA_BLOCK_RE.sub('', body)
        clean_body = re.sub(r'\n{3,}', '\n\n', clean_body)
        clean_parts.append(header + clean_body)
        i += 2

    clean = re.sub(r'\n{3,}', '\n\n', ''.join(clean_parts)).strip()
    qa    = '\n\n---\n\n'.join(qa_parts)
    return clean, qa


# ---------------------------------------------------------------------------
# Course processing
# ---------------------------------------------------------------------------

def find_course(name_query: str) -> Path | None:
    """Find a course folder by partial case-insensitive name match (recursive 2 levels)."""
    query = name_query.lower()
    # First try top-level
    for d in COURSES_DIR.iterdir():
        if d.is_dir() and query in d.name.lower():
            return d
    # Then try one level deeper (groups)
    for group in COURSES_DIR.iterdir():
        if not group.is_dir():
            continue
        for d in group.iterdir():
            if d.is_dir() and query in d.name.lower():
                return d
    return None


def process_course(
    client,
    client_type: str,
    course_dir: Path,
    limit: int,
    force: bool,
    only_files: list[str] | None = None,
    lang: str = "en",
    output_filename: str | None = None,
    length: str = "short",
) -> None:
    output_path = course_dir / (output_filename or OUTPUT_FILENAME)

    if output_path.exists() and not force:
        log.info("  ✓ Summary already exists: %s", output_path)
        log.info("  Use --force to regenerate.")
        return

    all_files = [
        f for f in course_dir.rglob("*")
        if f.is_file()
        and f.suffix.lower() in SUPPORTED_EXTENSIONS
        and ".summary" not in f.name
        and f.name != OUTPUT_FILENAME
    ]

    if not all_files:
        log.info("  No files found.")
        return

    if only_files:
        only_set = set(only_files)
        all_files = [
            f for f in all_files
            if f.name in only_set or str(f.relative_to(course_dir)) in only_set
        ]

    files_to_process = all_files[:limit]
    if len(all_files) > limit:
        log.info("  %d file(s) found, processing %d (--limit %d)", len(all_files), limit, limit)
    else:
        log.info("  %d file(s) found", len(all_files))

    file_contents = []
    for f in files_to_process:
        text = extract_text(f)
        if text.strip():
            file_contents.append({"name": f.name, "text": text})
            log.info("  ✎ %s", f.name)
        else:
            log.warning("  No text: %s — skipped", f.name)

    if not file_contents:
        log.warning("  No readable files.")
        return

    log.info("  → Sending %d file(s) to AI…", len(file_contents))
    try:
        summary = summarize_files(client, client_type, course_dir.name, file_contents, lang=lang, length=length)
        p = PROMPTS.get(lang, PROMPTS["en"])
        header = p["header"].format(
            name=course_dir.name,
            date=date.today(),
            n=len(file_contents),
            total=len(all_files),
        )
        clean_summary, qa_only = _split_qa(summary)
        output_path.write_text(header + clean_summary, encoding="utf-8")
        log.info("  ✓ Saved: %s", output_path)
        if qa_only:
            qa_path = output_path.parent / (output_path.stem + "_qa.md")
            qa_path.write_text(header + qa_only, encoding="utf-8")
            log.info("  ✓ Saved Q&A: %s", qa_path)
    except Exception as exc:
        log.error("  Error: %s", exc)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(description="Stud.IP Course Summarizer")
    parser.add_argument("--course", metavar="NAME", help="Course name (partial match)")
    parser.add_argument("--dir", metavar="PATH", help="Absolute path to course directory (overrides --course)")
    parser.add_argument("--limit", type=int, default=3, help="Max files per course (default: 3)")
    parser.add_argument("--force", action="store_true", help="Overwrite existing summary")
    parser.add_argument("--files", nargs="+", metavar="FILE", help="Only summarise these files")
    parser.add_argument("--lang", default="en", choices=["en", "de"], help="Summary language (default: en)")
    parser.add_argument("--out", metavar="FILENAME", default=None, help="Output filename (overrides default _zusammenfassung.md)")
    parser.add_argument("--length", default="short", choices=["short", "long"], help="Summary length: short (~300 words) or long (~1000 words)")
    args = parser.parse_args()

    # ── Build AI client — Anthropic takes priority, then any OpenAI-compatible provider ──
    anthropic_key = os.environ.get("ANTHROPIC_API_KEY")
    openai_key    = os.environ.get("OPENAI_API_KEY")

    if anthropic_key:
        client      = anthropic.Anthropic(api_key=anthropic_key)
        client_type = "anthropic"
        log.info("Using Anthropic (model: %s)", MODEL_ANTHROPIC)
    elif openai_key:
        try:
            from openai import OpenAI
        except ImportError:
            log.error("openai package not installed. Run: pip install openai")
            sys.exit(1)
        base_url    = os.environ.get("OPENAI_BASE_URL") or None
        client      = OpenAI(api_key=openai_key, base_url=base_url)
        client_type = "openai"
        log.info("Using OpenAI-compatible API (model: %s%s)", MODEL_OPENAI,
                 f", base_url: {base_url}" if base_url else "")
    else:
        log.error(
            "No API key found in .env.\n"
            "  Set ANTHROPIC_API_KEY  →  https://console.anthropic.com\n"
            "  or OPENAI_API_KEY      →  https://platform.openai.com/api-keys"
        )
        sys.exit(1)

    if args.dir:
        course_dir = Path(args.dir)
        if not course_dir.is_dir():
            log.error("Directory not found: %s", args.dir)
            sys.exit(1)
        log.info("── Course: %s", course_dir.name)
        process_course(client, client_type, course_dir, args.limit, args.force, only_files=args.files, lang=args.lang, output_filename=args.out, length=args.length)
    elif args.course:
        course_dir = find_course(args.course)
        if not course_dir:
            log.error("No course found for: '%s'", args.course)
            sys.exit(1)
        log.info("── Course: %s", course_dir.name)
        process_course(client, client_type, course_dir, args.limit, args.force, only_files=args.files, lang=args.lang, output_filename=args.out, length=args.length)
    else:
        for course_dir in sorted(COURSES_DIR.iterdir()):
            if course_dir.is_dir():
                log.info("── Course: %s", course_dir.name)
                process_course(client, client_type, course_dir, args.limit, args.force, lang=args.lang, length=args.length)

    log.info("Done.")


if __name__ == "__main__":
    main()
